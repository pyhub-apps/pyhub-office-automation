"""
pyhub-office-automation Typer 기반 CLI 명령어
PyInstaller 호환성을 위한 정적 명령어 등록
"""

import json
import os
import sys
from typing import Optional

# Windows 환경에서 UTF-8 인코딩 강제 설정
if sys.platform == "win32":
    # 환경 변수 설정
    os.environ.setdefault("PYTHONIOENCODING", "utf-8")
    os.environ.setdefault("PYTHONUTF8", "1")

    # stdout/stderr 인코딩 설정
    if hasattr(sys.stdout, "reconfigure"):
        try:
            sys.stdout.reconfigure(encoding="utf-8")
            sys.stderr.reconfigure(encoding="utf-8")
        except Exception:
            pass  # 설정 실패해도 계속 진행

import typer
from rich.console import Console
from rich.table import Table

from pyhub_office_automation.cli.ai_setup import ai_setup_app
from pyhub_office_automation.email.email_accounts import accounts_app

# Email 명령어 import
from pyhub_office_automation.email.email_send import email_send

# Chart 명령어 import
from pyhub_office_automation.excel.chart_add import chart_add
from pyhub_office_automation.excel.chart_configure import chart_configure
from pyhub_office_automation.excel.chart_delete import chart_delete
from pyhub_office_automation.excel.chart_export import chart_export
from pyhub_office_automation.excel.chart_list import chart_list
from pyhub_office_automation.excel.chart_pivot_create import chart_pivot_create
from pyhub_office_automation.excel.chart_position import chart_position

# Data 명령어 import (Issue #39)
from pyhub_office_automation.excel.data_analyze import data_analyze
from pyhub_office_automation.excel.data_transform import data_transform
from pyhub_office_automation.excel.metadata_generate import metadata_generate

# Pivot 명령어 import
from pyhub_office_automation.excel.pivot_configure import pivot_configure
from pyhub_office_automation.excel.pivot_create import pivot_create
from pyhub_office_automation.excel.pivot_delete import pivot_delete
from pyhub_office_automation.excel.pivot_list import pivot_list
from pyhub_office_automation.excel.pivot_refresh import pivot_refresh

# Excel 명령어 import
from pyhub_office_automation.excel.range_convert import range_convert
from pyhub_office_automation.excel.range_read import range_read
from pyhub_office_automation.excel.range_write import range_write

# Shape 명령어 import
from pyhub_office_automation.excel.shape_add import shape_add
from pyhub_office_automation.excel.shape_delete import shape_delete
from pyhub_office_automation.excel.shape_format import shape_format
from pyhub_office_automation.excel.shape_group import shape_group
from pyhub_office_automation.excel.shape_list import shape_list
from pyhub_office_automation.excel.sheet_activate import sheet_activate
from pyhub_office_automation.excel.sheet_add import sheet_add
from pyhub_office_automation.excel.sheet_delete import sheet_delete
from pyhub_office_automation.excel.sheet_rename import sheet_rename

# Slicer 명령어 import
from pyhub_office_automation.excel.slicer_add import slicer_add
from pyhub_office_automation.excel.slicer_connect import slicer_connect
from pyhub_office_automation.excel.slicer_list import slicer_list
from pyhub_office_automation.excel.slicer_position import slicer_position
from pyhub_office_automation.excel.table_analyze import table_analyze
from pyhub_office_automation.excel.table_create import table_create
from pyhub_office_automation.excel.table_list import table_list
from pyhub_office_automation.excel.table_read import table_read
from pyhub_office_automation.excel.table_sort import table_sort
from pyhub_office_automation.excel.table_sort_clear import table_sort_clear
from pyhub_office_automation.excel.table_sort_info import table_sort_info
from pyhub_office_automation.excel.table_write import table_write
from pyhub_office_automation.excel.textbox_add import textbox_add
from pyhub_office_automation.excel.workbook_create import workbook_create
from pyhub_office_automation.excel.workbook_info import workbook_info
from pyhub_office_automation.excel.workbook_list import workbook_list
from pyhub_office_automation.excel.workbook_open import workbook_open

# HWP 명령어 import
from pyhub_office_automation.hwp.hwp_export import hwp_export

# PowerPoint Phase 3 (Windows Advanced) 명령어 import (Issue #84 Phase 3)
from pyhub_office_automation.powerpoint.animation_add import animation_add

# PowerPoint Advanced Content 명령어 import (Issue #78)
from pyhub_office_automation.powerpoint.content_add_audio import content_add_audio
from pyhub_office_automation.powerpoint.content_add_chart import content_add_chart
from pyhub_office_automation.powerpoint.content_add_equation import content_add_equation
from pyhub_office_automation.powerpoint.content_add_excel_chart import content_add_excel_chart
from pyhub_office_automation.powerpoint.content_add_image import content_add_image
from pyhub_office_automation.powerpoint.content_add_shape import content_add_shape
from pyhub_office_automation.powerpoint.content_add_smartart import content_add_smartart
from pyhub_office_automation.powerpoint.content_add_table import content_add_table

# PowerPoint Content 명령어 import (Issue #77)
from pyhub_office_automation.powerpoint.content_add_text import content_add_text
from pyhub_office_automation.powerpoint.content_add_video import content_add_video
from pyhub_office_automation.powerpoint.content_update import content_update

# PowerPoint Export 명령어 import (Issue #80)
from pyhub_office_automation.powerpoint.export_images import export_images
from pyhub_office_automation.powerpoint.export_notes import export_notes
from pyhub_office_automation.powerpoint.export_pdf import export_pdf
from pyhub_office_automation.powerpoint.layout_apply import layout_apply

# PowerPoint Layout & Theme 명령어 import (Issue #79)
from pyhub_office_automation.powerpoint.layout_list import layout_list

# PowerPoint 명령어 import
from pyhub_office_automation.powerpoint.presentation_create import presentation_create
from pyhub_office_automation.powerpoint.presentation_info import presentation_info
from pyhub_office_automation.powerpoint.presentation_list import presentation_list
from pyhub_office_automation.powerpoint.presentation_open import presentation_open
from pyhub_office_automation.powerpoint.presentation_save import presentation_save
from pyhub_office_automation.powerpoint.run_macro import run_macro
from pyhub_office_automation.powerpoint.slide_add import slide_add
from pyhub_office_automation.powerpoint.slide_copy import slide_copy
from pyhub_office_automation.powerpoint.slide_delete import slide_delete
from pyhub_office_automation.powerpoint.slide_duplicate import slide_duplicate
from pyhub_office_automation.powerpoint.slide_list import slide_list
from pyhub_office_automation.powerpoint.slide_reorder import slide_reorder
from pyhub_office_automation.powerpoint.slideshow_control import slideshow_control
from pyhub_office_automation.powerpoint.slideshow_start import slideshow_start
from pyhub_office_automation.powerpoint.template_apply import template_apply
from pyhub_office_automation.powerpoint.theme_apply import theme_apply
from pyhub_office_automation.utils.resource_loader import load_llm_guide, load_welcome_message
from pyhub_office_automation.version import get_version, get_version_info

# Typer 앱 생성
app = typer.Typer(help="pyhub-office-automation: AI 에이전트를 위한 Office 자동화 도구")


def version_callback(value: bool):
    """--version 콜백 함수"""
    if value:
        version_info = get_version_info()
        typer.echo(f"pyhub-office-automation version {version_info['version']}")
        raise typer.Exit()


# 글로벌 --version 옵션 추가 및 기본 메시지 표시
@app.callback(invoke_without_command=True)
def main_callback(
    ctx: typer.Context,
    version: bool = typer.Option(False, "--version", "-v", callback=version_callback, help="버전 정보 출력"),
):
    """
    pyhub-office-automation: AI 에이전트를 위한 Office 자동화 도구
    """
    # 서브커맨드가 없고 버전 옵션도 아닌 경우 welcome 메시지 표시
    if ctx.invoked_subcommand is None:
        show_welcome_message()


def show_welcome_message():
    """Welcome 메시지를 표시합니다."""
    welcome_content = load_welcome_message()
    console.print(welcome_content)

    # LLM 가이드 안내 추가
    console.print("\n💡 [bold cyan]AI 에이전트 사용 시 상세 지침을 보려면:[/bold cyan]")
    console.print("   oa llm-guide")


# version 명령어 추가
@app.command()
def version():
    """버전 정보 출력"""
    version_info = get_version_info()
    typer.echo(f"pyhub-office-automation version {version_info['version']}")


@app.command()
def welcome(output_format: str = typer.Option("text", "--format", help="출력 형식 선택 (text, json)")):
    """환영 메시지 및 시작 가이드 출력"""
    welcome_content = load_welcome_message()

    if output_format == "json":
        # JSON 형식으로 출력
        welcome_data = {
            "message_type": "welcome",
            "content": welcome_content,
            "package_version": get_version(),
            "available_commands": {
                "info": "패키지 정보 및 설치 상태",
                "excel": "Excel 자동화 명령어들",
                "hwp": "HWP 자동화 명령어들 (Windows 전용)",
                "install-guide": "설치 가이드",
                "llm-guide": "AI 에이전트 사용 지침",
            },
        }
        try:
            json_output = json.dumps(welcome_data, ensure_ascii=False, indent=2)
            typer.echo(json_output)
        except UnicodeEncodeError:
            json_output = json.dumps(welcome_data, ensure_ascii=True, indent=2)
            typer.echo(json_output)
    else:
        console.print(welcome_content)


@app.command()
def llm_guide(
    ai_type: str = typer.Argument(
        "default", help="AI 어시스턴트 타입 [default|codex|claude|gemini|copilot]", case_sensitive=False, show_default=True
    ),
    format: str = typer.Option("json", "--format", "-f", help="출력 형식 [json|text|markdown]", case_sensitive=False),
    verbose: bool = typer.Option(False, "--verbose", "-v", help="상세 가이드 출력"),
    lang: str = typer.Option("ko", "--lang", "-l", help="언어 선택 [ko|en]", case_sensitive=False),
):
    """AI 어시스턴트별 맞춤형 사용 가이드 제공

    각 AI의 특성에 맞는 최적화된 가이드를 제공합니다:

    - default: 범용 AI를 위한 표준 가이드
    - codex: OpenAI Codex CLI (Less is More 원칙)
    - claude: Claude Code (체계적 워크플로우)
    - gemini: Gemini CLI (대화형 상호작용)
    - copilot: GitHub Copilot (IDE 통합형)
    """
    from pyhub_office_automation.cli.ai_guides import AIAssistant, AIGuideGenerator, OutputFormat

    # 지원되는 AI 타입 검증
    supported_ai_types = [e.value for e in AIAssistant]
    supported_formats = [e.value for e in OutputFormat]

    ai_type_lower = ai_type.lower()
    format_lower = format.lower()

    # AI 타입 검증
    if ai_type_lower not in supported_ai_types:
        typer.echo(f"Error: '{ai_type}'는 지원하지 않는 AI 타입입니다.")
        typer.echo(f"지원 타입: {', '.join(supported_ai_types)}")
        typer.echo("기본값 'default'를 사용하거나 지원되는 타입을 선택하세요.")
        raise typer.Exit(1)

    # 출력 형식 검증
    if format_lower not in supported_formats:
        typer.echo(f"Error: '{format}'는 지원하지 않는 출력 형식입니다.")
        typer.echo(f"지원 형식: {', '.join(supported_formats)}")
        format_lower = "json"  # 기본값으로 폴백

    # 언어 검증
    if lang.lower() not in ["ko", "en"]:
        typer.echo(f"Warning: '{lang}'는 지원하지 않는 언어입니다. 'ko'를 사용합니다.")
        lang = "ko"

    try:
        # 가이드 생성
        ai_enum = AIAssistant(ai_type_lower)
        format_enum = OutputFormat(format_lower)

        generator = AIGuideGenerator()
        guide = generator.generate(ai_type=ai_enum, verbose=verbose, lang=lang.lower())

        # 출력 형식에 따라 처리
        if format_enum == OutputFormat.json:
            try:
                json_output = json.dumps(guide, ensure_ascii=False, indent=2)
                typer.echo(json_output)
            except UnicodeEncodeError:
                json_output = json.dumps(guide, ensure_ascii=True, indent=2)
                typer.echo(json_output)
        elif format_enum == OutputFormat.markdown:
            markdown_output = generator.to_markdown(guide)
            typer.echo(markdown_output)
        else:  # text
            text_output = generator.to_text(guide)
            typer.echo(text_output)

    except ValueError as e:
        # Enum 변환 실패 등 값 관련 오류
        typer.echo(f"Error: 잘못된 값이 입력되었습니다: {e}")
        raise typer.Exit(1)
    except ImportError as e:
        # 모듈 import 실패
        typer.echo(f"Error: 필요한 모듈을 불러올 수 없습니다: {e}")
        typer.echo("패키지 설치 상태를 확인하세요: oa info")
        raise typer.Exit(1)
    except (OSError, IOError) as e:
        # 파일 읽기/쓰기 오류
        typer.echo(f"Error: 파일 작업 중 오류가 발생했습니다: {e}")
        raise typer.Exit(1)
    except KeyError as e:
        # 딕셔너리 키 누락 등
        typer.echo(f"Error: 가이드 데이터 구조에 문제가 있습니다: {e}")
        typer.echo("개발팀에 문의하세요.")
        raise typer.Exit(1)


excel_app = typer.Typer(help="Excel 자동화 명령어들", no_args_is_help=True)
hwp_app = typer.Typer(help="HWP 자동화 명령어들 (Windows 전용)", no_args_is_help=True)
ppt_app = typer.Typer(help="PowerPoint 자동화 명령어들", no_args_is_help=True)
email_app = typer.Typer(help="AI 기반 이메일 자동화 명령어들", no_args_is_help=True)

# Rich 콘솔 - UTF-8 인코딩 안전성 확보
try:
    # Windows 환경에서 UTF-8 출력 보장
    console = Console(force_terminal=True, force_jupyter=False, legacy_windows=False, width=None)  # 자동 감지
except Exception:
    # fallback to basic console
    console = Console(legacy_windows=True)

# Excel 명령어 등록 (단계적 테스트)
# Range Commands
excel_app.command("range-read")(range_read)
excel_app.command("range-write")(range_write)
excel_app.command("range-convert")(range_convert)

# Data Commands (Issue #39)
excel_app.command("data-analyze")(data_analyze)
excel_app.command("data-transform")(data_transform)

# Workbook Commands
excel_app.command("workbook-list")(workbook_list)
excel_app.command("workbook-open")(workbook_open)
excel_app.command("workbook-create")(workbook_create)
excel_app.command("workbook-info")(workbook_info)
excel_app.command("metadata-generate")(metadata_generate)

# Sheet Commands
excel_app.command("sheet-activate")(sheet_activate)
excel_app.command("sheet-add")(sheet_add)
excel_app.command("sheet-delete")(sheet_delete)
excel_app.command("sheet-rename")(sheet_rename)

# Table Commands
excel_app.command("table-create")(table_create)
excel_app.command("table-list")(table_list)
excel_app.command("table-read")(table_read)
excel_app.command("table-sort")(table_sort)
excel_app.command("table-sort-clear")(table_sort_clear)
excel_app.command("table-sort-info")(table_sort_info)
excel_app.command("table-write")(table_write)
excel_app.command("table-analyze")(table_analyze)

# Chart Commands
excel_app.command("chart-add")(chart_add)
excel_app.command("chart-configure")(chart_configure)
excel_app.command("chart-delete")(chart_delete)
excel_app.command("chart-export")(chart_export)
excel_app.command("chart-list")(chart_list)
excel_app.command("chart-pivot-create")(chart_pivot_create)
excel_app.command("chart-position")(chart_position)

# Pivot Commands
excel_app.command("pivot-configure")(pivot_configure)
excel_app.command("pivot-create")(pivot_create)
excel_app.command("pivot-delete")(pivot_delete)
excel_app.command("pivot-list")(pivot_list)
excel_app.command("pivot-refresh")(pivot_refresh)

# Shape Commands (이제 Typer로 전환 완료)
excel_app.command("shape-add")(shape_add)
excel_app.command("shape-delete")(shape_delete)
excel_app.command("shape-format")(shape_format)
excel_app.command("shape-group")(shape_group)
excel_app.command("shape-list")(shape_list)
excel_app.command("textbox-add")(textbox_add)

# Slicer Commands (이제 Typer로 전환 완료)
excel_app.command("slicer-add")(slicer_add)
excel_app.command("slicer-connect")(slicer_connect)
excel_app.command("slicer-list")(slicer_list)
excel_app.command("slicer-position")(slicer_position)

# Shell Command (Issue #85)
from pyhub_office_automation.shell.excel_shell import excel_shell

excel_app.command("shell")(excel_shell)


# Excel list command
@excel_app.command("list")
def excel_list_temp(
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택"),
):
    """Excel 자동화 명령어 목록 출력"""
    commands = [
        # Workbook Commands
        {"name": "workbook-list", "description": "열린 Excel 워크북 목록 조회", "category": "workbook"},
        {"name": "workbook-open", "description": "Excel 워크북 열기", "category": "workbook"},
        {"name": "workbook-create", "description": "새 Excel 워크북 생성", "category": "workbook"},
        {"name": "workbook-info", "description": "워크북 정보 조회", "category": "workbook"},
        {"name": "metadata-generate", "description": "워크북 전체 Excel Table 메타데이터 자동 생성", "category": "workbook"},
        # Sheet Commands
        {"name": "sheet-activate", "description": "시트 활성화", "category": "sheet"},
        {"name": "sheet-add", "description": "새 시트 추가", "category": "sheet"},
        {"name": "sheet-delete", "description": "시트 삭제", "category": "sheet"},
        {"name": "sheet-rename", "description": "시트 이름 변경", "category": "sheet"},
        # Range Commands
        {"name": "range-read", "description": "셀 범위 데이터 읽기", "category": "range"},
        {"name": "range-write", "description": "셀 범위에 데이터 쓰기", "category": "range"},
        {"name": "range-convert", "description": "셀 범위 데이터 형식 변환 (문자열 → 숫자)", "category": "range"},
        # Data Commands (Issue #39)
        {"name": "data-analyze", "description": "피벗테이블용 데이터 구조 분석", "category": "data"},
        {"name": "data-transform", "description": "피벗테이블용 형식으로 데이터 변환", "category": "data"},
        # Table Commands
        {"name": "table-create", "description": "기존 범위를 Excel Table로 변환", "category": "table"},
        {"name": "table-list", "description": "워크북의 Excel Table 목록 조회", "category": "table"},
        {"name": "table-read", "description": "테이블 데이터를 DataFrame으로 읽기", "category": "table"},
        {"name": "table-sort", "description": "테이블 정렬 적용", "category": "table"},
        {"name": "table-sort-clear", "description": "테이블 정렬 해제", "category": "table"},
        {"name": "table-sort-info", "description": "테이블 정렬 상태 확인", "category": "table"},
        {"name": "table-write", "description": "DataFrame을 Excel 테이블로 쓰기 (선택적 Table 생성)", "category": "table"},
        {"name": "table-analyze", "description": "Excel Table 분석 및 메타데이터 자동 생성", "category": "table"},
        # Chart Commands
        {"name": "chart-add", "description": "차트 추가", "category": "chart"},
        {"name": "chart-configure", "description": "차트 설정", "category": "chart"},
        {"name": "chart-delete", "description": "차트 삭제", "category": "chart"},
        {"name": "chart-export", "description": "차트 내보내기", "category": "chart"},
        {"name": "chart-list", "description": "차트 목록 조회", "category": "chart"},
        {"name": "chart-pivot-create", "description": "피벗 차트 생성", "category": "chart"},
        {"name": "chart-position", "description": "차트 위치 설정", "category": "chart"},
        # Pivot Commands
        {"name": "pivot-configure", "description": "피벗테이블 설정", "category": "pivot"},
        {"name": "pivot-create", "description": "피벗테이블 생성", "category": "pivot"},
        {"name": "pivot-delete", "description": "피벗테이블 삭제", "category": "pivot"},
        {"name": "pivot-list", "description": "피벗테이블 목록 조회", "category": "pivot"},
        {"name": "pivot-refresh", "description": "피벗테이블 새로고침", "category": "pivot"},
        # Shape Commands
        {"name": "shape-add", "description": "도형 추가", "category": "shape"},
        {"name": "shape-delete", "description": "도형 삭제", "category": "shape"},
        {"name": "shape-format", "description": "도형 서식 설정", "category": "shape"},
        {"name": "shape-group", "description": "도형 그룹화", "category": "shape"},
        {"name": "shape-list", "description": "도형 목록 조회", "category": "shape"},
        {"name": "textbox-add", "description": "텍스트 상자 추가", "category": "shape"},
        # Slicer Commands
        {"name": "slicer-add", "description": "슬라이서 추가", "category": "slicer"},
        {"name": "slicer-connect", "description": "슬라이서 연결", "category": "slicer"},
        {"name": "slicer-list", "description": "슬라이서 목록 조회", "category": "slicer"},
        {"name": "slicer-position", "description": "슬라이서 위치 설정", "category": "slicer"},
    ]

    excel_data = {
        "category": "excel",
        "description": "Excel 자동화 명령어들 (xlwings 기반)",
        "platform_requirement": "Windows (전체 기능) / macOS (제한적)",
        "commands": commands,
        "total_commands": len(commands),
        "package_version": get_version(),
    }

    if output_format == "json":
        try:
            json_output = json.dumps(excel_data, ensure_ascii=False, indent=2)
            typer.echo(json_output)
        except UnicodeEncodeError:
            json_output = json.dumps(excel_data, ensure_ascii=True, indent=2)
            typer.echo(json_output)
    else:
        console.print("=== Excel 자동화 명령어 목록 ===", style="bold green")
        console.print(f"Platform: {excel_data['platform_requirement']}")
        console.print(f"Total: {excel_data['total_commands']} commands")
        console.print()

        categories = {}
        for cmd in commands:
            category = cmd["category"]
            if category not in categories:
                categories[category] = []
            categories[category].append(cmd)

        for category, cmds in categories.items():
            console.print(f"[bold blue]{category.upper()} Commands:[/bold blue]")
            for cmd in cmds:
                console.print(f"  • oa excel {cmd['name']}")
                console.print(f"    {cmd['description']}")
            console.print()

        console.print("📚 [bold yellow]더 자세한 사용 지침은 다음 명령어를 참고하세요:[/bold yellow]", style="bright_yellow")
        console.print("   [bold cyan]oa llm-guide[/bold cyan] - AI 에이전트를 위한 상세 가이드")
        console.print("   [bold cyan]oa excel <command> --help[/bold cyan] - 특정 명령어 도움말")
        console.print()


# Email 명령어 등록
email_app.command("send")(email_send)
email_app.add_typer(accounts_app, name="accounts")


# Email list command
@email_app.command("list")
def email_list(
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택"),
):
    """Email 자동화 명령어 목록 출력"""
    commands = [
        {"name": "send", "description": "AI 기반 이메일 생성 및 발송", "category": "core"},
    ]

    email_data = {
        "category": "email",
        "description": "AI 기반 이메일 자동화 명령어들",
        "platform_requirement": "Windows (Outlook COM) / 크로스 플랫폼 (SMTP)",
        "commands": commands,
        "total_commands": len(commands),
        "ai_providers": ["claude", "codex", "gemini", "openai", "anthropic"],
        "package_version": get_version(),
    }

    if output_format == "json":
        try:
            json_output = json.dumps(email_data, ensure_ascii=False, indent=2)
            typer.echo(json_output)
        except UnicodeEncodeError:
            json_output = json.dumps(email_data, ensure_ascii=True, indent=2)
            typer.echo(json_output)
    else:
        console.print("=== Email 자동화 명령어 목록 ===", style="bold green")
        console.print(f"Platform: {email_data['platform_requirement']}")
        console.print(f"Total: {email_data['total_commands']} commands")
        console.print(f"AI Providers: {', '.join(email_data['ai_providers'])}")
        console.print()

        for cmd in commands:
            console.print(f"  • oa email {cmd['name']}")
            console.print(f"    {cmd['description']}")
        console.print()

        console.print("📚 [bold yellow]AI 프롬프트 사용 예시:[/bold yellow]")
        console.print('   [bold cyan]oa email send --to "user@example.com" --prompt "회의 일정 변경 안내"[/bold cyan]')
        console.print('   [bold cyan]oa email send --to "team@company.com" --prompt "프로젝트 완료 보고"[/bold cyan]')


# 서브 앱을 메인 앱에 등록
app.add_typer(excel_app, name="excel")
app.add_typer(hwp_app, name="hwp")
app.add_typer(ppt_app, name="ppt")
app.add_typer(email_app, name="email")
app.add_typer(ai_setup_app, name="ai-setup")


@app.command()
def info(
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택"),
):
    """패키지 정보 및 설치 상태 출력"""
    try:
        version_info = get_version_info()
        dependencies = check_dependencies()

        info_data = {
            "package": "pyhub-office-automation",
            "version": version_info,
            "platform": sys.platform,
            "python_version": sys.version,
            "dependencies": dependencies,
            "status": "installed",
        }

        if output_format == "json":
            try:
                json_output = json.dumps(info_data, ensure_ascii=False, indent=2)
                typer.echo(json_output)
            except UnicodeEncodeError:
                json_output = json.dumps(info_data, ensure_ascii=True, indent=2)
                typer.echo(json_output)
        else:
            console.print(f"Package: {info_data['package']}", style="bold green")
            console.print(f"Version: {info_data['version']['version']}")
            console.print(f"Platform: {info_data['platform']}")
            console.print(f"Python: {info_data['python_version']}")
            console.print("Dependencies:", style="bold")

            table = Table()
            table.add_column("Package", style="cyan")
            table.add_column("Status", style="green")
            table.add_column("Version")

            for dep, status in info_data["dependencies"].items():
                status_mark = "✓" if status["available"] else "✗"
                version = status["version"] or "Not installed"
                table.add_row(dep, status_mark, version)

            console.print(table)

    except Exception as e:
        error_data = {"error": str(e), "command": "info", "version": get_version()}
        typer.echo(json.dumps(error_data, ensure_ascii=False, indent=2), err=True)
        raise typer.Exit(1)


@app.command()
def install_guide(
    output_format: str = typer.Option("text", "--format", help="출력 형식 선택"),
):
    """설치 가이드 출력"""
    guide_steps = [
        {
            "step": 1,
            "title": "Python 설치",
            "description": "Python 3.13 이상을 설치하세요",
            "url": "https://www.python.org/downloads/",
            "command": None,
        },
        {
            "step": 2,
            "title": "패키지 설치",
            "description": "pip를 사용하여 pyhub-office-automation을 설치하세요",
            "command": "pip install pyhub-office-automation",
        },
        {"step": 3, "title": "설치 확인", "description": "oa 명령어가 정상 동작하는지 확인하세요", "command": "oa info"},
        {
            "step": 4,
            "title": "Excel 사용 시 (선택사항)",
            "description": "Microsoft Excel이 설치되어 있어야 합니다",
            "note": "xlwings는 Excel이 설치된 환경에서 동작합니다",
        },
        {
            "step": 5,
            "title": "HWP 사용 시 (선택사항, Windows 전용)",
            "description": "한글(HWP) 프로그램이 설치되어 있어야 합니다",
            "note": "pyhwpx는 Windows COM을 통해 HWP와 연동됩니다",
        },
    ]

    guide_data = {
        "title": "pyhub-office-automation 설치 가이드",
        "version": get_version(),
        "platform_requirement": "Windows 10/11 (추천)",
        "python_requirement": "Python 3.13+",
        "steps": guide_steps,
    }

    if output_format == "json":
        try:
            # JSON 출력 시 한글 인코딩 문제 해결
            json_output = json.dumps(guide_data, ensure_ascii=False, indent=2)
            typer.echo(json_output)
        except UnicodeEncodeError:
            # Windows 콘솔 인코딩 문제 시 ensure_ascii=True로 폴백
            json_output = json.dumps(guide_data, ensure_ascii=True, indent=2)
            typer.echo(json_output)
    else:
        console.print(f"=== {guide_data['title']} ===", style="bold blue")
        console.print(f"Version: {guide_data['version']}")
        console.print(f"Platform: {guide_data['platform_requirement']}")
        console.print(f"Python: {guide_data['python_requirement']}")
        console.print()

        for step in guide_steps:
            console.print(f"Step {step['step']}: {step['title']}", style="bold yellow")
            console.print(f"  {step['description']}")
            if step.get("command"):
                console.print(f"  Command: [green]{step['command']}[/green]")
            if step.get("url"):
                console.print(f"  URL: [blue]{step['url']}[/blue]")
            if step.get("note"):
                console.print(f"  Note: [dim]{step['note']}[/dim]")
            console.print()


@hwp_app.command("list")
def hwp_list(
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택"),
):
    """HWP 자동화 명령어 목록 출력"""
    commands = [
        {"name": "export", "description": "HWP 파일을 HTML로 변환", "version": "1.0.0", "status": "available"},
        {"name": "open", "description": "HWP 파일 열기", "version": "1.0.0", "status": "planned"},
        {"name": "save", "description": "HWP 파일 저장", "version": "1.0.0", "status": "planned"},
    ]

    hwp_data = {
        "category": "hwp",
        "description": "HWP 자동화 명령어들 (pyhwpx 기반, Windows 전용)",
        "platform_requirement": "Windows + HWP 프로그램 설치 필요",
        "commands": commands,
        "total_commands": len(commands),
        "package_version": get_version(),
    }

    if output_format == "json":
        try:
            json_output = json.dumps(hwp_data, ensure_ascii=False, indent=2)
            typer.echo(json_output)
        except UnicodeEncodeError:
            json_output = json.dumps(hwp_data, ensure_ascii=True, indent=2)
            typer.echo(json_output)
    else:
        console.print("=== HWP 자동화 명령어 목록 ===", style="bold blue")
        console.print(f"Platform: {hwp_data['platform_requirement']}")
        console.print(f"Total: {hwp_data['total_commands']} commands")
        console.print()

        for cmd in commands:
            status_mark = "✓" if cmd["status"] == "available" else "○"
            console.print(f"  {status_mark} oa hwp {cmd['name']}")
            console.print(f"     {cmd['description']} (v{cmd['version']})")


# PowerPoint 명령어 등록
# Presentation 관리
ppt_app.command("presentation-create")(presentation_create)
ppt_app.command("presentation-open")(presentation_open)
ppt_app.command("presentation-save")(presentation_save)
ppt_app.command("presentation-list")(presentation_list)
ppt_app.command("presentation-info")(presentation_info)

# Slide 관리 (Issue #76)
ppt_app.command("slide-list")(slide_list)
ppt_app.command("slide-add")(slide_add)
ppt_app.command("slide-delete")(slide_delete)
ppt_app.command("slide-duplicate")(slide_duplicate)
ppt_app.command("slide-copy")(slide_copy)
ppt_app.command("slide-reorder")(slide_reorder)

# Content 추가 (Issue #77)
ppt_app.command("content-add-text")(content_add_text)
ppt_app.command("content-add-image")(content_add_image)
ppt_app.command("content-add-shape")(content_add_shape)
ppt_app.command("content-add-table")(content_add_table)

# Advanced Content 추가 및 업데이트 (Issue #78)
ppt_app.command("content-add-chart")(content_add_chart)
ppt_app.command("content-add-video")(content_add_video)
ppt_app.command("content-add-smartart")(content_add_smartart)
ppt_app.command("content-update")(content_update)
ppt_app.command("content-add-excel-chart")(content_add_excel_chart)
ppt_app.command("content-add-audio")(content_add_audio)
ppt_app.command("content-add-equation")(content_add_equation)

# Layout & Theme 관리 (Issue #79)
ppt_app.command("layout-list")(layout_list)
ppt_app.command("layout-apply")(layout_apply)
ppt_app.command("template-apply")(template_apply)
ppt_app.command("theme-apply")(theme_apply)

# Export 기능 (Issue #80)
ppt_app.command("export-pdf")(export_pdf)
ppt_app.command("export-images")(export_images)
ppt_app.command("export-notes")(export_notes)

# Windows Advanced Features (Issue #84 Phase 3)
ppt_app.command("slideshow-start")(slideshow_start)
ppt_app.command("slideshow-control")(slideshow_control)
ppt_app.command("run-macro")(run_macro)
ppt_app.command("animation-add")(animation_add)

# Shell Command (Issue #85 Phase 5)
from pyhub_office_automation.shell.ppt_shell import ppt_shell

ppt_app.command("shell")(ppt_shell)


@ppt_app.command("list")
def ppt_list(
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택"),
):
    """PowerPoint 자동화 명령어 목록 출력"""
    commands = [
        {"name": "presentation-create", "description": "새 프레젠테이션 생성", "category": "presentation"},
        {"name": "presentation-open", "description": "프레젠테이션 파일 열기", "category": "presentation"},
        {"name": "presentation-save", "description": "프레젠테이션 저장", "category": "presentation"},
        {"name": "presentation-list", "description": "열린 프레젠테이션 목록 (Windows COM 전용)", "category": "presentation"},
        {"name": "presentation-info", "description": "프레젠테이션 상세 정보", "category": "presentation"},
        {"name": "slide-list", "description": "슬라이드 목록 조회", "category": "slide"},
        {"name": "slide-add", "description": "새 슬라이드 추가", "category": "slide"},
        {"name": "slide-delete", "description": "슬라이드 삭제", "category": "slide"},
        {"name": "slide-duplicate", "description": "슬라이드 복제", "category": "slide"},
        {"name": "slide-copy", "description": "슬라이드 복사", "category": "slide"},
        {"name": "slide-reorder", "description": "슬라이드 순서 변경", "category": "slide"},
        {"name": "content-add-text", "description": "슬라이드에 텍스트 추가", "category": "content"},
        {"name": "content-add-image", "description": "슬라이드에 이미지 추가", "category": "content"},
        {"name": "content-add-shape", "description": "슬라이드에 도형 추가", "category": "content"},
        {"name": "content-add-table", "description": "슬라이드에 표 추가", "category": "content"},
        {"name": "content-add-chart", "description": "슬라이드에 차트 추가", "category": "content"},
        {"name": "content-add-video", "description": "슬라이드에 비디오 추가", "category": "content"},
        {"name": "content-add-smartart", "description": "슬라이드에 SmartArt 추가", "category": "content"},
        {"name": "content-update", "description": "슬라이드 콘텐츠 수정", "category": "content"},
        {"name": "layout-list", "description": "사용 가능한 레이아웃 목록", "category": "layout"},
        {"name": "layout-apply", "description": "슬라이드에 레이아웃 적용", "category": "layout"},
        {"name": "template-apply", "description": "템플릿 적용", "category": "theme"},
        {"name": "theme-apply", "description": "테마 적용", "category": "theme"},
        {"name": "slideshow-start", "description": "슬라이드쇼 시작 (Windows COM 전용)", "category": "advanced"},
        {"name": "slideshow-control", "description": "슬라이드쇼 제어 (Windows COM 전용)", "category": "advanced"},
        {"name": "run-macro", "description": "VBA 매크로 실행 (Windows COM 전용)", "category": "advanced"},
        {"name": "animation-add", "description": "애니메이션 효과 추가 (Windows COM 전용)", "category": "advanced"},
    ]

    ppt_data = {
        "category": "ppt",
        "description": "PowerPoint 자동화 명령어들 (python-pptx 기반)",
        "platform_requirement": "Windows (전체 기능) / macOS (85%+ 기능)",
        "commands": commands,
        "total_commands": len(commands),
        "package_version": get_version(),
        "status": "layout_theme_ready",
        "note": "Issue #79 완료 - Layout & Theme 관리 명령어 4개 구현됨",
    }

    if output_format == "json":
        try:
            json_output = json.dumps(ppt_data, ensure_ascii=False, indent=2)
            typer.echo(json_output)
        except UnicodeEncodeError:
            json_output = json.dumps(ppt_data, ensure_ascii=True, indent=2)
            typer.echo(json_output)
    else:
        console.print("=== PowerPoint 자동화 명령어 목록 ===", style="bold green")
        console.print(f"Platform: {ppt_data['platform_requirement']}")
        console.print(f"Status: {ppt_data['status']}")
        console.print(f"Total: {ppt_data['total_commands']} commands")
        console.print()

        categories = {}
        for cmd in commands:
            category = cmd["category"]
            if category not in categories:
                categories[category] = []
            categories[category].append(cmd)

        for category, cmds in categories.items():
            console.print(f"[bold cyan]{category.upper()} Commands:[/bold cyan]")
            for cmd in cmds:
                console.print(f"  • oa ppt {cmd['name']}")
                console.print(f"    {cmd['description']}")
            console.print()

        console.print("📚 [bold yellow]더 자세한 사용법:[/bold yellow]")
        console.print("   [bold cyan]oa ppt <command> --help[/bold cyan] - 특정 명령어 도움말")
        console.print()


@hwp_app.command("export")
def hwp_export_command(
    file_path: str = typer.Option(..., "--file-path", help="변환할 HWP 파일의 절대 경로"),
    format_type: str = typer.Option("html", "--format", help="출력 형식 (현재 html만 지원)"),
    output_file: Optional[str] = typer.Option(None, "--output-file", help="HTML 저장 경로 (선택, 미지정시 표준출력)"),
    encoding: str = typer.Option("utf-8", "--encoding", help="출력 인코딩 (기본값: utf-8)"),
    include_css: bool = typer.Option(
        False, "--include-css/--no-include-css", help="CSS 스타일 포함 여부 (기본값: False, 모든 CSS 제거)"
    ),
    include_images: bool = typer.Option(
        False, "--include-images/--no-include-images", help="이미지 포함 여부 (기본값: False, Base64 인코딩으로 포함)"
    ),
    temp_cleanup: bool = typer.Option(True, "--temp-cleanup/--no-temp-cleanup", help="임시 파일 자동 정리 (기본값: True)"),
    output_format: str = typer.Option("json", "--output-format", help="응답 출력 형식 (json)"),
):
    """HWP 파일을 HTML 형식으로 변환"""
    hwp_export(
        file_path=file_path,
        format_type=format_type,
        output_file=output_file,
        encoding=encoding,
        include_css=include_css,
        include_images=include_images,
        temp_cleanup=temp_cleanup,
        output_format=output_format,
    )


@app.command()
def get_help(
    category: str = typer.Argument(..., help="명령어 카테고리 (excel, hwp)"),
    command_name: str = typer.Argument(..., help="명령어 이름"),
    output_format: str = typer.Option("text", "--format", help="출력 형식 선택"),
):
    """특정 명령어의 도움말 조회"""
    help_data = {
        "category": category,
        "command": command_name,
        "help": f"oa {category} {command_name} 명령어 도움말 (구현 예정)",
        "usage": f"oa {category} {command_name} [OPTIONS]",
        "status": "planned",
        "version": get_version(),
    }

    if output_format == "json":
        try:
            json_output = json.dumps(help_data, ensure_ascii=False, indent=2)
            typer.echo(json_output)
        except UnicodeEncodeError:
            json_output = json.dumps(help_data, ensure_ascii=True, indent=2)
            typer.echo(json_output)
    else:
        console.print(f"Command: oa {category} {command_name}", style="bold")
        console.print(f"Usage: {help_data['usage']}")
        console.print(f"Status: {help_data['status']}")
        console.print()
        console.print(help_data["help"])


def check_dependencies():
    """의존성 패키지 설치 상태 확인"""
    dependencies = {}

    # xlwings 확인
    try:
        import xlwings

        dependencies["xlwings"] = {"available": True, "version": xlwings.__version__}
    except ImportError:
        dependencies["xlwings"] = {"available": False, "version": None}

    # pyhwpx 확인 (Windows 전용)
    try:
        import pyhwpx

        dependencies["pyhwpx"] = {"available": True, "version": getattr(pyhwpx, "__version__", "unknown")}
    except ImportError:
        dependencies["pyhwpx"] = {"available": False, "version": None}

    # pandas 확인
    try:
        import pandas

        dependencies["pandas"] = {"available": True, "version": pandas.__version__}
    except ImportError:
        dependencies["pandas"] = {"available": False, "version": None}

    # python-pptx 확인
    try:
        import pptx

        dependencies["python-pptx"] = {"available": True, "version": pptx.__version__}
    except ImportError:
        dependencies["python-pptx"] = {"available": False, "version": None}

    return dependencies


def main():
    """메인 엔트리포인트"""
    app()


if __name__ == "__main__":
    main()
