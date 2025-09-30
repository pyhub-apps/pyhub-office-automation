"""
PowerPoint 오디오 추가 명령어 (COM-First)
슬라이드에 오디오 파일을 삽입합니다.
"""

import json
from pathlib import Path
from typing import Optional

import typer

from pyhub_office_automation.version import get_version

from .utils import (
    PowerPointBackend,
    create_error_response,
    create_success_response,
    get_or_open_presentation,
    get_powerpoint_backend,
    normalize_path,
    validate_slide_number,
)

# 지원되는 오디오 형식
SUPPORTED_AUDIO_FORMATS = [".mp3", ".wav", ".m4a", ".wma", ".aac", ".flac", ".ogg"]


def content_add_audio(
    slide_number: int = typer.Option(..., "--slide-number", help="오디오를 추가할 슬라이드 번호 (1부터 시작)"),
    audio_path: str = typer.Option(..., "--audio-path", help="추가할 오디오 파일 경로"),
    left: Optional[float] = typer.Option(None, "--left", help="오디오 아이콘 왼쪽 위치 (인치)"),
    top: Optional[float] = typer.Option(None, "--top", help="오디오 아이콘 상단 위치 (인치)"),
    width: Optional[float] = typer.Option(1.0, "--width", help="오디오 아이콘 너비 (인치, 기본값: 1.0)"),
    height: Optional[float] = typer.Option(1.0, "--height", help="오디오 아이콘 높이 (인치, 기본값: 1.0)"),
    center: bool = typer.Option(False, "--center", help="슬라이드 중앙에 배치 (--left, --top 무시)"),
    autoplay: bool = typer.Option(False, "--autoplay", help="슬라이드 표시 시 자동 재생"),
    loop: bool = typer.Option(False, "--loop", help="반복 재생"),
    hide_icon: bool = typer.Option(False, "--hide-icon", help="재생 중 아이콘 숨기기"),
    file_path: Optional[str] = typer.Option(None, "--file-path", help="PowerPoint 파일 경로"),
    presentation_name: Optional[str] = typer.Option(None, "--presentation-name", help="열려있는 프레젠테이션 이름 (COM 전용)"),
    backend: str = typer.Option("auto", "--backend", help="백엔드 선택 (auto/com/python-pptx)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 슬라이드에 오디오를 추가합니다.

    COM-First: Windows에서는 COM 백엔드 우선, python-pptx는 fallback

    **백엔드 선택**:
    - auto (기본): 자동으로 최적 백엔드 선택 (Windows COM 우선)
    - com: Windows COM 강제 사용 (완전한 기능)
    - python-pptx: python-pptx 강제 사용 (제한적 기능)

    **COM 백엔드 (Windows) - 완전한 기능!**:
    - ✅ 오디오 삽입 및 위치/크기 조정
    - ✅ Shapes.AddMediaObject2() 사용
    - ✅ 자동 재생, 반복, 아이콘 숨김 설정 가능
    - 열려있는 프레젠테이션에서 직접 작업

    **python-pptx 백엔드**:
    - ⚠️ 파일 저장 필수 (--file-path 필수)
    - 오디오 삽입 가능
    - 제한적 재생 옵션

    **지원 오디오 형식**:
      MP3, WAV, M4A, WMA, AAC, FLAC, OGG

    **위치 지정**:
      --center: 슬라이드 중앙에 배치
      --left, --top: 특정 위치에 배치

    예제:
        # COM 백엔드 (활성 프레젠테이션, 중앙 배치)
        oa ppt content-add-audio --slide-number 2 --audio-path "bgm.mp3" --center

        # COM 백엔드 (자동 재생 + 반복)
        oa ppt content-add-audio --slide-number 3 --audio-path "narration.wav" --left 1 --top 1 --autoplay --loop

        # COM 백엔드 (아이콘 숨김)
        oa ppt content-add-audio --slide-number 4 --audio-path "sfx.m4a" --center --hide-icon --presentation-name "demo.pptx"

        # python-pptx 백엔드
        oa ppt content-add-audio --slide-number 5 --audio-path "audio.mp3" --file-path "report.pptx" --backend python-pptx
    """

    try:
        # 입력 검증
        if not center and (left is None or top is None):
            result = create_error_response(
                command="content-add-audio",
                error="--center를 사용하지 않는 경우 --left와 --top을 모두 지정해야 합니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 오디오 파일 검증
        normalized_audio_path = normalize_path(audio_path)
        audio_file = Path(normalized_audio_path).resolve()

        if not audio_file.exists():
            result = create_error_response(
                command="content-add-audio",
                error=f"오디오 파일을 찾을 수 없습니다: {audio_path}",
                error_type="FileNotFoundError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 오디오 형식 검증
        audio_ext = audio_file.suffix.lower()
        if audio_ext not in SUPPORTED_AUDIO_FORMATS:
            supported_str = ", ".join(SUPPORTED_AUDIO_FORMATS)
            result = create_error_response(
                command="content-add-audio",
                error=f"지원하지 않는 오디오 형식: {audio_ext}. 지원 형식: {supported_str}",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드 결정
        try:
            selected_backend = get_powerpoint_backend(force_backend=backend if backend != "auto" else None)
        except (ValueError, RuntimeError) as e:
            result = create_error_response(
                command="content-add-audio",
                error=str(e),
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 오디오 파일 크기
        audio_size_mb = audio_file.stat().st_size / (1024 * 1024)

        # 프레젠테이션 가져오기
        try:
            backend_inst, prs = get_or_open_presentation(
                file_path=file_path,
                presentation_name=presentation_name,
                backend=selected_backend,
            )
        except Exception as e:
            result = create_error_response(
                command="content-add-audio",
                error=f"프레젠테이션을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드별 처리
        if selected_backend == PowerPointBackend.COM.value:
            # COM 백엔드: 완전한 오디오 추가 기능
            try:
                total_slides = prs.Slides.Count

                # 슬라이드 번호 검증 (COM은 1-based)
                if slide_number < 1 or slide_number > total_slides:
                    result = create_error_response(
                        command="content-add-audio",
                        error=f"슬라이드 번호가 범위를 벗어났습니다: {slide_number} (1-{total_slides})",
                        error_type="ValueError",
                    )
                    typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                    raise typer.Exit(1)

                slide = prs.Slides(slide_number)

                # 위치 계산
                if center:
                    # 슬라이드 크기 가져오기 (포인트 단위)
                    slide_width_pt = prs.PageSetup.SlideWidth
                    slide_height_pt = prs.PageSetup.SlideHeight

                    # 포인트를 인치로 변환
                    slide_width_in = slide_width_pt / 72
                    slide_height_in = slide_height_pt / 72

                    # 중앙 배치 위치 계산
                    final_left = (slide_width_in - width) / 2
                    final_top = (slide_height_in - height) / 2
                else:
                    final_left = left
                    final_top = top

                # 인치를 포인트로 변환
                left_pt = final_left * 72
                top_pt = final_top * 72
                width_pt = width * 72
                height_pt = height * 72

                # 오디오 추가
                # AddMediaObject2(FileName, LinkToFile, SaveWithDocument, Left, Top, Width, Height)
                shape = slide.Shapes.AddMediaObject2(
                    str(audio_file),
                    0,  # LinkToFile = msoFalse (파일 포함)
                    -1,  # SaveWithDocument = msoTrue
                    left_pt,
                    top_pt,
                    width_pt,
                    height_pt,
                )

                # 재생 설정 (AnimationSettings 사용)
                if hasattr(shape, "AnimationSettings"):
                    anim_settings = shape.AnimationSettings

                    # 자동 재생 설정
                    if autoplay:
                        anim_settings.PlaySettings.PlayOnEntry = True

                    # 반복 재생 설정
                    if loop:
                        anim_settings.PlaySettings.LoopUntilStopped = True

                    # 아이콘 숨김 설정
                    if hide_icon:
                        anim_settings.PlaySettings.HideWhileNotPlaying = True

                # 성공 응답
                result_data = {
                    "backend": "com",
                    "slide_number": slide_number,
                    "audio_file": str(audio_file),
                    "audio_format": audio_ext,
                    "audio_size_mb": round(audio_size_mb, 2),
                    "position": {
                        "left": round(final_left, 2),
                        "top": round(final_top, 2),
                        "width": width,
                        "height": height,
                    },
                    "centered": center,
                    "autoplay": autoplay,
                    "loop": loop,
                    "hide_icon": hide_icon,
                }

                message = f"오디오 추가 완료 (COM): 슬라이드 {slide_number}"
                if autoplay:
                    message += ", 자동 재생"
                if loop:
                    message += ", 반복"

            except Exception as e:
                result = create_error_response(
                    command="content-add-audio",
                    error=f"오디오 추가 실패: {str(e)}",
                    error_type=type(e).__name__,
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

        else:
            # python-pptx 백엔드
            if not file_path:
                result = create_error_response(
                    command="content-add-audio",
                    error="python-pptx 백엔드는 --file-path 옵션이 필수입니다",
                    error_type="ValueError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

            # 슬라이드 번호 검증
            slide_idx = validate_slide_number(slide_number, len(prs.slides))
            slide = prs.slides[slide_idx]

            # 위치 계산
            if center:
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                # EMU를 인치로 변환
                slide_width_in = slide_width / 914400
                slide_height_in = slide_height / 914400

                final_left = (slide_width_in - width) / 2
                final_top = (slide_height_in - height) / 2
            else:
                final_left = left
                final_top = top

            # 오디오 추가 (python-pptx는 add_movie 사용)
            from pptx.util import Inches

            movie = slide.shapes.add_movie(
                str(audio_file),
                Inches(final_left),
                Inches(final_top),
                Inches(width),
                Inches(height),
            )

            # 저장
            pptx_path = Path(normalize_path(file_path)).resolve()
            prs.save(str(pptx_path))

            # 결과 데이터
            result_data = {
                "backend": "python-pptx",
                "file": str(pptx_path),
                "file_name": pptx_path.name,
                "slide_number": slide_number,
                "audio_file": str(audio_file),
                "audio_format": audio_ext,
                "audio_size_mb": round(audio_size_mb, 2),
                "position": {
                    "left": round(final_left, 2),
                    "top": round(final_top, 2),
                    "width": width,
                    "height": height,
                },
                "centered": center,
                "note": "python-pptx는 자동재생/반복 설정을 지원하지 않습니다",
            }

            message = f"오디오 추가 완료 (python-pptx): 슬라이드 {slide_number}"

        # 성공 응답
        response = create_success_response(
            data=result_data,
            command="content-add-audio",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✅ {message}")
            typer.echo(f"📍 슬라이드: {slide_number}")
            typer.echo(f"🎵 오디오: {audio_file.name}")
            typer.echo(f"📦 형식: {audio_ext.upper()}")
            typer.echo(f"💾 크기: {result_data['audio_size_mb']} MB")
            typer.echo(f"📐 위치: {result_data['position']['left']}in × {result_data['position']['top']}in")
            typer.echo(f"📏 크기: {width}in × {height}in")
            if selected_backend == PowerPointBackend.COM.value:
                if autoplay:
                    typer.echo("▶️ 자동 재생: 켜짐")
                if loop:
                    typer.echo("🔁 반복: 켜짐")
                if hide_icon:
                    typer.echo("👁️ 아이콘: 재생 중 숨김")

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="content-add-audio",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)
    finally:
        # python-pptx는 자동 정리, COM은 유지
        pass


if __name__ == "__main__":
    typer.run(content_add_audio)
