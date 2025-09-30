"""
PowerPoint 차트 추가 명령어 (COM-First)
Excel 데이터 또는 CSV 파일로부터 차트를 생성합니다.
"""

import json
from pathlib import Path
from typing import Optional

import typer

from pyhub_office_automation.version import get_version

from .utils import (
    PowerPointBackend,
    create_error_response,
    create_success_response,
    get_or_open_presentation,
    get_powerpoint_backend,
    normalize_path,
    validate_slide_number,
)

# COM 차트 타입 상수 매핑 (XlChartType)
CHART_TYPE_MAP_COM = {
    "column": 51,  # xlColumnClustered
    "bar": 57,  # xlBarClustered
    "line": 4,  # xlLine
    "pie": 5,  # xlPie
    "area": 1,  # xlArea
    "scatter": -4169,  # xlXYScatter
    "doughnut": -4120,  # xlDoughnut
}


def content_add_chart(
    slide_number: int = typer.Option(..., "--slide-number", help="차트를 추가할 슬라이드 번호 (1부터 시작)"),
    chart_type: str = typer.Option(..., "--chart-type", help="차트 타입 (column/bar/line/pie/area/scatter/doughnut)"),
    csv_data: Optional[str] = typer.Option(None, "--csv-data", help="CSV 파일 경로"),
    excel_data: Optional[str] = typer.Option(None, "--excel-data", help="Excel 데이터 참조 (예: data.xlsx!A1:C10)"),
    left: Optional[float] = typer.Option(None, "--left", help="차트 왼쪽 위치 (인치)"),
    top: Optional[float] = typer.Option(None, "--top", help="차트 상단 위치 (인치)"),
    width: Optional[float] = typer.Option(6.0, "--width", help="차트 너비 (인치, 기본값: 6.0)"),
    height: Optional[float] = typer.Option(4.5, "--height", help="차트 높이 (인치, 기본값: 4.5)"),
    title: Optional[str] = typer.Option(None, "--title", help="차트 제목"),
    center: bool = typer.Option(False, "--center", help="슬라이드 중앙에 배치 (--left, --top 무시)"),
    show_legend: bool = typer.Option(True, "--show-legend/--no-legend", help="범례 표시 여부 (기본값: 표시)"),
    file_path: Optional[str] = typer.Option(None, "--file-path", help="PowerPoint 파일 경로"),
    presentation_name: Optional[str] = typer.Option(None, "--presentation-name", help="열려있는 프레젠테이션 이름 (COM 전용)"),
    backend: str = typer.Option("auto", "--backend", help="백엔드 선택 (auto/com/python-pptx)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 슬라이드에 데이터 기반 차트를 추가합니다.

    COM-First: Windows에서는 COM 백엔드 우선, python-pptx는 fallback

    **백엔드 선택**:
    - auto (기본): 자동으로 최적 백엔드 선택 (Windows COM 우선)
    - com: Windows COM 강제 사용 (완전한 기능)
    - python-pptx: python-pptx 강제 사용 (제한적 기능)

    **COM 백엔드 (Windows) - 완전한 기능!**:
    - ✅ Shapes.AddChart() 사용
    - ✅ 네이티브 Excel 통합
    - ✅ 더 많은 차트 타입 지원 가능
    - 열려있는 프레젠테이션에서 직접 작업

    **python-pptx 백엔드**:
    - ⚠️ 파일 저장 필수 (--file-path 필수)
    - 기본 차트 기능만 지원

    데이터 소스 (둘 중 하나만 지정):
      --csv-data: CSV 파일 경로
      --excel-data: Excel 참조 (예: "data.xlsx!A1:C10" 또는 "data.xlsx!Sheet1!A1:C10")

    차트 타입:
      column, bar, line, pie, area, scatter, doughnut

    위치 지정:
      --center: 슬라이드 중앙에 배치
      --left, --top: 특정 위치에 배치

    예제:
        # COM 백엔드 (활성 프레젠테이션)
        oa ppt content-add-chart --slide-number 2 --chart-type column --csv-data "sales.csv" --center --title "판매 현황"

        # COM 백엔드 (특정 프레젠테이션)
        oa ppt content-add-chart --slide-number 3 --chart-type pie --excel-data "data.xlsx!A1:C10" --left 1 --top 2 --presentation-name "report.pptx"

        # python-pptx 백엔드
        oa ppt content-add-chart --slide-number 4 --chart-type bar --csv-data "sales.csv" --file-path "report.pptx" --backend python-pptx
    """

    try:
        # 입력 검증
        if not csv_data and not excel_data:
            result = create_error_response(
                command="content-add-chart",
                error="--csv-data 또는 --excel-data 중 하나는 반드시 지정해야 합니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        if csv_data and excel_data:
            result = create_error_response(
                command="content-add-chart",
                error="--csv-data와 --excel-data는 동시에 사용할 수 없습니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        if not center and (left is None or top is None):
            result = create_error_response(
                command="content-add-chart",
                error="--center를 사용하지 않는 경우 --left와 --top을 모두 지정해야 합니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 차트 타입 검증
        chart_type_lower = chart_type.lower()
        if chart_type_lower not in CHART_TYPE_MAP_COM:
            available_types = ", ".join(CHART_TYPE_MAP_COM.keys())
            result = create_error_response(
                command="content-add-chart",
                error=f"지원하지 않는 차트 타입: {chart_type}\n사용 가능: {available_types}",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 데이터 파일 검증
        data_source_path = None
        if csv_data:
            normalized_csv_path = normalize_path(csv_data)
            data_source_path = Path(normalized_csv_path).resolve()
            data_source_name = str(data_source_path)
        else:
            # Excel 데이터 참조 파싱
            from .utils import parse_excel_range

            excel_ref = parse_excel_range(excel_data)
            normalized_excel_path = normalize_path(excel_ref["file_path"])
            data_source_path = Path(normalized_excel_path).resolve()
            data_source_name = excel_data

        if not data_source_path.exists():
            result = create_error_response(
                command="content-add-chart",
                error=f"데이터 파일을 찾을 수 없습니다: {data_source_path}",
                error_type="FileNotFoundError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드 결정
        try:
            selected_backend = get_powerpoint_backend(force_backend=backend if backend != "auto" else None)
        except (ValueError, RuntimeError) as e:
            result = create_error_response(
                command="content-add-chart",
                error=str(e),
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 데이터 로드 (pandas DataFrame)
        import pandas as pd

        if csv_data:
            df = pd.read_csv(str(data_source_path))
        else:
            # Excel 데이터 로드
            from .utils import load_data_from_excel

            df = load_data_from_excel(
                file_path=excel_ref["file_path"], sheet_name=excel_ref["sheet"], range_addr=excel_ref["range"]
            )

        # 데이터 검증
        if df is None or df.empty:
            result = create_error_response(
                command="content-add-chart",
                error="데이터가 비어있습니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        if len(df.columns) < 2:
            result = create_error_response(
                command="content-add-chart",
                error=f"차트를 생성하려면 최소 2개의 열이 필요합니다 (현재: {len(df.columns)}개)",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 프레젠테이션 가져오기
        try:
            backend_inst, prs = get_or_open_presentation(
                file_path=file_path,
                presentation_name=presentation_name,
                backend=selected_backend,
            )
        except Exception as e:
            result = create_error_response(
                command="content-add-chart",
                error=f"프레젠테이션을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드별 처리
        if selected_backend == PowerPointBackend.COM.value:
            # COM 백엔드: 완전한 차트 추가 기능
            try:
                total_slides = prs.Slides.Count

                # 슬라이드 번호 검증 (COM은 1-based)
                if slide_number < 1 or slide_number > total_slides:
                    result = create_error_response(
                        command="content-add-chart",
                        error=f"슬라이드 번호가 범위를 벗어났습니다: {slide_number} (1-{total_slides})",
                        error_type="ValueError",
                    )
                    typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                    raise typer.Exit(1)

                slide = prs.Slides(slide_number)

                # 위치 계산
                if center:
                    # 슬라이드 크기 가져오기 (포인트 단위)
                    slide_width_pt = prs.PageSetup.SlideWidth
                    slide_height_pt = prs.PageSetup.SlideHeight

                    # 포인트를 인치로 변환
                    slide_width_in = slide_width_pt / 72
                    slide_height_in = slide_height_pt / 72

                    # 중앙 배치 위치 계산
                    final_left = (slide_width_in - width) / 2
                    final_top = (slide_height_in - height) / 2
                else:
                    final_left = left
                    final_top = top

                # 인치를 포인트로 변환
                left_pt = final_left * 72
                top_pt = final_top * 72
                width_pt = width * 72
                height_pt = height * 72

                # 차트 추가
                chart_type_const = CHART_TYPE_MAP_COM[chart_type_lower]
                shape = slide.Shapes.AddChart2(
                    -1,  # Style (-1 = default)
                    chart_type_const,
                    left_pt,
                    top_pt,
                    width_pt,
                    height_pt,
                )

                chart = shape.Chart

                # 차트 데이터 설정
                chart.ChartData.Activate()
                workbook = chart.ChartData.Workbook
                worksheet = workbook.Worksheets(1)

                # DataFrame을 2D 배열로 변환 (헤더 포함)
                data_with_header = [df.columns.tolist()] + df.values.tolist()
                num_rows = len(data_with_header)
                num_cols = len(data_with_header[0])

                # Excel 범위에 데이터 쓰기
                cell_range = worksheet.Range(worksheet.Cells(1, 1), worksheet.Cells(num_rows, num_cols))

                # 2D 리스트를 1D 튜플로 변환 (COM에서 요구하는 형식)
                flat_data = tuple(tuple(row) for row in data_with_header)
                cell_range.Value = flat_data

                # 차트 데이터 범위 설정
                chart.SetSourceData(worksheet.Range(worksheet.Cells(1, 1), worksheet.Cells(num_rows, num_cols)))

                # 차트 제목 설정
                if title:
                    chart.HasTitle = True
                    chart.ChartTitle.Text = title

                # 범례 설정
                chart.HasLegend = show_legend

                # 워크북 닫기
                workbook.Close()

                # 성공 응답
                result_data = {
                    "backend": "com",
                    "slide_number": slide_number,
                    "chart_type": chart_type_lower,
                    "data_source": data_source_name,
                    "data_shape": {"rows": len(df), "columns": len(df.columns)},
                    "series_count": len(df.columns) - 1,
                    "position": {
                        "left": round(final_left, 2),
                        "top": round(final_top, 2),
                        "width": width,
                        "height": height,
                    },
                    "centered": center,
                    "has_title": title is not None,
                    "has_legend": show_legend,
                }

                if title:
                    result_data["title"] = title

                message = f"차트 추가 완료 (COM): 슬라이드 {slide_number}, {chart_type_lower} 차트"
                if title:
                    message += f" (제목: {title})"

            except Exception as e:
                result = create_error_response(
                    command="content-add-chart",
                    error=f"차트 추가 실패: {str(e)}",
                    error_type=type(e).__name__,
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

        else:
            # python-pptx 백엔드
            if not file_path:
                result = create_error_response(
                    command="content-add-chart",
                    error="python-pptx 백엔드는 --file-path 옵션이 필수입니다",
                    error_type="ValueError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

            # 슬라이드 번호 검증
            slide_idx = validate_slide_number(slide_number, len(prs.slides))
            slide = prs.slides[slide_idx]

            # 위치 계산
            if center:
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                # EMU를 인치로 변환
                slide_width_in = slide_width / 914400
                slide_height_in = slide_height / 914400

                final_left = (slide_width_in - width) / 2
                final_top = (slide_height_in - height) / 2
            else:
                final_left = left
                final_top = top

            # ChartData 생성
            from pptx.chart.data import ChartData
            from pptx.enum.chart import XL_CHART_TYPE
            from pptx.util import Inches

            chart_data = ChartData()

            # 카테고리 설정 (첫 번째 열)
            categories = df.iloc[:, 0].tolist()
            chart_data.categories = categories

            # 시리즈 추가 (두 번째 열부터)
            for col in df.columns[1:]:
                chart_data.add_series(col, df[col].tolist())

            # 차트 타입 매핑 (python-pptx)
            chart_type_map_pptx = {
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
                "area": XL_CHART_TYPE.AREA,
                "scatter": XL_CHART_TYPE.XY_SCATTER,
                "doughnut": XL_CHART_TYPE.DOUGHNUT,
            }

            chart_type_const = chart_type_map_pptx[chart_type_lower]

            # 차트 추가
            graphic_frame = slide.shapes.add_chart(
                chart_type_const, Inches(final_left), Inches(final_top), Inches(width), Inches(height), chart_data
            )

            chart = graphic_frame.chart

            # 차트 제목 설정
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title

            # 범례 설정
            chart.has_legend = show_legend

            # 저장
            pptx_path = Path(normalize_path(file_path)).resolve()
            prs.save(str(pptx_path))

            # 결과 데이터
            result_data = {
                "backend": "python-pptx",
                "file": str(pptx_path),
                "file_name": pptx_path.name,
                "slide_number": slide_number,
                "chart_type": chart_type_lower,
                "data_source": data_source_name,
                "data_shape": {"rows": len(df), "columns": len(df.columns)},
                "series_count": len(df.columns) - 1,
                "position": {
                    "left": round(final_left, 2),
                    "top": round(final_top, 2),
                    "width": width,
                    "height": height,
                },
                "centered": center,
                "has_title": title is not None,
                "has_legend": show_legend,
            }

            if title:
                result_data["title"] = title

            message = f"차트 추가 완료 (python-pptx): 슬라이드 {slide_number}, {chart_type_lower} 차트"
            if title:
                message += f" (제목: {title})"

        # 성공 응답
        response = create_success_response(
            data=result_data,
            command="content-add-chart",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✅ {message}")
            typer.echo(f"📍 슬라이드: {slide_number}")
            typer.echo(f"📊 차트 타입: {chart_type_lower}")
            typer.echo(f"📈 데이터 소스: {data_source_name}")
            typer.echo(f"📏 데이터 크기: {result_data['data_shape']['rows']}행 × {result_data['data_shape']['columns']}열")
            typer.echo(f"📐 위치: {result_data['position']['left']}in × {result_data['position']['top']}in")
            typer.echo(f"📏 크기: {width}in × {height}in")
            if title:
                typer.echo(f"🏷️ 제목: {title}")
            typer.echo(f"📊 시리즈 개수: {result_data['series_count']}")
            typer.echo(f"📖 범례: {'표시' if show_legend else '숨김'}")

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="content-add-chart",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)
    finally:
        # python-pptx는 자동 정리, COM은 유지
        pass


if __name__ == "__main__":
    typer.run(content_add_chart)
