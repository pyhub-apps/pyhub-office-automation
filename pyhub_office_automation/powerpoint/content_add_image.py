"""
PowerPoint 이미지 추가 명령어 (COM-First)
슬라이드에 이미지를 추가합니다.
"""

import json
from pathlib import Path
from typing import Optional

import typer
from PIL import Image

from pyhub_office_automation.version import get_version

from .utils import (
    PowerPointBackend,
    calculate_aspect_ratio_size,
    create_error_response,
    create_success_response,
    get_or_open_presentation,
    get_powerpoint_backend,
    normalize_path,
    validate_slide_number,
)


def content_add_image(
    slide_number: int = typer.Option(..., "--slide-number", help="이미지를 추가할 슬라이드 번호 (1부터 시작)"),
    image_path: str = typer.Option(..., "--image-path", help="추가할 이미지 파일 경로"),
    left: Optional[float] = typer.Option(None, "--left", help="이미지 왼쪽 위치 (인치)"),
    top: Optional[float] = typer.Option(None, "--top", help="이미지 상단 위치 (인치)"),
    width: Optional[float] = typer.Option(None, "--width", help="이미지 너비 (인치) - 미지정시 원본 비율 유지"),
    height: Optional[float] = typer.Option(None, "--height", help="이미지 높이 (인치) - 미지정시 원본 비율 유지"),
    center: bool = typer.Option(False, "--center", help="슬라이드 중앙에 배치 (--left, --top 무시)"),
    file_path: Optional[str] = typer.Option(None, "--file-path", help="PowerPoint 파일 경로"),
    presentation_name: Optional[str] = typer.Option(None, "--presentation-name", help="열려있는 프레젠테이션 이름 (COM 전용)"),
    backend: str = typer.Option("auto", "--backend", help="백엔드 선택 (auto/com/python-pptx)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 슬라이드에 이미지를 추가합니다.

    COM-First: Windows에서는 COM 백엔드 우선, python-pptx는 fallback

    **백엔드 선택**:
    - auto (기본): 자동으로 최적 백엔드 선택 (Windows COM 우선)
    - com: Windows COM 강제 사용 (완전한 기능)
    - python-pptx: python-pptx 강제 사용 (제한적 기능)

    **COM 백엔드 (Windows) - 완전한 기능!**:
    - ✅ 이미지 추가 및 위치/크기 조정
    - Shapes.AddPicture() 사용
    - 열려있는 프레젠테이션에서 직접 작업

    **python-pptx 백엔드**:
    - ⚠️ 파일 저장 필수 (--file-path 필수)
    - 이미지 추가 가능

    **위치 지정 방법**:
      --center: 슬라이드 중앙에 배치
      --left, --top: 특정 위치에 배치

    **크기 지정**:
      --width, --height: 둘 다 지정하면 지정된 크기로
      --width만 지정: 너비 기준으로 비율 유지하여 높이 자동 계산
      --height만 지정: 높이 기준으로 비율 유지하여 너비 자동 계산
      미지정: 원본 크기 (DPI 기준 인치 변환)

    예제:
        # COM 백엔드 (활성 프레젠테이션, 중앙 배치)
        oa ppt content-add-image --slide-number 1 --image-path "logo.png" --center --width 2

        # COM 백엔드 (특정 프레젠테이션, 위치 지정)
        oa ppt content-add-image --slide-number 2 --image-path "chart.jpg" --left 1 --top 2 --height 3 --presentation-name "report.pptx"

        # python-pptx 백엔드
        oa ppt content-add-image --slide-number 3 --image-path "photo.png" --center --file-path "report.pptx" --backend python-pptx
    """
    backend_inst = None

    try:
        # 입력 검증
        if not center and (left is None or top is None):
            result = create_error_response(
                command="content-add-image",
                error="--center를 사용하지 않는 경우 --left와 --top을 모두 지정해야 합니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드 결정
        try:
            selected_backend = get_powerpoint_backend(force_backend=backend if backend != "auto" else None)
        except (ValueError, RuntimeError) as e:
            result = create_error_response(
                command="content-add-image",
                error=str(e),
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 이미지 경로 검증
        normalized_image_path = normalize_path(image_path)
        img_path = Path(normalized_image_path).resolve()

        if not img_path.exists():
            result = create_error_response(
                command="content-add-image",
                error=f"이미지 파일을 찾을 수 없습니다: {image_path}",
                error_type="FileNotFoundError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 이미지 정보 읽기 (PIL 사용)
        try:
            with Image.open(str(img_path)) as img:
                original_width_px, original_height_px = img.size
                dpi = img.info.get("dpi", (96, 96))
                if isinstance(dpi, tuple):
                    dpi_x, dpi_y = dpi
                else:
                    dpi_x = dpi_y = dpi

                # 픽셀을 인치로 변환
                original_width_in = original_width_px / dpi_x
                original_height_in = original_height_px / dpi_y
        except Exception as e:
            result = create_error_response(
                command="content-add-image",
                error=f"이미지 파일을 읽을 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 크기 계산 (aspect ratio 유지)
        if width is None and height is None:
            final_width = original_width_in
            final_height = original_height_in
        elif width is not None and height is not None:
            final_width = width
            final_height = height
        elif width is not None:
            final_width, final_height = calculate_aspect_ratio_size(original_width_in, original_height_in, target_width=width)
        else:
            final_width, final_height = calculate_aspect_ratio_size(
                original_width_in, original_height_in, target_height=height
            )

        # 프레젠테이션 가져오기
        try:
            backend_inst, prs = get_or_open_presentation(
                file_path=file_path,
                presentation_name=presentation_name,
                backend=selected_backend,
            )
        except Exception as e:
            result = create_error_response(
                command="content-add-image",
                error=f"프레젠테이션을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드별 처리
        if selected_backend == PowerPointBackend.COM.value:
            # COM 백엔드: 완전한 이미지 추가 기능
            try:
                total_slides = prs.Slides.Count

                # 슬라이드 번호 검증 (COM은 1-based)
                if slide_number < 1 or slide_number > total_slides:
                    result = create_error_response(
                        command="content-add-image",
                        error=f"슬라이드 번호가 범위를 벗어났습니다: {slide_number} (1-{total_slides})",
                        error_type="ValueError",
                    )
                    typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                    raise typer.Exit(1)

                slide = prs.Slides(slide_number)

                # 위치 계산
                if center:
                    # 슬라이드 크기 가져오기 (포인트 단위)
                    slide_width_pt = prs.PageSetup.SlideWidth
                    slide_height_pt = prs.PageSetup.SlideHeight

                    # 포인트를 인치로 변환 (1 inch = 72 points)
                    slide_width_in = slide_width_pt / 72
                    slide_height_in = slide_height_pt / 72

                    # 중앙 배치 위치 계산
                    final_left = (slide_width_in - final_width) / 2
                    final_top = (slide_height_in - final_height) / 2
                else:
                    final_left = left
                    final_top = top

                # 인치를 포인트로 변환 (COM API는 포인트 사용)
                left_pt = final_left * 72
                top_pt = final_top * 72
                width_pt = final_width * 72
                height_pt = final_height * 72

                # 이미지 추가
                picture = slide.Shapes.AddPicture(
                    FileName=str(img_path),
                    LinkToFile=0,  # msoFalse - 파일 링크 없음
                    SaveWithDocument=-1,  # msoTrue - 문서에 포함
                    Left=left_pt,
                    Top=top_pt,
                    Width=width_pt,
                    Height=height_pt,
                )

                # 성공 응답
                result_data = {
                    "backend": "com",
                    "slide_number": slide_number,
                    "image_file": str(img_path),
                    "image_name": img_path.name,
                    "position": {
                        "left": round(final_left, 2),
                        "top": round(final_top, 2),
                        "width": round(final_width, 2),
                        "height": round(final_height, 2),
                    },
                    "original_size": {
                        "width_px": original_width_px,
                        "height_px": original_height_px,
                        "width_in": round(original_width_in, 2),
                        "height_in": round(original_height_in, 2),
                    },
                    "centered": center,
                }

                message = f"이미지 추가 완료 (COM): 슬라이드 {slide_number}"
                if center:
                    message += ", 중앙 배치"
                else:
                    message += f", 위치 {final_left}in × {final_top}in"

            except Exception as e:
                result = create_error_response(
                    command="content-add-image",
                    error=f"이미지 추가 실패: {str(e)}",
                    error_type=type(e).__name__,
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

        else:
            # python-pptx 백엔드
            if not file_path:
                result = create_error_response(
                    command="content-add-image",
                    error="python-pptx 백엔드는 --file-path 옵션이 필수입니다",
                    error_type="ValueError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

            # 슬라이드 번호 검증
            slide_idx = validate_slide_number(slide_number, len(prs.slides))
            slide = prs.slides[slide_idx]

            # 위치 계산
            if center:
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                # EMU를 인치로 변환 (1 inch = 914400 EMU)
                slide_width_in = slide_width / 914400
                slide_height_in = slide_height / 914400

                final_left = (slide_width_in - final_width) / 2
                final_top = (slide_height_in - final_height) / 2
            else:
                final_left = left
                final_top = top

            # 이미지 추가
            from pptx.util import Inches

            picture = slide.shapes.add_picture(
                str(img_path), Inches(final_left), Inches(final_top), width=Inches(final_width), height=Inches(final_height)
            )

            # 저장
            pptx_path = Path(normalize_path(file_path)).resolve()
            prs.save(str(pptx_path))

            # 결과 데이터
            result_data = {
                "backend": "python-pptx",
                "file": str(pptx_path),
                "file_name": pptx_path.name,
                "slide_number": slide_number,
                "image_file": str(img_path),
                "image_name": img_path.name,
                "position": {
                    "left": round(final_left, 2),
                    "top": round(final_top, 2),
                    "width": round(final_width, 2),
                    "height": round(final_height, 2),
                },
                "original_size": {
                    "width_px": original_width_px,
                    "height_px": original_height_px,
                    "width_in": round(original_width_in, 2),
                    "height_in": round(original_height_in, 2),
                },
                "centered": center,
            }

            message = f"이미지 추가 완료 (python-pptx): 슬라이드 {slide_number}"
            if center:
                message += ", 중앙 배치"
            else:
                message += f", 위치 {final_left}in × {final_top}in"

        # 성공 응답
        response = create_success_response(
            data=result_data,
            command="content-add-image",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✅ {message}")
            typer.echo(f"📍 슬라이드: {slide_number}")
            typer.echo(f"🖼️ 이미지: {result_data['image_name']}")
            typer.echo(f"📐 위치: {result_data['position']['left']}in × {result_data['position']['top']}in")
            typer.echo(f"📏 크기: {result_data['position']['width']}in × {result_data['position']['height']}in")
            typer.echo(
                f"🎨 원본: {result_data['original_size']['width_px']}px × {result_data['original_size']['height_px']}px"
            )

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="content-add-image",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)
    finally:
        # python-pptx는 자동 정리, COM은 유지
        pass


if __name__ == "__main__":
    typer.run(content_add_image)
