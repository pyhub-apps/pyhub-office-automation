"""
PowerPoint 도형 추가 명령어 (COM-First)
슬라이드에 도형을 추가합니다.
"""

import json
from pathlib import Path
from typing import Optional

import typer

from pyhub_office_automation.version import get_version

from .utils import (
    PowerPointBackend,
    create_error_response,
    create_success_response,
    get_or_open_presentation,
    get_powerpoint_backend,
    normalize_path,
    parse_color,
    validate_slide_number,
)

# COM 도형 타입 상수 매핑
SHAPE_TYPE_MAP_COM = {
    "rectangle": 1,  # msoShapeRectangle
    "rounded-rectangle": 5,  # msoShapeRoundedRectangle
    "ellipse": 9,  # msoShapeOval
    "arrow-right": 33,  # msoShapeRightArrow
    "arrow-left": 34,  # msoShapeLeftArrow
    "arrow-up": 35,  # msoShapeUpArrow
    "arrow-down": 36,  # msoShapeDownArrow
    "star": 12,  # msoShape5pointStar
    "pentagon": 56,  # msoShapePentagon
    "hexagon": 10,  # msoShapeHexagon
}


def content_add_shape(
    slide_number: int = typer.Option(..., "--slide-number", help="도형을 추가할 슬라이드 번호 (1부터 시작)"),
    shape_type: str = typer.Option(..., "--shape-type", help="도형 유형 (rectangle/ellipse/star/arrow 등)"),
    left: float = typer.Option(..., "--left", help="도형 왼쪽 위치 (인치)"),
    top: float = typer.Option(..., "--top", help="도형 상단 위치 (인치)"),
    width: float = typer.Option(..., "--width", help="도형 너비 (인치)"),
    height: float = typer.Option(..., "--height", help="도형 높이 (인치)"),
    fill_color: Optional[str] = typer.Option(None, "--fill-color", help="채우기 색상 (색상명 또는 #RGB/#RRGGBB)"),
    line_color: Optional[str] = typer.Option(None, "--line-color", help="테두리 색상 (색상명 또는 #RGB/#RRGGBB)"),
    line_width: Optional[float] = typer.Option(None, "--line-width", help="테두리 두께 (포인트)"),
    text: Optional[str] = typer.Option(None, "--text", help="도형 내부 텍스트"),
    file_path: Optional[str] = typer.Option(None, "--file-path", help="PowerPoint 파일 경로"),
    presentation_name: Optional[str] = typer.Option(None, "--presentation-name", help="열려있는 프레젠테이션 이름 (COM 전용)"),
    backend: str = typer.Option("auto", "--backend", help="백엔드 선택 (auto/com/python-pptx)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 슬라이드에 도형을 추가합니다.

    COM-First: Windows에서는 COM 백엔드 우선, python-pptx는 fallback

    **백엔드 선택**:
    - auto (기본): 자동으로 최적 백엔드 선택 (Windows COM 우선)
    - com: Windows COM 강제 사용 (완전한 기능)
    - python-pptx: python-pptx 강제 사용 (제한적 기능)

    **COM 백엔드 (Windows) - 완전한 기능!**:
    - ✅ 도형 생성 및 스타일 설정
    - Shapes.AddShape() 사용
    - 열려있는 프레젠테이션에서 직접 작업

    **python-pptx 백엔드**:
    - ⚠️ 파일 저장 필수 (--file-path 필수)
    - 도형 생성 및 스타일 설정 가능

    **지원 도형 유형**:
      - rectangle: 사각형
      - rounded-rectangle: 둥근 사각형
      - ellipse: 타원
      - arrow-right: 오른쪽 화살표
      - arrow-left: 왼쪽 화살표
      - arrow-up: 위쪽 화살표
      - arrow-down: 아래쪽 화살표
      - star: 별
      - pentagon: 오각형
      - hexagon: 육각형

    예제:
        # COM 백엔드 (활성 프레젠테이션)
        oa ppt content-add-shape --slide-number 1 --shape-type rectangle --left 1 --top 2 --width 3 --height 2 --fill-color blue

        # COM 백엔드 (특정 프레젠테이션)
        oa ppt content-add-shape --slide-number 2 --shape-type star --left 2 --top 3 --width 1.5 --height 1.5 --fill-color "#FFD700" --text "중요" --presentation-name "report.pptx"

        # python-pptx 백엔드
        oa ppt content-add-shape --slide-number 3 --shape-type arrow-right --left 1 --top 1 --width 2 --height 1 --fill-color red --line-color black --line-width 2 --file-path "report.pptx" --backend python-pptx
    """

    try:
        # 도형 유형 검증
        if shape_type not in SHAPE_TYPE_MAP_COM:
            result = create_error_response(
                command="content-add-shape",
                error=f"지원하지 않는 도형 유형: {shape_type}. 사용 가능: {', '.join(SHAPE_TYPE_MAP_COM.keys())}",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드 결정
        try:
            selected_backend = get_powerpoint_backend(force_backend=backend if backend != "auto" else None)
        except (ValueError, RuntimeError) as e:
            result = create_error_response(
                command="content-add-shape",
                error=str(e),
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 프레젠테이션 가져오기
        try:
            backend_inst, prs = get_or_open_presentation(
                file_path=file_path,
                presentation_name=presentation_name,
                backend=selected_backend,
            )
        except Exception as e:
            result = create_error_response(
                command="content-add-shape",
                error=f"프레젠테이션을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드별 처리
        if selected_backend == PowerPointBackend.COM.value:
            # COM 백엔드: 완전한 도형 추가 기능
            try:
                total_slides = prs.Slides.Count

                # 슬라이드 번호 검증 (COM은 1-based)
                if slide_number < 1 or slide_number > total_slides:
                    result = create_error_response(
                        command="content-add-shape",
                        error=f"슬라이드 번호가 범위를 벗어났습니다: {slide_number} (1-{total_slides})",
                        error_type="ValueError",
                    )
                    typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                    raise typer.Exit(1)

                slide = prs.Slides(slide_number)

                # 인치를 포인트로 변환 (COM API는 포인트 사용)
                left_pt = left * 72
                top_pt = top * 72
                width_pt = width * 72
                height_pt = height * 72

                # 도형 추가
                mso_shape_type = SHAPE_TYPE_MAP_COM[shape_type]
                shape = slide.Shapes.AddShape(Type=mso_shape_type, Left=left_pt, Top=top_pt, Width=width_pt, Height=height_pt)

                # 채우기 색상 설정
                if fill_color is not None:
                    rgb = parse_color(fill_color)
                    # RGB 튜플을 정수로 변환 (COM에서 사용하는 형식)
                    rgb_value = rgb.red + (rgb.green * 256) + (rgb.blue * 256 * 256)
                    shape.Fill.Solid()
                    shape.Fill.ForeColor.RGB = rgb_value

                # 테두리 색상 설정
                if line_color is not None:
                    rgb = parse_color(line_color)
                    rgb_value = rgb.red + (rgb.green * 256) + (rgb.blue * 256 * 256)
                    shape.Line.ForeColor.RGB = rgb_value

                # 테두리 두께 설정
                if line_width is not None:
                    shape.Line.Weight = line_width

                # 텍스트 추가
                if text is not None:
                    if shape.HasTextFrame:
                        shape.TextFrame.TextRange.Text = text

                # 성공 응답
                result_data = {
                    "backend": "com",
                    "slide_number": slide_number,
                    "shape_type": shape_type,
                    "position": {
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                    },
                }

                if fill_color:
                    result_data["fill_color"] = fill_color
                if line_color:
                    result_data["line_color"] = line_color
                if line_width:
                    result_data["line_width"] = line_width
                if text:
                    result_data["text"] = text

                message = f"도형 추가 완료 (COM): 슬라이드 {slide_number}, {shape_type}"

            except Exception as e:
                result = create_error_response(
                    command="content-add-shape",
                    error=f"도형 추가 실패: {str(e)}",
                    error_type=type(e).__name__,
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

        else:
            # python-pptx 백엔드
            if not file_path:
                result = create_error_response(
                    command="content-add-shape",
                    error="python-pptx 백엔드는 --file-path 옵션이 필수입니다",
                    error_type="ValueError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

            # python-pptx 도형 타입 매핑
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

            shape_type_map_pptx = {
                "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                "rounded-rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
                "arrow-right": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                "arrow-left": MSO_AUTO_SHAPE_TYPE.LEFT_ARROW,
                "arrow-up": MSO_AUTO_SHAPE_TYPE.UP_ARROW,
                "arrow-down": MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
                "star": MSO_AUTO_SHAPE_TYPE.STAR_5,
                "pentagon": MSO_AUTO_SHAPE_TYPE.PENTAGON,
                "hexagon": MSO_AUTO_SHAPE_TYPE.HEXAGON,
            }

            # 슬라이드 번호 검증
            slide_idx = validate_slide_number(slide_number, len(prs.slides))
            slide = prs.slides[slide_idx]

            # 도형 추가
            from pptx.util import Inches, Pt

            mso_shape_type = shape_type_map_pptx[shape_type]
            shape = slide.shapes.add_shape(mso_shape_type, Inches(left), Inches(top), Inches(width), Inches(height))

            # 채우기 색상 설정
            if fill_color is not None:
                rgb = parse_color(fill_color)
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb

            # 테두리 설정
            if line_color is not None:
                rgb = parse_color(line_color)
                shape.line.color.rgb = rgb

            if line_width is not None:
                shape.line.width = Pt(line_width)

            # 텍스트 추가
            if text is not None:
                if hasattr(shape, "text_frame"):
                    text_frame = shape.text_frame
                    text_frame.clear()
                    paragraph = text_frame.paragraphs[0]
                    run = paragraph.add_run()
                    run.text = text

            # 저장
            pptx_path = Path(normalize_path(file_path)).resolve()
            prs.save(str(pptx_path))

            # 결과 데이터
            result_data = {
                "backend": "python-pptx",
                "file": str(pptx_path),
                "file_name": pptx_path.name,
                "slide_number": slide_number,
                "shape_type": shape_type,
                "position": {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                },
            }

            if fill_color:
                result_data["fill_color"] = fill_color
            if line_color:
                result_data["line_color"] = line_color
            if line_width:
                result_data["line_width"] = line_width
            if text:
                result_data["text"] = text

            message = f"도형 추가 완료 (python-pptx): 슬라이드 {slide_number}, {shape_type}"

        # 성공 응답
        response = create_success_response(
            data=result_data,
            command="content-add-shape",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✅ {message}")
            typer.echo(f"📍 슬라이드: {slide_number}")
            typer.echo(f"🔷 도형: {shape_type}")
            typer.echo(f"📐 위치: {left}in × {top}in")
            typer.echo(f"📏 크기: {width}in × {height}in")
            if fill_color:
                typer.echo(f"🎨 채우기: {fill_color}")
            if line_color:
                typer.echo(f"✏️ 테두리: {line_color}")
            if text:
                typer.echo(f"📝 텍스트: {text}")

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="content-add-shape",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)
    finally:
        # python-pptx는 자동 정리, COM은 유지
        pass


if __name__ == "__main__":
    typer.run(content_add_shape)
