"""
PowerPoint SmartArt 추가 명령어
JSON 데이터로부터 SmartArt 스타일의 다이어그램을 생성합니다.
python-pptx는 네이티브 SmartArt를 지원하지 않으므로 도형 조합으로 구현합니다.
"""

import json
from pathlib import Path
from typing import List, Optional

import typer
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from pyhub_office_automation.version import get_version

from .utils import create_error_response, create_success_response, normalize_path, parse_color, validate_slide_number


def content_add_smartart(
    file_path: str = typer.Option(..., "--file-path", help="PowerPoint 파일 경로"),
    slide_number: int = typer.Option(..., "--slide-number", help="SmartArt를 추가할 슬라이드 번호 (1부터 시작)"),
    diagram_type: str = typer.Option(
        ..., "--diagram-type", help="다이어그램 타입 (process/hierarchy/list/cycle/relationship)"
    ),
    data_file: str = typer.Option(..., "--data-file", help="데이터 JSON 파일 경로"),
    left: Optional[float] = typer.Option(1.0, "--left", help="다이어그램 왼쪽 위치 (인치, 기본값: 1.0)"),
    top: Optional[float] = typer.Option(2.0, "--top", help="다이어그램 상단 위치 (인치, 기본값: 2.0)"),
    width: Optional[float] = typer.Option(8.0, "--width", help="다이어그램 너비 (인치, 기본값: 8.0)"),
    height: Optional[float] = typer.Option(4.0, "--height", help="다이어그램 높이 (인치, 기본값: 4.0)"),
    shape_color: Optional[str] = typer.Option("blue", "--shape-color", help="도형 색상 (색상명 또는 #RGB)"),
    text_color: Optional[str] = typer.Option("white", "--text-color", help="텍스트 색상 (색상명 또는 #RGB)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 슬라이드에 SmartArt 스타일 다이어그램을 추가합니다.

    다이어그램 타입:
      process: 프로세스 흐름 (순차적 단계)
      hierarchy: 계층 구조 (조직도 등)
      list: 목록형 (항목 나열)
      cycle: 순환 구조 (반복 프로세스)
      relationship: 관계도 (연결 관계)

    데이터 JSON 형식:
      {
        "items": [
          {"text": "항목 1", "description": "설명 (선택)"},
          {"text": "항목 2", "description": "설명 (선택)"}
        ]
      }

    예제:
        oa ppt content-add-smartart --file-path "presentation.pptx" --slide-number 2 --diagram-type process --data-file "steps.json"
        oa ppt content-add-smartart --file-path "presentation.pptx" --slide-number 3 --diagram-type hierarchy --data-file "org.json" --shape-color "#4472C4"

    JSON 예제 (steps.json):
      {
        "items": [
          {"text": "계획", "description": "목표 설정"},
          {"text": "실행", "description": "작업 수행"},
          {"text": "평가", "description": "결과 검토"}
        ]
      }
    """
    try:
        # 다이어그램 타입 검증
        supported_types = ["process", "hierarchy", "list", "cycle", "relationship"]
        if diagram_type.lower() not in supported_types:
            raise ValueError(f"지원하지 않는 다이어그램 타입: {diagram_type}\n지원 타입: {', '.join(supported_types)}")

        # 파일 경로 정규화 및 존재 확인
        normalized_pptx_path = normalize_path(file_path)
        pptx_path = Path(normalized_pptx_path).resolve()

        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint 파일을 찾을 수 없습니다: {pptx_path}")

        normalized_data_path = normalize_path(data_file)
        json_path = Path(normalized_data_path).resolve()

        if not json_path.exists():
            raise FileNotFoundError(f"데이터 JSON 파일을 찾을 수 없습니다: {json_path}")

        # JSON 데이터 로드
        with open(str(json_path), "r", encoding="utf-8") as f:
            data = json.load(f)

        if "items" not in data or not isinstance(data["items"], list):
            raise ValueError("JSON 데이터에 'items' 배열이 필요합니다")

        items = data["items"]
        if len(items) == 0:
            raise ValueError("items 배열이 비어있습니다")

        # 색상 파싱
        fill_color = parse_color(shape_color)
        font_color = parse_color(text_color)

        # 프레젠테이션 열기
        prs = Presentation(str(pptx_path))

        # 슬라이드 번호 검증
        slide_idx = validate_slide_number(slide_number, len(prs.slides))
        slide = prs.slides[slide_idx]

        # 다이어그램 타입별 레이아웃 생성
        shapes_created = []

        if diagram_type.lower() == "process":
            shapes_created = create_process_diagram(slide, items, left, top, width, height, fill_color, font_color)

        elif diagram_type.lower() == "hierarchy":
            shapes_created = create_hierarchy_diagram(slide, items, left, top, width, height, fill_color, font_color)

        elif diagram_type.lower() == "list":
            shapes_created = create_list_diagram(slide, items, left, top, width, height, fill_color, font_color)

        elif diagram_type.lower() == "cycle":
            shapes_created = create_cycle_diagram(slide, items, left, top, width, height, fill_color, font_color)

        elif diagram_type.lower() == "relationship":
            shapes_created = create_relationship_diagram(slide, items, left, top, width, height, fill_color, font_color)

        # 저장
        prs.save(str(pptx_path))

        # 결과 데이터 구성
        result_data = {
            "file": str(pptx_path),
            "slide_number": slide_number,
            "diagram_type": diagram_type.lower(),
            "data_file": str(json_path),
            "item_count": len(items),
            "shapes_created": shapes_created,
            "position": {
                "left": left,
                "top": top,
                "width": width,
                "height": height,
            },
            "colors": {
                "shape": shape_color,
                "text": text_color,
            },
        }

        # 성공 응답
        message = f"슬라이드 {slide_number}에 {diagram_type} 다이어그램을 추가했습니다 ({len(items)}개 항목)"

        response = create_success_response(
            data=result_data,
            command="content-add-smartart",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✅ {message}")
            typer.echo(f"📄 파일: {pptx_path.name}")
            typer.echo(f"📍 슬라이드: {slide_number}")
            typer.echo(f"🎨 다이어그램 타입: {diagram_type}")
            typer.echo(f"📊 항목 개수: {len(items)}")
            typer.echo(f"📐 위치: {left}in × {top}in")
            typer.echo(f"📏 크기: {width}in × {height}in")
            typer.echo(f"🎨 도형 색상: {shape_color}")
            typer.echo(f"📝 생성된 도형: {shapes_created}개")

    except FileNotFoundError as e:
        error_response = create_error_response(e, "content-add-smartart")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ {str(e)}", err=True)
        raise typer.Exit(1)

    except ValueError as e:
        error_response = create_error_response(e, "content-add-smartart")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ {str(e)}", err=True)
        raise typer.Exit(1)

    except json.JSONDecodeError as e:
        error_response = create_error_response(f"JSON 파일 파싱 실패: {str(e)}", "content-add-smartart")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ JSON 파일 파싱 실패: {str(e)}", err=True)
        raise typer.Exit(1)

    except Exception as e:
        error_response = create_error_response(e, "content-add-smartart")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ 예기치 않은 오류: {str(e)}", err=True)
        raise typer.Exit(1)


# 다이어그램 생성 함수들


def create_process_diagram(
    slide, items: List[dict], left: float, top: float, width: float, height: float, fill_color, font_color
) -> int:
    """프로세스 흐름 다이어그램 생성 (왼쪽에서 오른쪽으로 화살표)"""
    num_items = len(items)
    box_width = (width / num_items) * 0.8  # 20% 간격
    spacing = (width - (box_width * num_items)) / (num_items + 1)

    shapes_count = 0

    for i, item in enumerate(items):
        # 박스 위치 계산
        box_left = left + spacing + (i * (box_width + spacing))
        box_top = top + (height / 2) - 0.75  # 1.5 inches height centered

        # 둥근 사각형 추가
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(box_left), Inches(box_top), Inches(box_width), Inches(1.5)
        )

        # 도형 스타일
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = fill_color

        # 텍스트 추가
        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = item.get("text", "")
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = font_color

        # 설명 추가 (있는 경우)
        if "description" in item and item["description"]:
            p2 = text_frame.add_paragraph()
            run2 = p2.add_run()
            run2.text = item["description"]
            run2.font.size = Pt(10)
            run2.font.color.rgb = font_color

        shapes_count += 1

        # 화살표 추가 (마지막 항목 제외)
        if i < num_items - 1:
            arrow_left = box_left + box_width
            arrow_top = top + (height / 2)
            arrow_width = spacing

            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(arrow_left), Inches(arrow_top - 0.25), Inches(arrow_width), Inches(0.5)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = fill_color
            arrow.line.color.rgb = fill_color
            shapes_count += 1

    return shapes_count


def create_hierarchy_diagram(
    slide, items: List[dict], left: float, top: float, width: float, height: float, fill_color, font_color
) -> int:
    """계층 구조 다이어그램 생성 (상단에서 하단으로)"""
    # 간단한 2단계 계층 구조
    # 첫 번째 항목은 상단 (루트), 나머지는 하단 (자식)
    shapes_count = 0

    # 루트 노드
    root_width = 3.0
    root_left = left + (width / 2) - (root_width / 2)
    root_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(root_left), Inches(top), Inches(root_width), Inches(1.0)
    )

    root_shape.fill.solid()
    root_shape.fill.fore_color.rgb = fill_color
    root_shape.line.color.rgb = fill_color

    text_frame = root_shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = items[0].get("text", "")
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = font_color

    shapes_count += 1

    # 자식 노드들
    if len(items) > 1:
        children = items[1:]
        num_children = len(children)
        child_width = (width / num_children) * 0.8
        spacing = (width - (child_width * num_children)) / (num_children + 1)

        for i, item in enumerate(children):
            child_left = left + spacing + (i * (child_width + spacing))
            child_top = top + 2.5

            child_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(child_left),
                Inches(child_top),
                Inches(child_width),
                Inches(1.0),
            )

            child_shape.fill.solid()
            child_shape.fill.fore_color.rgb = fill_color
            child_shape.line.color.rgb = fill_color

            text_frame = child_shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = item.get("text", "")
            run.font.size = Pt(12)
            run.font.color.rgb = font_color

            shapes_count += 1

    return shapes_count


def create_list_diagram(
    slide, items: List[dict], left: float, top: float, width: float, height: float, fill_color, font_color
) -> int:
    """목록형 다이어그램 생성 (세로 나열)"""
    num_items = len(items)
    box_height = (height / num_items) * 0.8
    spacing = (height - (box_height * num_items)) / (num_items + 1)

    shapes_count = 0

    for i, item in enumerate(items):
        box_top = top + spacing + (i * (box_height + spacing))

        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(box_top), Inches(width), Inches(box_height)
        )

        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = fill_color

        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = f"{i+1}. {item.get('text', '')}"
        run.font.size = Pt(14)
        run.font.color.rgb = font_color

        if "description" in item and item["description"]:
            p2 = text_frame.add_paragraph()
            run2 = p2.add_run()
            run2.text = item["description"]
            run2.font.size = Pt(10)
            run2.font.color.rgb = font_color

        shapes_count += 1

    return shapes_count


def create_cycle_diagram(
    slide, items: List[dict], left: float, top: float, width: float, height: float, fill_color, font_color
) -> int:
    """순환 구조 다이어그램 생성 (원형 배치)"""
    # 간단한 원형 배치 (4개 항목 가정)
    import math

    num_items = len(items)
    center_x = left + (width / 2)
    center_y = top + (height / 2)
    radius = min(width, height) / 3

    shapes_count = 0

    for i, item in enumerate(items):
        # 원형 배치 각도 계산
        angle = (2 * math.pi * i / num_items) - (math.pi / 2)  # -90도 시작 (상단)
        box_left = center_x + (radius * math.cos(angle)) - 0.75
        box_top = center_y + (radius * math.sin(angle)) - 0.5

        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(box_left), Inches(box_top), Inches(1.5), Inches(1.0)
        )

        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = fill_color

        text_frame = shape.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        run = p.add_run()
        run.text = item.get("text", "")
        run.font.size = Pt(12)
        run.font.color.rgb = font_color

        shapes_count += 1

    return shapes_count


def create_relationship_diagram(
    slide, items: List[dict], left: float, top: float, width: float, height: float, fill_color, font_color
) -> int:
    """관계도 다이어그램 생성 (중앙 허브 방식)"""
    # 첫 번째 항목을 중앙에, 나머지를 주변에 배치
    shapes_count = 0

    # 중앙 허브
    hub_size = 2.0
    hub_left = left + (width / 2) - (hub_size / 2)
    hub_top = top + (height / 2) - (hub_size / 2)

    hub_shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(hub_left), Inches(hub_top), Inches(hub_size), Inches(hub_size)
    )

    hub_shape.fill.solid()
    hub_shape.fill.fore_color.rgb = fill_color
    hub_shape.line.color.rgb = fill_color

    text_frame = hub_shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = items[0].get("text", "")
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = font_color

    shapes_count += 1

    # 주변 노드들
    if len(items) > 1:
        import math

        satellites = items[1:]
        num_satellites = len(satellites)
        radius = min(width, height) / 3

        for i, item in enumerate(satellites):
            angle = 2 * math.pi * i / num_satellites
            sat_left = hub_left + hub_size / 2 + (radius * math.cos(angle)) - 0.75
            sat_top = hub_top + hub_size / 2 + (radius * math.sin(angle)) - 0.5

            sat_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(sat_left), Inches(sat_top), Inches(1.5), Inches(1.0)
            )

            sat_shape.fill.solid()
            sat_shape.fill.fore_color.rgb = fill_color
            sat_shape.line.color.rgb = fill_color

            text_frame = sat_shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = item.get("text", "")
            run.font.size = Pt(11)
            run.font.color.rgb = font_color

            shapes_count += 1

    return shapes_count


if __name__ == "__main__":
    typer.run(content_add_smartart)
