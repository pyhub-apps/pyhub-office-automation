"""
PowerPoint 표 추가 명령어 (COM-First)
슬라이드에 표를 추가하고 데이터를 채웁니다.
"""

import csv
import json
from pathlib import Path
from typing import Optional

import typer

from pyhub_office_automation.version import get_version

from .utils import (
    PowerPointBackend,
    create_error_response,
    create_success_response,
    get_or_open_presentation,
    get_powerpoint_backend,
    normalize_path,
    validate_slide_number,
)


def content_add_table(
    slide_number: int = typer.Option(..., "--slide-number", help="표를 추가할 슬라이드 번호 (1부터 시작)"),
    rows: int = typer.Option(..., "--rows", help="표 행 수"),
    cols: int = typer.Option(..., "--cols", help="표 열 수"),
    left: float = typer.Option(..., "--left", help="표 왼쪽 위치 (인치)"),
    top: float = typer.Option(..., "--top", help="표 상단 위치 (인치)"),
    width: float = typer.Option(..., "--width", help="표 너비 (인치)"),
    height: float = typer.Option(..., "--height", help="표 높이 (인치)"),
    data: Optional[str] = typer.Option(None, "--data", help="표 데이터 (JSON 2차원 배열)"),
    data_file: Optional[str] = typer.Option(None, "--data-file", help="표 데이터 파일 (.csv 또는 .json)"),
    first_row_header: bool = typer.Option(False, "--first-row-header", help="첫 행을 헤더로 처리"),
    file_path: Optional[str] = typer.Option(None, "--file-path", help="PowerPoint 파일 경로"),
    presentation_name: Optional[str] = typer.Option(None, "--presentation-name", help="열려있는 프레젠테이션 이름 (COM 전용)"),
    backend: str = typer.Option("auto", "--backend", help="백엔드 선택 (auto/com/python-pptx)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 슬라이드에 표를 추가하고 데이터를 채웁니다.

    COM-First: Windows에서는 COM 백엔드 우선, python-pptx는 fallback

    **백엔드 선택**:
    - auto (기본): 자동으로 최적 백엔드 선택 (Windows COM 우선)
    - com: Windows COM 강제 사용 (완전한 기능)
    - python-pptx: python-pptx 강제 사용 (제한적 기능)

    **COM 백엔드 (Windows) - 완전한 기능!**:
    - ✅ 표 생성 및 데이터 채우기
    - Shapes.AddTable(), Table.Cell() 사용
    - 열려있는 프레젠테이션에서 직접 작업

    **python-pptx 백엔드**:
    - ⚠️ 파일 저장 필수 (--file-path 필수)
    - 표 생성 및 데이터 채우기 가능

    **데이터 입력 방법**:
      --data: JSON 2차원 배열로 직접 입력 (예: '[["Name", "Age"], ["Alice", "25"]]')
      --data-file: CSV 또는 JSON 파일에서 데이터 읽기

    예제:
        # COM 백엔드 (활성 프레젠테이션, JSON 데이터)
        oa ppt content-add-table --slide-number 1 --rows 3 --cols 2 --left 1 --top 2 --width 5 --height 3 --data '[["Name", "Age"], ["Alice", "25"], ["Bob", "30"]]' --first-row-header

        # COM 백엔드 (특정 프레젠테이션, CSV 파일)
        oa ppt content-add-table --slide-number 2 --rows 5 --cols 3 --left 1 --top 1 --width 7 --height 4 --data-file "data.csv" --first-row-header --presentation-name "report.pptx"

        # python-pptx 백엔드
        oa ppt content-add-table --slide-number 3 --rows 4 --cols 4 --left 0.5 --top 1.5 --width 9 --height 3.5 --file-path "report.pptx" --backend python-pptx
    """
    backend_inst = None

    try:
        # 입력 검증
        if data and data_file:
            result = create_error_response(
                command="content-add-table",
                error="--data와 --data-file은 동시에 사용할 수 없습니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        if rows < 1 or cols < 1:
            result = create_error_response(
                command="content-add-table",
                error="행과 열은 최소 1 이상이어야 합니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드 결정
        try:
            selected_backend = get_powerpoint_backend(force_backend=backend if backend != "auto" else None)
        except (ValueError, RuntimeError) as e:
            result = create_error_response(
                command="content-add-table",
                error=str(e),
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 데이터 로드
        table_data = None
        if data:
            try:
                table_data = json.loads(data)
                if not isinstance(table_data, list):
                    raise ValueError("데이터는 2차원 배열이어야 합니다")
                if table_data and not isinstance(table_data[0], list):
                    raise ValueError("데이터는 2차원 배열이어야 합니다")
            except json.JSONDecodeError as e:
                result = create_error_response(
                    command="content-add-table",
                    error=f"JSON 데이터 형식이 잘못되었습니다: {str(e)}",
                    error_type="ValueError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

        elif data_file:
            data_file_path = Path(normalize_path(data_file)).resolve()
            if not data_file_path.exists():
                result = create_error_response(
                    command="content-add-table",
                    error=f"데이터 파일을 찾을 수 없습니다: {data_file}",
                    error_type="FileNotFoundError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

            if data_file_path.suffix.lower() == ".csv":
                with open(data_file_path, "r", encoding="utf-8") as f:
                    reader = csv.reader(f)
                    table_data = list(reader)
            elif data_file_path.suffix.lower() == ".json":
                with open(data_file_path, "r", encoding="utf-8") as f:
                    table_data = json.load(f)
                    if not isinstance(table_data, list):
                        raise ValueError("JSON 데이터는 2차원 배열이어야 합니다")
                    if table_data and not isinstance(table_data[0], list):
                        raise ValueError("JSON 데이터는 2차원 배열이어야 합니다")
            else:
                result = create_error_response(
                    command="content-add-table",
                    error="데이터 파일은 .csv 또는 .json 형식이어야 합니다",
                    error_type="ValueError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

        # 데이터 크기 검증
        if table_data:
            data_rows = len(table_data)
            data_cols = max(len(row) for row in table_data) if table_data else 0

            if data_rows > rows:
                result = create_error_response(
                    command="content-add-table",
                    error=f"데이터 행 수({data_rows})가 표 행 수({rows})보다 큽니다",
                    error_type="ValueError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)
            if data_cols > cols:
                result = create_error_response(
                    command="content-add-table",
                    error=f"데이터 열 수({data_cols})가 표 열 수({cols})보다 큽니다",
                    error_type="ValueError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

        # 프레젠테이션 가져오기
        try:
            backend_inst, prs = get_or_open_presentation(
                file_path=file_path,
                presentation_name=presentation_name,
                backend=selected_backend,
            )
        except Exception as e:
            result = create_error_response(
                command="content-add-table",
                error=f"프레젠테이션을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드별 처리
        if selected_backend == PowerPointBackend.COM.value:
            # COM 백엔드: 완전한 테이블 추가 기능
            try:
                total_slides = prs.Slides.Count

                # 슬라이드 번호 검증 (COM은 1-based)
                if slide_number < 1 or slide_number > total_slides:
                    result = create_error_response(
                        command="content-add-table",
                        error=f"슬라이드 번호가 범위를 벗어났습니다: {slide_number} (1-{total_slides})",
                        error_type="ValueError",
                    )
                    typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                    raise typer.Exit(1)

                slide = prs.Slides(slide_number)

                # 인치를 포인트로 변환 (COM API는 포인트 사용)
                left_pt = left * 72
                top_pt = top * 72
                width_pt = width * 72
                height_pt = height * 72

                # 테이블 추가
                table_shape = slide.Shapes.AddTable(
                    NumRows=rows, NumColumns=cols, Left=left_pt, Top=top_pt, Width=width_pt, Height=height_pt
                )
                table = table_shape.Table

                # 데이터 채우기 (COM은 1-based index)
                if table_data:
                    for row_idx, row_data in enumerate(table_data, start=1):
                        for col_idx, cell_data in enumerate(row_data, start=1):
                            if row_idx <= rows and col_idx <= cols:
                                cell = table.Cell(row_idx, col_idx)
                                cell.Shape.TextFrame.TextRange.Text = str(cell_data)

                # 헤더 스타일 적용 (첫 행)
                if first_row_header and rows > 0:
                    for col_idx in range(1, cols + 1):
                        cell = table.Cell(1, col_idx)
                        text_range = cell.Shape.TextFrame.TextRange
                        text_range.Font.Bold = True

                # 성공 응답
                result_data = {
                    "backend": "com",
                    "slide_number": slide_number,
                    "table_size": {"rows": rows, "cols": cols},
                    "position": {
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                    },
                    "first_row_header": first_row_header,
                }

                if table_data:
                    result_data["data_filled"] = {
                        "rows": len(table_data),
                        "cols": max(len(row) for row in table_data) if table_data else 0,
                    }

                message = f"테이블 추가 완료 (COM): 슬라이드 {slide_number}, {rows}×{cols}"
                if table_data:
                    message += f", 데이터 {len(table_data)}행 채움"

            except Exception as e:
                result = create_error_response(
                    command="content-add-table",
                    error=f"테이블 추가 실패: {str(e)}",
                    error_type=type(e).__name__,
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

        else:
            # python-pptx 백엔드
            if not file_path:
                result = create_error_response(
                    command="content-add-table",
                    error="python-pptx 백엔드는 --file-path 옵션이 필수입니다",
                    error_type="ValueError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

            # 슬라이드 번호 검증
            slide_idx = validate_slide_number(slide_number, len(prs.slides))
            slide = prs.slides[slide_idx]

            # 테이블 추가
            from pptx.util import Inches

            table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
            table = table_shape.table

            # 데이터 채우기 (python-pptx는 0-based index)
            if table_data:
                for row_idx, row_data in enumerate(table_data):
                    for col_idx, cell_data in enumerate(row_data):
                        if row_idx < rows and col_idx < cols:
                            table.cell(row_idx, col_idx).text = str(cell_data)

            # 헤더 스타일 적용 (첫 행)
            if first_row_header and rows > 0:
                for col_idx in range(cols):
                    cell = table.cell(0, col_idx)
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.bold = True

            # 저장
            pptx_path = Path(normalize_path(file_path)).resolve()
            prs.save(str(pptx_path))

            # 결과 데이터
            result_data = {
                "backend": "python-pptx",
                "file": str(pptx_path),
                "file_name": pptx_path.name,
                "slide_number": slide_number,
                "table_size": {"rows": rows, "cols": cols},
                "position": {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                },
                "first_row_header": first_row_header,
            }

            if table_data:
                result_data["data_filled"] = {
                    "rows": len(table_data),
                    "cols": max(len(row) for row in table_data) if table_data else 0,
                }

            message = f"테이블 추가 완료 (python-pptx): 슬라이드 {slide_number}, {rows}×{cols}"
            if table_data:
                message += f", 데이터 {len(table_data)}행 채움"

        # 성공 응답
        response = create_success_response(
            data=result_data,
            command="content-add-table",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✅ {message}")
            typer.echo(f"📍 슬라이드: {slide_number}")
            typer.echo(f"📊 표 크기: {rows}행 × {cols}열")
            typer.echo(f"📐 위치: {left}in × {top}in")
            typer.echo(f"📏 크기: {width}in × {height}in")
            if first_row_header:
                typer.echo("🎯 첫 행: 헤더로 처리됨")
            if table_data:
                typer.echo(f"💾 데이터: {len(table_data)}행 채워짐")

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="content-add-table",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)
    finally:
        # python-pptx는 자동 정리, COM은 유지
        pass


if __name__ == "__main__":
    typer.run(content_add_table)
