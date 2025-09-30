"""
PowerPoint 텍스트 추가 명령어 (COM-First)
슬라이드에 텍스트를 추가합니다 (플레이스홀더 또는 자유 위치).
"""

import json
from pathlib import Path
from typing import Optional

import typer

from pyhub_office_automation.version import get_version

from .utils import (
    PlaceholderType,
    PowerPointBackend,
    create_error_response,
    create_success_response,
    get_or_open_presentation,
    get_placeholder_by_type,
    get_powerpoint_backend,
    normalize_path,
    parse_color,
    validate_slide_number,
)


def content_add_text(
    slide_number: int = typer.Option(..., "--slide-number", help="텍스트를 추가할 슬라이드 번호 (1부터 시작)"),
    placeholder: Optional[str] = typer.Option(
        None, "--placeholder", help="플레이스홀더 유형 (title/body/subtitle) - 이 옵션 사용 시 위치 옵션 무시"
    ),
    text: Optional[str] = typer.Option(None, "--text", help="추가할 텍스트 (직접 입력)"),
    text_file: Optional[str] = typer.Option(None, "--text-file", help="텍스트 파일 경로 (.txt)"),
    left: Optional[float] = typer.Option(None, "--left", help="텍스트 박스 왼쪽 위치 (인치)"),
    top: Optional[float] = typer.Option(None, "--top", help="텍스트 박스 상단 위치 (인치)"),
    width: Optional[float] = typer.Option(3.0, "--width", help="텍스트 박스 너비 (인치, 기본값: 3.0)"),
    height: Optional[float] = typer.Option(1.0, "--height", help="텍스트 박스 높이 (인치, 기본값: 1.0)"),
    font_size: Optional[int] = typer.Option(None, "--font-size", help="글꼴 크기 (포인트)"),
    font_color: Optional[str] = typer.Option(None, "--font-color", help="글꼴 색상 (색상명 또는 #RGB/#RRGGBB)"),
    bold: bool = typer.Option(False, "--bold", help="굵게 적용"),
    italic: bool = typer.Option(False, "--italic", help="기울임꼴 적용"),
    file_path: Optional[str] = typer.Option(None, "--file-path", help="PowerPoint 파일 경로"),
    presentation_name: Optional[str] = typer.Option(None, "--presentation-name", help="열려있는 프레젠테이션 이름 (COM 전용)"),
    backend: str = typer.Option("auto", "--backend", help="백엔드 선택 (auto/com/python-pptx)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 슬라이드에 텍스트를 추가합니다.

    COM-First: Windows에서는 COM 백엔드 우선, python-pptx는 fallback

    **백엔드 선택**:
    - auto (기본): 자동으로 최적 백엔드 선택 (Windows COM 우선)
    - com: Windows COM 강제 사용 (완전한 기능)
    - python-pptx: python-pptx 강제 사용 (제한적 기능)

    **COM 백엔드 (Windows) - 완전한 기능!**:
    - ✅ 플레이스홀더 및 자유 위치 텍스트 추가
    - ✅ 스마트 플레이스홀더 자동 감지
    - Shapes.AddTextbox(), TextFrame.TextRange 사용
    - 열려있는 프레젠테이션에서 직접 작업

    **python-pptx 백엔드**:
    - ⚠️ 파일 저장 필수 (--file-path 필수)
    - ✅ 플레이스홀더 자동 감지 지원
    - 플레이스홀더 및 자유 위치 텍스트 추가 가능

    **사용 모드**:
    1. **스마트 자동 감지 모드 (옵션 없음 - 추천!)**:
       - 레이아웃의 플레이스홀더를 자동으로 찾아 사용
       - 우선순위: Body > Title > Subtitle
       - 플레이스홀더가 없으면 중앙에 텍스트박스 생성
       - 템플릿 디자인을 최대한 활용!

    2. **플레이스홀더 모드 (--placeholder 지정)**:
       - title, body, subtitle 중 하나를 명시적으로 지정

    3. **자유 위치 모드 (--left, --top 지정)**:
       - 지정된 위치에 텍스트박스 생성

    **텍스트 입력**:
      --text: 직접 텍스트 입력
      --text-file: 파일에서 텍스트 읽기 (.txt)

    예제:
        # 🌟 자동 감지 모드 (권장) - 템플릿 디자인 활용
        oa ppt content-add-text --slide-number 1 --text "제목"

        # 플레이스홀더 명시적 지정
        oa ppt content-add-text --slide-number 1 --placeholder title --text "제목"

        # 자유 위치 지정
        oa ppt content-add-text --slide-number 2 --left 1 --top 2 --text "본문" --font-size 18

        # python-pptx 백엔드 (자동 감지)
        oa ppt content-add-text --slide-number 3 --text-file "content.txt" --file-path "report.pptx" --backend python-pptx
    """
    backend_inst = None

    try:
        # 입력 검증
        if not text and not text_file:
            result = create_error_response(
                command="content-add-text",
                error="--text 또는 --text-file 중 하나는 반드시 지정해야 합니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        if text and text_file:
            result = create_error_response(
                command="content-add-text",
                error="--text와 --text-file은 동시에 사용할 수 없습니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        if placeholder and (left is not None or top is not None):
            result = create_error_response(
                command="content-add-text",
                error="--placeholder와 --left/--top은 동시에 사용할 수 없습니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 스마트 자동 감지 모드: 옵션이 없으면 슬라이드 레이아웃의 플레이스홀더 자동 사용
        auto_detect_mode = False
        if not placeholder and left is None and top is None:
            auto_detect_mode = True
        elif not placeholder and (left is None or top is None):
            # left와 top 중 하나만 지정된 경우 에러
            result = create_error_response(
                command="content-add-text",
                error="--left와 --top은 함께 지정해야 합니다 (또는 --placeholder 사용, 또는 모두 생략하여 자동 감지)",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        if placeholder and placeholder not in [PlaceholderType.TITLE, PlaceholderType.BODY, PlaceholderType.SUBTITLE]:
            result = create_error_response(
                command="content-add-text",
                error=f"잘못된 플레이스홀더 유형: {placeholder}. 사용 가능: title, body, subtitle",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드 결정
        try:
            selected_backend = get_powerpoint_backend(force_backend=backend if backend != "auto" else None)
        except (ValueError, RuntimeError) as e:
            result = create_error_response(
                command="content-add-text",
                error=str(e),
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 텍스트 로드
        text_content = None
        if text_file:
            text_file_path = Path(normalize_path(text_file)).resolve()
            if not text_file_path.exists():
                result = create_error_response(
                    command="content-add-text",
                    error=f"텍스트 파일을 찾을 수 없습니다: {text_file}",
                    error_type="FileNotFoundError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)
            with open(text_file_path, "r", encoding="utf-8") as f:
                text_content = f.read()
        else:
            text_content = text

        # 프레젠테이션 가져오기
        try:
            backend_inst, prs = get_or_open_presentation(
                file_path=file_path,
                presentation_name=presentation_name,
                backend=selected_backend,
            )
        except Exception as e:
            result = create_error_response(
                command="content-add-text",
                error=f"프레젠테이션을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드별 처리
        mode = "placeholder" if placeholder else "position"
        auto_detected_placeholder = None

        if selected_backend == PowerPointBackend.COM.value:
            # COM 백엔드: 완전한 텍스트 추가 기능
            try:
                total_slides = prs.Slides.Count

                # 슬라이드 번호 검증 (COM은 1-based)
                if slide_number < 1 or slide_number > total_slides:
                    result = create_error_response(
                        command="content-add-text",
                        error=f"슬라이드 번호가 범위를 벗어났습니다: {slide_number} (1-{total_slides})",
                        error_type="ValueError",
                    )
                    typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                    raise typer.Exit(1)

                slide = prs.Slides(slide_number)

                # 자동 감지 모드: 슬라이드의 플레이스홀더 자동 선택
                if auto_detect_mode:
                    # 우선순위: Body(2) > Title(1,3) > Subtitle(10)
                    for shape in slide.Shapes:
                        if shape.Type == 14:  # msoPlaceholder
                            ph_type = shape.PlaceholderFormat.Type
                            if ph_type == 2:  # Body
                                placeholder = PlaceholderType.BODY
                                auto_detected_placeholder = "body"
                                mode = "placeholder"
                                break

                    if not auto_detected_placeholder:
                        for shape in slide.Shapes:
                            if shape.Type == 14:  # msoPlaceholder
                                ph_type = shape.PlaceholderFormat.Type
                                if ph_type in [1, 3]:  # Title, CenterTitle
                                    placeholder = PlaceholderType.TITLE
                                    auto_detected_placeholder = "title"
                                    mode = "placeholder"
                                    break

                    if not auto_detected_placeholder:
                        for shape in slide.Shapes:
                            if shape.Type == 14:  # msoPlaceholder
                                ph_type = shape.PlaceholderFormat.Type
                                if ph_type == 10:  # Subtitle
                                    placeholder = PlaceholderType.SUBTITLE
                                    auto_detected_placeholder = "subtitle"
                                    mode = "placeholder"
                                    break

                    # 플레이스홀더가 없으면 중앙에 텍스트박스 생성
                    if not auto_detected_placeholder:
                        # 슬라이드 크기 가져오기 (표준: 10" x 7.5")
                        slide_width = prs.PageSetup.SlideWidth / 72  # points to inches
                        slide_height = prs.PageSetup.SlideHeight / 72

                        # 중앙에 배치 (슬라이드의 40% 너비, 30% 높이)
                        left = slide_width * 0.3
                        top = slide_height * 0.35
                        width = slide_width * 0.4
                        height = slide_height * 0.3
                        mode = "position"

                if placeholder:
                    # 플레이스홀더 모드
                    placeholder_found = False
                    target_shape = None

                    # PlaceholderFormat.Type으로 플레이스홀더 찾기
                    for shape in slide.Shapes:
                        if shape.Type == 14:  # msoPlaceholder
                            ph_type = shape.PlaceholderFormat.Type
                            # 1=Title, 2=Body, 3=CenterTitle, 10=Subtitle
                            if placeholder == PlaceholderType.TITLE and ph_type in [1, 3]:
                                target_shape = shape
                                placeholder_found = True
                                break
                            elif placeholder == PlaceholderType.BODY and ph_type == 2:
                                target_shape = shape
                                placeholder_found = True
                                break
                            elif placeholder == PlaceholderType.SUBTITLE and ph_type == 10:
                                target_shape = shape
                                placeholder_found = True
                                break

                    if not placeholder_found:
                        result = create_error_response(
                            command="content-add-text",
                            error=f"슬라이드 {slide_number}에 '{placeholder}' 플레이스홀더가 없습니다",
                            error_type="ValueError",
                        )
                        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                        raise typer.Exit(1)

                    # 플레이스홀더에 텍스트 설정
                    text_frame = target_shape.TextFrame
                    text_range = text_frame.TextRange
                    text_range.Text = text_content

                    # 스타일 적용
                    if font_size is not None:
                        text_range.Font.Size = font_size
                    if bold:
                        text_range.Font.Bold = True
                    if italic:
                        text_range.Font.Italic = True
                    if font_color is not None:
                        color = parse_color(font_color)
                        # RGB를 COM 컬러 형식으로 변환 (B, G, R 순서)
                        text_range.Font.Color.RGB = color.red + (color.green << 8) + (color.blue << 16)

                else:
                    # 자유 위치 모드 - 텍스트 박스 추가
                    # 인치를 포인트로 변환 (1 inch = 72 points)
                    left_pt = left * 72
                    top_pt = top * 72
                    width_pt = width * 72
                    height_pt = height * 72

                    text_box = slide.Shapes.AddTextbox(
                        Orientation=1,  # msoTextOrientationHorizontal
                        Left=left_pt,
                        Top=top_pt,
                        Width=width_pt,
                        Height=height_pt,
                    )

                    text_frame = text_box.TextFrame
                    text_range = text_frame.TextRange
                    text_range.Text = text_content

                    # 스타일 적용
                    if font_size is not None:
                        text_range.Font.Size = font_size
                    if bold:
                        text_range.Font.Bold = True
                    if italic:
                        text_range.Font.Italic = True
                    if font_color is not None:
                        color = parse_color(font_color)
                        text_range.Font.Color.RGB = color.red + (color.green << 8) + (color.blue << 16)

                # 성공 응답
                result_data = {
                    "backend": "com",
                    "slide_number": slide_number,
                    "mode": mode,
                    "text_length": len(text_content),
                    "text_preview": text_content[:100] + "..." if len(text_content) > 100 else text_content,
                }

                # 자동 감지 정보 추가
                if auto_detected_placeholder:
                    result_data["auto_detected"] = True
                    result_data["auto_detected_placeholder"] = auto_detected_placeholder

                if placeholder:
                    result_data["placeholder"] = placeholder
                else:
                    result_data["position"] = {
                        "left": left,
                        "top": top,
                        "width": width,
                        "height": height,
                    }

                if font_size is not None:
                    result_data["font_size"] = font_size
                if font_color is not None:
                    result_data["font_color"] = font_color
                result_data["bold"] = bold
                result_data["italic"] = italic

                message = f"텍스트 추가 완료 (COM): 슬라이드 {slide_number}"
                if auto_detected_placeholder:
                    message += f", 자동 감지된 플레이스홀더 {auto_detected_placeholder}"
                elif placeholder:
                    message += f", 플레이스홀더 {placeholder}"
                else:
                    message += f", 위치 {left:.2f}in × {top:.2f}in"

            except Exception as e:
                result = create_error_response(
                    command="content-add-text",
                    error=f"텍스트 추가 실패: {str(e)}",
                    error_type=type(e).__name__,
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

        else:
            # python-pptx 백엔드
            if not file_path:
                result = create_error_response(
                    command="content-add-text",
                    error="python-pptx 백엔드는 --file-path 옵션이 필수입니다",
                    error_type="ValueError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

            # 슬라이드 번호 검증
            slide_idx = validate_slide_number(slide_number, len(prs.slides))
            slide = prs.slides[slide_idx]

            # 자동 감지 모드: 슬라이드의 플레이스홀더 자동 선택
            if auto_detect_mode:
                from pptx.enum.shapes import MSO_SHAPE_TYPE

                # 우선순위: Body > Title > Subtitle
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        if shape.placeholder_format.type == 2:  # BODY
                            placeholder = PlaceholderType.BODY
                            auto_detected_placeholder = "body"
                            mode = "placeholder"
                            break

                if not auto_detected_placeholder:
                    for shape in slide.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            if shape.placeholder_format.type in [1, 3]:  # TITLE, CENTER_TITLE
                                placeholder = PlaceholderType.TITLE
                                auto_detected_placeholder = "title"
                                mode = "placeholder"
                                break

                if not auto_detected_placeholder:
                    for shape in slide.shapes:
                        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                            if shape.placeholder_format.type == 10:  # SUBTITLE
                                placeholder = PlaceholderType.SUBTITLE
                                auto_detected_placeholder = "subtitle"
                                mode = "placeholder"
                                break

                # 플레이스홀더가 없으면 중앙에 텍스트박스 생성
                if not auto_detected_placeholder:
                    from pptx.util import Inches

                    # 슬라이드 크기 가져오기
                    slide_width = prs.slide_width.inches
                    slide_height = prs.slide_height.inches

                    # 중앙에 배치 (슬라이드의 40% 너비, 30% 높이)
                    left = slide_width * 0.3
                    top = slide_height * 0.35
                    width = slide_width * 0.4
                    height = slide_height * 0.3
                    mode = "position"

            # 텍스트 추가 처리
            if placeholder:
                # 플레이스홀더 모드
                shape = get_placeholder_by_type(slide, placeholder)
                if shape is None:
                    result = create_error_response(
                        command="content-add-text",
                        error=f"슬라이드 {slide_number}에 '{placeholder}' 플레이스홀더가 없습니다",
                        error_type="ValueError",
                    )
                    typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                    raise typer.Exit(1)
                text_frame = shape.text_frame
            else:
                # 자유 위치 모드
                from pptx.util import Inches

                shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                text_frame = shape.text_frame

            # 텍스트 설정
            text_frame.clear()
            paragraph = text_frame.paragraphs[0]
            run = paragraph.add_run()
            run.text = text_content

            # 스타일 적용
            if font_size is not None:
                from pptx.util import Pt

                run.font.size = Pt(font_size)

            if font_color is not None:
                color = parse_color(font_color)
                run.font.color.rgb = color

            if bold:
                run.font.bold = True

            if italic:
                run.font.italic = True

            # 저장
            pptx_path = Path(normalize_path(file_path)).resolve()
            prs.save(str(pptx_path))

            # 결과 데이터
            result_data = {
                "backend": "python-pptx",
                "file": str(pptx_path),
                "file_name": pptx_path.name,
                "slide_number": slide_number,
                "mode": mode,
                "text_length": len(text_content),
                "text_preview": text_content[:100] + "..." if len(text_content) > 100 else text_content,
            }

            # 자동 감지 정보 추가
            if auto_detected_placeholder:
                result_data["auto_detected"] = True
                result_data["auto_detected_placeholder"] = auto_detected_placeholder

            if placeholder:
                result_data["placeholder"] = placeholder
            else:
                result_data["position"] = {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height,
                }

            if font_size is not None:
                result_data["font_size"] = font_size
            if font_color is not None:
                result_data["font_color"] = font_color
            result_data["bold"] = bold
            result_data["italic"] = italic

            message = f"텍스트 추가 완료 (python-pptx): 슬라이드 {slide_number}"
            if auto_detected_placeholder:
                message += f", 자동 감지된 플레이스홀더 {auto_detected_placeholder}"
            elif placeholder:
                message += f", 플레이스홀더 {placeholder}"
            else:
                message += f", 위치 {left:.2f}in × {top:.2f}in"

        # 성공 응답
        response = create_success_response(
            data=result_data,
            command="content-add-text",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✅ {message}")
            typer.echo(f"📍 슬라이드: {slide_number}")
            if placeholder:
                typer.echo(f"🎯 플레이스홀더: {placeholder}")
            else:
                typer.echo(f"📐 위치: {left}in × {top}in")
                typer.echo(f"📏 크기: {width}in × {height}in")
            typer.echo(f"📝 텍스트 길이: {len(text_content)}자")
            typer.echo(f"📄 미리보기: {result_data['text_preview']}")

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="content-add-text",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)
    finally:
        # python-pptx는 자동 정리, COM은 유지
        pass


if __name__ == "__main__":
    typer.run(content_add_text)
