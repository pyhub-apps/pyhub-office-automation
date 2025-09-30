"""
PowerPoint 비디오 추가 명령어
슬라이드에 비디오 파일을 삽입합니다.
"""

import json
from pathlib import Path
from typing import Optional

import typer
from pptx import Presentation
from pptx.util import Inches

from pyhub_office_automation.version import get_version

from .utils import create_error_response, create_success_response, normalize_path, validate_slide_number

# 지원되는 비디오 형식
SUPPORTED_VIDEO_FORMATS = [".mp4", ".avi", ".wmv", ".mov", ".m4v", ".mpg", ".mpeg"]


def content_add_video(
    file_path: str = typer.Option(..., "--file-path", help="PowerPoint 파일 경로"),
    slide_number: int = typer.Option(..., "--slide-number", help="비디오를 추가할 슬라이드 번호 (1부터 시작)"),
    video_path: str = typer.Option(..., "--video-path", help="추가할 비디오 파일 경로"),
    left: Optional[float] = typer.Option(None, "--left", help="비디오 왼쪽 위치 (인치)"),
    top: Optional[float] = typer.Option(None, "--top", help="비디오 상단 위치 (인치)"),
    width: Optional[float] = typer.Option(6.0, "--width", help="비디오 너비 (인치, 기본값: 6.0)"),
    height: Optional[float] = typer.Option(4.5, "--height", help="비디오 높이 (인치, 기본값: 4.5)"),
    poster_frame: Optional[str] = typer.Option(None, "--poster-frame", help="포스터 프레임 이미지 경로 (선택)"),
    center: bool = typer.Option(False, "--center", help="슬라이드 중앙에 배치 (--left, --top 무시)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 슬라이드에 비디오를 추가합니다.

    지원 비디오 형식:
      MP4, AVI, WMV, MOV, M4V, MPG, MPEG

    위치 지정:
      --center: 슬라이드 중앙에 배치
      --left, --top: 특정 위치에 배치

    포스터 프레임:
      --poster-frame: 비디오 재생 전 표시될 이미지 (PNG, JPG 등)
      미지정 시: 비디오 첫 프레임 사용

    예제:
        oa ppt content-add-video --file-path "presentation.pptx" --slide-number 2 --video-path "demo.mp4" --center
        oa ppt content-add-video --file-path "presentation.pptx" --slide-number 3 --video-path "tutorial.mp4" --left 1 --top 2 --poster-frame "thumbnail.png"
    """
    try:
        # 입력 검증
        if not center and (left is None or top is None):
            raise ValueError("--center를 사용하지 않는 경우 --left와 --top을 모두 지정해야 합니다")

        # 파일 경로 정규화 및 존재 확인
        normalized_pptx_path = normalize_path(file_path)
        pptx_path = Path(normalized_pptx_path).resolve()

        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint 파일을 찾을 수 없습니다: {pptx_path}")

        normalized_video_path = normalize_path(video_path)
        video_file = Path(normalized_video_path).resolve()

        if not video_file.exists():
            raise FileNotFoundError(f"비디오 파일을 찾을 수 없습니다: {video_file}")

        # 비디오 형식 검증
        video_ext = video_file.suffix.lower()
        if video_ext not in SUPPORTED_VIDEO_FORMATS:
            supported_str = ", ".join(SUPPORTED_VIDEO_FORMATS)
            raise ValueError(f"지원하지 않는 비디오 형식: {video_ext}\n지원 형식: {supported_str}")

        # 포스터 프레임 검증
        poster_file = None
        if poster_frame:
            normalized_poster_path = normalize_path(poster_frame)
            poster_file = Path(normalized_poster_path).resolve()
            if not poster_file.exists():
                raise FileNotFoundError(f"포스터 프레임 이미지를 찾을 수 없습니다: {poster_file}")

        # 프레젠테이션 열기
        prs = Presentation(str(pptx_path))

        # 슬라이드 번호 검증
        slide_idx = validate_slide_number(slide_number, len(prs.slides))
        slide = prs.slides[slide_idx]

        # 위치 계산
        if center:
            # 슬라이드 크기 가져오기 (EMU 단위)
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # 인치 단위로 변환
            slide_width_in = slide_width / 914400  # 1 inch = 914400 EMU
            slide_height_in = slide_height / 914400

            # 중앙 배치 위치 계산
            final_left = (slide_width_in - width) / 2
            final_top = (slide_height_in - height) / 2
        else:
            final_left = left
            final_top = top

        # 비디오 추가
        if poster_file:
            # 포스터 프레임과 함께 비디오 추가
            movie = slide.shapes.add_movie(
                str(video_file),
                Inches(final_left),
                Inches(final_top),
                Inches(width),
                Inches(height),
                poster_frame_image=str(poster_file),
            )
        else:
            # 포스터 프레임 없이 비디오 추가
            movie = slide.shapes.add_movie(
                str(video_file), Inches(final_left), Inches(final_top), Inches(width), Inches(height)
            )

        # 비디오 파일 정보
        video_size_mb = video_file.stat().st_size / (1024 * 1024)

        # 저장
        prs.save(str(pptx_path))

        # 결과 데이터 구성
        result_data = {
            "file": str(pptx_path),
            "slide_number": slide_number,
            "video_file": str(video_file),
            "video_format": video_ext,
            "video_size_mb": round(video_size_mb, 2),
            "position": {
                "left": round(final_left, 2),
                "top": round(final_top, 2),
                "width": width,
                "height": height,
            },
            "centered": center,
            "has_poster_frame": poster_file is not None,
        }

        if poster_file:
            result_data["poster_frame"] = str(poster_file)

        # 성공 응답
        message = f"슬라이드 {slide_number}에 비디오를 추가했습니다"
        if center:
            message += " (중앙 배치)"

        response = create_success_response(
            data=result_data,
            command="content-add-video",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✅ {message}")
            typer.echo(f"📄 파일: {pptx_path.name}")
            typer.echo(f"📍 슬라이드: {slide_number}")
            typer.echo(f"🎬 비디오: {video_file.name}")
            typer.echo(f"📦 형식: {video_ext.upper()}")
            typer.echo(f"💾 크기: {result_data['video_size_mb']} MB")
            typer.echo(f"📐 위치: {result_data['position']['left']}in × {result_data['position']['top']}in")
            typer.echo(f"📏 크기: {width}in × {height}in")
            if poster_file:
                typer.echo(f"🖼️ 포스터 프레임: {poster_file.name}")

    except FileNotFoundError as e:
        error_response = create_error_response(e, "content-add-video")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ {str(e)}", err=True)
        raise typer.Exit(1)

    except ValueError as e:
        error_response = create_error_response(e, "content-add-video")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ {str(e)}", err=True)
        raise typer.Exit(1)

    except Exception as e:
        error_response = create_error_response(e, "content-add-video")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ 예기치 않은 오류: {str(e)}", err=True)
        raise typer.Exit(1)


if __name__ == "__main__":
    typer.run(content_add_video)
