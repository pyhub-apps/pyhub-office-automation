"""
PowerPoint 콘텐츠 업데이트 명령어 (COM-First)
슬라이드의 기존 콘텐츠(텍스트, 이미지, 위치 등)를 수정합니다.
"""

import json
from pathlib import Path
from typing import Optional

import typer

from pyhub_office_automation.version import get_version

from .utils import (
    PowerPointBackend,
    create_error_response,
    create_success_response,
    get_or_open_presentation,
    get_powerpoint_backend,
    normalize_path,
    validate_slide_number,
)


def content_update(
    slide_number: int = typer.Option(..., "--slide-number", help="콘텐츠를 업데이트할 슬라이드 번호 (1부터 시작)"),
    shape_index: Optional[int] = typer.Option(None, "--shape-index", help="Shape 인덱스 (1부터 시작)"),
    shape_name: Optional[str] = typer.Option(None, "--shape-name", help="Shape 이름"),
    text: Optional[str] = typer.Option(None, "--text", help="업데이트할 텍스트 내용"),
    image_path: Optional[str] = typer.Option(None, "--image-path", help="교체할 이미지 파일 경로"),
    left: Optional[float] = typer.Option(None, "--left", help="새 위치 - 왼쪽 (인치)"),
    top: Optional[float] = typer.Option(None, "--top", help="새 위치 - 상단 (인치)"),
    width: Optional[float] = typer.Option(None, "--width", help="새 크기 - 너비 (인치)"),
    height: Optional[float] = typer.Option(None, "--height", help="새 크기 - 높이 (인치)"),
    file_path: Optional[str] = typer.Option(None, "--file-path", help="PowerPoint 파일 경로"),
    presentation_name: Optional[str] = typer.Option(None, "--presentation-name", help="열려있는 프레젠테이션 이름 (COM 전용)"),
    backend: str = typer.Option("auto", "--backend", help="백엔드 선택 (auto/com/python-pptx)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 슬라이드의 기존 콘텐츠를 업데이트합니다.

    COM-First: Windows에서는 COM 백엔드 우선, python-pptx는 fallback

    **백엔드 선택**:
    - auto (기본): 자동으로 최적 백엔드 선택 (Windows COM 우선)
    - com: Windows COM 강제 사용 (완전한 기능)
    - python-pptx: python-pptx 강제 사용 (제한적 기능)

    **COM 백엔드 (Windows) - 완전한 기능!**:
    - ✅ 텍스트, 위치, 크기, 색상 등 모든 속성 수정 가능
    - ✅ 이미지 교체 (Delete + AddPicture)
    - ✅ Shape 인덱스/이름으로 찾기
    - 열려있는 프레젠테이션에서 직접 작업

    **python-pptx 백엔드**:
    - ⚠️ 파일 저장 필수 (--file-path 필수)
    - 제한적 속성 수정만 가능

    Shape 선택 (둘 중 하나만 지정):
      --shape-index: Shape 인덱스 (1부터 시작)
      --shape-name: Shape 이름

    업데이트 옵션:
      --text: 텍스트 내용 변경
      --image-path: 이미지 교체 (picture shape만 가능)
      --left, --top: 위치 이동
      --width, --height: 크기 조정

    예제:
        # COM 백엔드 (활성 프레젠테이션)
        oa ppt content-update --slide-number 1 --shape-index 1 --text "새 제목"

        # COM 백엔드 (이미지 교체)
        oa ppt content-update --slide-number 2 --shape-name "Picture 1" --image-path "new_image.png"

        # COM 백엔드 (위치/크기 조정)
        oa ppt content-update --slide-number 3 --shape-index 2 --left 2 --top 2 --width 4 --height 3 --presentation-name "report.pptx"

        # python-pptx 백엔드
        oa ppt content-update --slide-number 4 --shape-index 1 --text "새 텍스트" --file-path "report.pptx" --backend python-pptx
    """

    try:
        # 입력 검증
        if shape_index is None and shape_name is None:
            result = create_error_response(
                command="content-update",
                error="--shape-index 또는 --shape-name 중 하나는 반드시 지정해야 합니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        if shape_index is not None and shape_name is not None:
            result = create_error_response(
                command="content-update",
                error="--shape-index와 --shape-name은 동시에 사용할 수 없습니다",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        if text is None and image_path is None and left is None and top is None and width is None and height is None:
            result = create_error_response(
                command="content-update",
                error="업데이트할 내용을 지정해야 합니다 (--text, --image-path, --left, --top, --width, --height 중 하나 이상)",
                error_type="ValueError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 이미지 경로 검증
        img_path = None
        if image_path:
            normalized_image_path = normalize_path(image_path)
            img_path = Path(normalized_image_path).resolve()
            if not img_path.exists():
                result = create_error_response(
                    command="content-update",
                    error=f"이미지 파일을 찾을 수 없습니다: {img_path}",
                    error_type="FileNotFoundError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

        # 백엔드 결정
        try:
            selected_backend = get_powerpoint_backend(force_backend=backend if backend != "auto" else None)
        except (ValueError, RuntimeError) as e:
            result = create_error_response(
                command="content-update",
                error=str(e),
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 프레젠테이션 가져오기
        try:
            backend_inst, prs = get_or_open_presentation(
                file_path=file_path,
                presentation_name=presentation_name,
                backend=selected_backend,
            )
        except Exception as e:
            result = create_error_response(
                command="content-update",
                error=f"프레젠테이션을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 백엔드별 처리
        if selected_backend == PowerPointBackend.COM.value:
            # COM 백엔드: 완전한 콘텐츠 업데이트 기능
            try:
                total_slides = prs.Slides.Count

                # 슬라이드 번호 검증 (COM은 1-based)
                if slide_number < 1 or slide_number > total_slides:
                    result = create_error_response(
                        command="content-update",
                        error=f"슬라이드 번호가 범위를 벗어났습니다: {slide_number} (1-{total_slides})",
                        error_type="ValueError",
                    )
                    typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                    raise typer.Exit(1)

                slide = prs.Slides(slide_number)

                # Shape 찾기
                if shape_index is not None:
                    # 인덱스로 찾기 (COM은 1-based)
                    if shape_index < 1 or shape_index > slide.Shapes.Count:
                        result = create_error_response(
                            command="content-update",
                            error=f"Shape 인덱스가 범위를 벗어났습니다: {shape_index} (1-{slide.Shapes.Count})",
                            error_type="ValueError",
                        )
                        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                        raise typer.Exit(1)
                    shape = slide.Shapes(shape_index)
                    shape_identifier = f"index_{shape_index}"
                else:
                    # 이름으로 찾기
                    try:
                        shape = slide.Shapes(shape_name)
                        shape_identifier = shape_name
                    except Exception:
                        result = create_error_response(
                            command="content-update",
                            error=f"Shape를 찾을 수 없습니다: {shape_name}",
                            error_type="ValueError",
                        )
                        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                        raise typer.Exit(1)

                # 업데이트 내용 기록
                updates = []

                # 텍스트 업데이트
                if text is not None:
                    if not shape.HasTextFrame:
                        result = create_error_response(
                            command="content-update",
                            error=f"이 Shape는 텍스트를 지원하지 않습니다 (Type: {shape.Type})",
                            error_type="ValueError",
                        )
                        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                        raise typer.Exit(1)

                    shape.TextFrame.TextRange.Text = text
                    updates.append(f"텍스트 변경: '{text[:50]}{'...' if len(text) > 50 else ''}' ({len(text)}자)")

                # 이미지 교체
                if image_path:
                    # Shape Type: 13 = msoPicture
                    if shape.Type != 13:
                        result = create_error_response(
                            command="content-update",
                            error=f"이미지 교체는 picture shape만 가능합니다 (현재 Type: {shape.Type})",
                            error_type="ValueError",
                        )
                        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                        raise typer.Exit(1)

                    # 기존 이미지 정보 저장
                    old_left = shape.Left
                    old_top = shape.Top
                    old_width = shape.Width
                    old_height = shape.Height

                    # Shape 삭제
                    shape.Delete()

                    # 새 이미지 추가 (같은 위치와 크기)
                    new_left = left * 72 if left is not None else old_left
                    new_top = top * 72 if top is not None else old_top
                    new_width = width * 72 if width is not None else old_width
                    new_height = height * 72 if height is not None else old_height

                    shape = slide.Shapes.AddPicture(
                        str(img_path),
                        0,  # LinkToFile = msoFalse
                        -1,  # SaveWithDocument = msoTrue
                        new_left,
                        new_top,
                        new_width,
                        new_height,
                    )
                    updates.append(f"이미지 교체: {img_path.name}")

                # 위치 업데이트
                if left is not None and image_path is None:
                    shape.Left = left * 72
                    updates.append(f"위치 변경 (left): {left}in")
                if top is not None and image_path is None:
                    shape.Top = top * 72
                    updates.append(f"위치 변경 (top): {top}in")

                # 크기 업데이트
                if width is not None and image_path is None:
                    shape.Width = width * 72
                    updates.append(f"크기 변경 (width): {width}in")
                if height is not None and image_path is None:
                    shape.Height = height * 72
                    updates.append(f"크기 변경 (height): {height}in")

                # 현재 위치/크기 정보 (포인트를 인치로 변환)
                current_position = {
                    "left": round(shape.Left / 72, 2),
                    "top": round(shape.Top / 72, 2),
                    "width": round(shape.Width / 72, 2),
                    "height": round(shape.Height / 72, 2),
                }

                # 성공 응답
                result_data = {
                    "backend": "com",
                    "slide_number": slide_number,
                    "shape_identifier": shape_identifier,
                    "shape_type": shape.Type,
                    "updates": updates,
                    "update_count": len(updates),
                    "current_position": current_position,
                }

                message = f"콘텐츠 업데이트 완료 (COM): 슬라이드 {slide_number}, {len(updates)}개 항목"

            except Exception as e:
                result = create_error_response(
                    command="content-update",
                    error=f"콘텐츠 업데이트 실패: {str(e)}",
                    error_type=type(e).__name__,
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

        else:
            # python-pptx 백엔드
            if not file_path:
                result = create_error_response(
                    command="content-update",
                    error="python-pptx 백엔드는 --file-path 옵션이 필수입니다",
                    error_type="ValueError",
                )
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                raise typer.Exit(1)

            # 슬라이드 번호 검증
            slide_idx = validate_slide_number(slide_number, len(prs.slides))
            slide = prs.slides[slide_idx]

            # Shape 찾기
            from .utils import get_shape_by_index_or_name

            identifier = shape_index if shape_index is not None else shape_name
            # python-pptx는 0-based이므로 shape_index를 0-based로 변환
            if shape_index is not None:
                identifier = shape_index - 1  # COM은 1-based, python-pptx는 0-based

            shape = get_shape_by_index_or_name(slide, identifier)

            # 업데이트 내용 기록
            updates = []

            # 텍스트 업데이트
            if text is not None:
                if not hasattr(shape, "text_frame"):
                    result = create_error_response(
                        command="content-update",
                        error=f"이 Shape는 텍스트를 지원하지 않습니다: {shape.shape_type}",
                        error_type="ValueError",
                    )
                    typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                    raise typer.Exit(1)

                shape.text_frame.clear()
                paragraph = shape.text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = text
                updates.append(f"텍스트 변경: '{text[:50]}{'...' if len(text) > 50 else ''}' ({len(text)}자)")

            # 이미지 교체
            if image_path:
                from pptx.enum.shapes import MSO_SHAPE_TYPE
                from pptx.util import Inches

                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    result = create_error_response(
                        command="content-update",
                        error=f"이미지 교체는 picture shape만 가능합니다 (현재: {shape.shape_type})",
                        error_type="ValueError",
                    )
                    typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
                    raise typer.Exit(1)

                # 기존 이미지 정보 저장
                old_left = shape.left
                old_top = shape.top
                old_width = shape.width
                old_height = shape.height

                # Shape 삭제
                sp = shape.element
                sp.getparent().remove(sp)

                # 새 이미지 추가
                new_left = Inches(left) if left is not None else old_left
                new_top = Inches(top) if top is not None else old_top
                new_width = Inches(width) if width is not None else old_width
                new_height = Inches(height) if height is not None else old_height

                picture = slide.shapes.add_picture(str(img_path), new_left, new_top, width=new_width, height=new_height)
                shape = picture
                updates.append(f"이미지 교체: {img_path.name}")

            # 위치 업데이트
            if left is not None and image_path is None:
                from pptx.util import Inches

                shape.left = Inches(left)
                updates.append(f"위치 변경 (left): {left}in")
            if top is not None and image_path is None:
                from pptx.util import Inches

                shape.top = Inches(top)
                updates.append(f"위치 변경 (top): {top}in")

            # 크기 업데이트
            if width is not None and image_path is None:
                from pptx.util import Inches

                shape.width = Inches(width)
                updates.append(f"크기 변경 (width): {width}in")
            if height is not None and image_path is None:
                from pptx.util import Inches

                shape.height = Inches(height)
                updates.append(f"크기 변경 (height): {height}in")

            # 저장
            pptx_path = Path(normalize_path(file_path)).resolve()
            prs.save(str(pptx_path))

            # 결과 데이터
            shape_identifier = shape_name if shape_name else f"index_{shape_index}"
            result_data = {
                "backend": "python-pptx",
                "file": str(pptx_path),
                "file_name": pptx_path.name,
                "slide_number": slide_number,
                "shape_identifier": shape_identifier,
                "shape_type": str(shape.shape_type),
                "updates": updates,
                "update_count": len(updates),
                "current_position": {
                    "left": round(shape.left / 914400, 2),
                    "top": round(shape.top / 914400, 2),
                    "width": round(shape.width / 914400, 2),
                    "height": round(shape.height / 914400, 2),
                },
            }

            message = f"콘텐츠 업데이트 완료 (python-pptx): 슬라이드 {slide_number}, {len(updates)}개 항목"

        # 성공 응답
        response = create_success_response(
            data=result_data,
            command="content-update",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✅ {message}")
            typer.echo(f"📍 슬라이드: {slide_number}")
            typer.echo(f"🎯 Shape: {result_data['shape_identifier']}")
            typer.echo(f"📦 Shape 타입: {result_data['shape_type']}")
            typer.echo(f"\n🔄 업데이트 내역:")
            for update in updates:
                typer.echo(f"  • {update}")
            typer.echo(f"\n📐 현재 위치/크기:")
            pos = result_data["current_position"]
            typer.echo(f"  위치: {pos['left']}in × {pos['top']}in")
            typer.echo(f"  크기: {pos['width']}in × {pos['height']}in")

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="content-update",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)
    finally:
        # python-pptx는 자동 정리, COM은 유지
        pass


if __name__ == "__main__":
    typer.run(content_update)
