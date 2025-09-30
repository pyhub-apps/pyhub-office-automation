"""
PowerPoint 레이아웃 적용 명령어
슬라이드에 특정 레이아웃을 적용합니다.
"""

import json
from pathlib import Path
from typing import Optional, Union

import typer
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

from pyhub_office_automation.version import get_version

from .utils import (
    create_error_response,
    create_success_response,
    get_layout_by_name_or_index,
    normalize_path,
    validate_slide_number,
)


def layout_apply(
    file_path: str = typer.Option(..., "--file-path", help="PowerPoint 파일 경로"),
    slide_number: int = typer.Option(..., "--slide-number", help="레이아웃을 적용할 슬라이드 번호 (1부터 시작)"),
    layout: str = typer.Option(..., "--layout", help="레이아웃 이름 또는 인덱스 (예: 'Title Slide' 또는 0)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 슬라이드에 특정 레이아웃을 적용합니다.

    ⚠️  제약사항:
        python-pptx는 기존 슬라이드의 레이아웃 직접 변경을 지원하지 않습니다.
        이 명령어는 레이아웃 정보를 조회하고 참조 목적으로만 사용됩니다.
        실제 레이아웃 변경은 PowerPoint에서 직접 수행해야 합니다.

    레이아웃은 이름 또는 인덱스로 지정할 수 있습니다.

    예제:
        oa ppt layout-apply --file-path "presentation.pptx" --slide-number 1 --layout "Title Slide"
        oa ppt layout-apply --file-path "report.pptx" --slide-number 2 --layout 1
    """
    try:
        # 파일 경로 정규화 및 존재 확인
        normalized_path = normalize_path(file_path)
        pptx_path = Path(normalized_path).resolve()

        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint 파일을 찾을 수 없습니다: {pptx_path}")

        # 프레젠테이션 열기
        prs = Presentation(str(pptx_path))

        # 슬라이드 번호 검증
        slide_idx = validate_slide_number(slide_number, len(prs.slides))
        slide = prs.slides[slide_idx]

        # 기존 레이아웃 정보 저장
        old_layout_name = slide.slide_layout.name

        # 레이아웃 찾기 (이름 또는 인덱스)
        # 숫자로 변환 시도 (인덱스인 경우)
        try:
            layout_identifier = int(layout)
        except ValueError:
            # 문자열 그대로 (이름인 경우)
            layout_identifier = layout

        # 레이아웃 가져오기
        new_layout = get_layout_by_name_or_index(prs, layout_identifier)

        # python-pptx 제약: 기존 슬라이드의 레이아웃 직접 변경 불가
        # 대신 레이아웃 정보만 조회하고 반환
        limitation_message = (
            "python-pptx는 기존 슬라이드의 레이아웃 직접 변경을 지원하지 않습니다. "
            "PowerPoint에서 직접 변경하거나, 새 슬라이드를 생성할 때 원하는 레이아웃을 지정하세요."
        )

        # 결과 데이터 구성
        result_data = {
            "file": str(pptx_path),
            "file_name": pptx_path.name,
            "slide_number": slide_number,
            "current_layout": old_layout_name,
            "requested_layout": new_layout.name,
            "layout_index": prs.slide_layouts.index(new_layout),
            "applied": False,
            "limitation": limitation_message,
            "alternative": "oa ppt slide-add --file-path 'file.pptx' --layout 'Title and Content'",
        }

        # 경고 응답
        message = f"레이아웃 '{new_layout.name}' 정보를 조회했습니다 (슬라이드 {slide_number}, 현재: {old_layout_name})"
        response = create_success_response(
            data=result_data,
            command="layout-apply",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"⚠️  {message}")
            typer.echo(f"📄 파일: {pptx_path.name}")
            typer.echo(f"📍 슬라이드: {slide_number}")
            typer.echo(f"📐 현재 레이아웃: {old_layout_name}")
            typer.echo(f"📐 요청 레이아웃: {new_layout.name}")
            typer.echo(f"\n💡 제약사항: {limitation_message}")
            typer.echo(f"💡 대안: {result_data['alternative']}")

    except FileNotFoundError as e:
        error_response = create_error_response(e, "layout-apply")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ {str(e)}", err=True)
        raise typer.Exit(1)

    except (ValueError, IndexError) as e:
        error_response = create_error_response(e, "layout-apply")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ {str(e)}", err=True)
        raise typer.Exit(1)

    except Exception as e:
        error_response = create_error_response(e, "layout-apply")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ 예기치 않은 오류: {str(e)}", err=True)
        raise typer.Exit(1)


if __name__ == "__main__":
    typer.run(layout_apply)
