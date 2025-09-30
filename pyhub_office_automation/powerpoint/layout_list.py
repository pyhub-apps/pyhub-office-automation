"""
PowerPoint 레이아웃 목록 조회 명령어
프레젠테이션의 사용 가능한 모든 레이아웃 정보를 제공합니다.
"""

import json
from pathlib import Path
from typing import Any, Dict, List

import typer
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE

from pyhub_office_automation.version import get_version

from .utils import create_error_response, create_success_response, normalize_path


def get_placeholder_info(layout) -> List[Dict[str, Any]]:
    """
    레이아웃의 placeholder 정보를 추출합니다.

    Args:
        layout: SlideLayout 객체

    Returns:
        List[Dict]: placeholder 정보 리스트
    """
    placeholders = []

    try:
        for shape in layout.placeholders:
            try:
                # Placeholder 타입 파싱
                placeholder_type = "unknown"
                if hasattr(shape, "placeholder_format"):
                    ph_type = shape.placeholder_format.type
                    if ph_type == PP_PLACEHOLDER_TYPE.TITLE:
                        placeholder_type = "title"
                    elif ph_type == PP_PLACEHOLDER_TYPE.BODY:
                        placeholder_type = "body"
                    elif ph_type == PP_PLACEHOLDER_TYPE.SUBTITLE:
                        placeholder_type = "subtitle"
                    elif ph_type == PP_PLACEHOLDER_TYPE.CENTER_TITLE:
                        placeholder_type = "center_title"
                    elif ph_type == PP_PLACEHOLDER_TYPE.PICTURE:
                        placeholder_type = "picture"
                    elif ph_type == PP_PLACEHOLDER_TYPE.CHART:
                        placeholder_type = "chart"
                    elif ph_type == PP_PLACEHOLDER_TYPE.TABLE:
                        placeholder_type = "table"
                    elif ph_type == PP_PLACEHOLDER_TYPE.OBJECT:
                        placeholder_type = "object"
                    else:
                        placeholder_type = str(ph_type)

                placeholder_info = {
                    "idx": shape.placeholder_format.idx,
                    "type": placeholder_type,
                }

                # Shape 이름 추가
                if hasattr(shape, "name"):
                    placeholder_info["name"] = shape.name

                placeholders.append(placeholder_info)

            except Exception:
                # 개별 placeholder 처리 실패 시 스킵
                continue

    except Exception:
        # placeholder 접근 실패 시 빈 리스트 반환
        pass

    return placeholders


def layout_list(
    file_path: str = typer.Option(..., "--file-path", help="PowerPoint 파일 경로"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 프레젠테이션의 사용 가능한 레이아웃 목록을 조회합니다.

    각 레이아웃의 인덱스, 이름, placeholder 정보를 제공합니다.

    예제:
        oa ppt layout-list --file-path "presentation.pptx"
        oa ppt layout-list --file-path "report.pptx" --format text
    """
    try:
        # 파일 경로 정규화 및 존재 확인
        normalized_path = normalize_path(file_path)
        pptx_path = Path(normalized_path).resolve()

        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint 파일을 찾을 수 없습니다: {pptx_path}")

        # 프레젠테이션 열기
        prs = Presentation(str(pptx_path))

        # 레이아웃 정보 수집
        layouts_info = []
        for idx, layout in enumerate(prs.slide_layouts):
            try:
                layout_name = layout.name
            except Exception:
                layout_name = f"Layout {idx}"

            # Placeholder 정보 추출
            placeholders = get_placeholder_info(layout)

            layout_data = {
                "index": idx,
                "name": layout_name,
                "placeholders": placeholders,
                "placeholder_count": len(placeholders),
            }

            layouts_info.append(layout_data)

        # 결과 데이터 구성
        result_data = {
            "file": str(pptx_path),
            "file_name": pptx_path.name,
            "total_layouts": len(layouts_info),
            "layouts": layouts_info,
        }

        # 성공 응답
        message = f"총 {len(layouts_info)}개의 레이아웃을 찾았습니다"
        response = create_success_response(
            data=result_data,
            command="layout-list",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✅ {message}")
            typer.echo(f"📄 파일: {pptx_path.name}")
            typer.echo(f"\n📐 사용 가능한 레이아웃:")
            for layout in layouts_info:
                placeholders_str = ", ".join([ph["type"] for ph in layout["placeholders"]])
                typer.echo(f"  [{layout['index']}] {layout['name']}")
                if placeholders_str:
                    typer.echo(f"      Placeholders: {placeholders_str}")

    except FileNotFoundError as e:
        error_response = create_error_response(e, "layout-list")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ {str(e)}", err=True)
        raise typer.Exit(1)

    except Exception as e:
        error_response = create_error_response(e, "layout-list")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ 예기치 않은 오류: {str(e)}", err=True)
        raise typer.Exit(1)


if __name__ == "__main__":
    typer.run(layout_list)
