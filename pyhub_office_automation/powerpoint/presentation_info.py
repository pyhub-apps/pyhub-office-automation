"""
PowerPoint 프레젠테이션 정보 조회 명령어 (Typer 버전)
프레젠테이션의 상세 정보 제공
"""

import datetime
import json
import sys
from pathlib import Path

import typer

from pyhub_office_automation.version import get_version

from .utils import create_error_response, create_success_response, normalize_path


def presentation_info(
    file_path: str = typer.Option(..., "--file-path", help="정보를 조회할 프레젠테이션 파일 경로"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택"),
):
    """
    PowerPoint 프레젠테이션의 상세 정보를 조회합니다.

    슬라이드 수, 크기, 레이아웃, 파일 정보 등을 제공합니다.

    예제:
        oa ppt presentation-info --file-path "report.pptx"
        oa ppt presentation-info --file-path "C:/Work/presentation.pptx"
    """
    try:
        # python-pptx import
        try:
            from pptx import Presentation
        except ImportError:
            result = create_error_response(
                command="presentation-info",
                error="python-pptx 패키지가 설치되지 않았습니다",
                error_type="ImportError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 경로 정규화 및 검증
        file_path_normalized = normalize_path(file_path)
        file_path_obj = Path(file_path_normalized).resolve()

        if not file_path_obj.exists():
            result = create_error_response(
                command="presentation-info",
                error=f"프레젠테이션 파일을 찾을 수 없습니다: {file_path}",
                error_type="FileNotFoundError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 프레젠테이션 열기
        try:
            prs = Presentation(str(file_path_obj))
        except Exception as e:
            result = create_error_response(
                command="presentation-info",
                error=f"프레젠테이션 파일을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 슬라이드 정보
        slide_count = len(prs.slides)

        # 슬라이드 크기 (EMU -> inches 변환, 1 inch = 914400 EMU)
        slide_width_inches = prs.slide_width / 914400
        slide_height_inches = prs.slide_height / 914400

        # 레이아웃 정보
        layouts_info = []
        for idx, layout in enumerate(prs.slide_layouts):
            try:
                layout_name = layout.name
            except:
                layout_name = f"Layout {idx}"

            layouts_info.append(
                {
                    "index": idx,
                    "name": layout_name,
                }
            )

        # 파일 정보
        file_stat = file_path_obj.stat()
        last_modified = datetime.datetime.fromtimestamp(file_stat.st_mtime).isoformat()

        # 슬라이드별 간단한 정보
        slides_info = []
        for idx, slide in enumerate(prs.slides):
            try:
                # 슬라이드 레이아웃 이름
                layout_name = slide.slide_layout.name
            except:
                layout_name = "Unknown"

            # 슬라이드의 도형 수
            shape_count = len(slide.shapes)

            slides_info.append(
                {
                    "slide_number": idx + 1,
                    "layout": layout_name,
                    "shape_count": shape_count,
                }
            )

        # 성공 응답 데이터
        data = {
            "file_name": file_path_obj.name,
            "file_path": str(file_path_obj),
            "file_size_bytes": file_stat.st_size,
            "last_modified": last_modified,
            "slide_count": slide_count,
            "slide_width_inches": round(slide_width_inches, 2),
            "slide_height_inches": round(slide_height_inches, 2),
            "layouts_count": len(layouts_info),
            "layouts": layouts_info,
            "slides": slides_info,
        }

        result = create_success_response(
            command="presentation-info",
            data=data,
            message=f"프레젠테이션 정보: {file_path_obj.name}",
        )

        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="presentation-info",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)
