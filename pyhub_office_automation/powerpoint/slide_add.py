"""
PowerPoint 슬라이드 추가 명령어 (Typer 버전)
새 슬라이드를 지정된 레이아웃으로 추가
"""

import json
from pathlib import Path
from typing import Optional

import typer

from pyhub_office_automation.version import get_version

from .utils import (
    create_error_response,
    create_success_response,
    get_layout_by_name_or_index,
    normalize_path,
    validate_slide_number,
)


def slide_add(
    file_path: str = typer.Option(..., "--file-path", help="프레젠테이션 파일 경로"),
    layout: Optional[str] = typer.Option(None, "--layout", help="레이아웃 이름 또는 인덱스 (기본: 1=Title and Content)"),
    position: Optional[int] = typer.Option(None, "--position", help="삽입 위치 (1-based, 기본: 끝에 추가)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택 (json/text)"),
):
    """
    PowerPoint 프레젠테이션에 새 슬라이드를 추가합니다.

    레이아웃과 위치를 지정할 수 있습니다.

    예제:
        oa ppt slide-add --file-path "report.pptx"
        oa ppt slide-add --file-path "report.pptx" --layout "Blank" --position 3
        oa ppt slide-add --file-path "report.pptx" --layout 0  # 인덱스로 지정
    """
    try:
        # python-pptx import
        try:
            from pptx import Presentation
        except ImportError:
            result = create_error_response(
                command="slide-add",
                error="python-pptx 패키지가 설치되지 않았습니다",
                error_type="ImportError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 경로 정규화 및 검증
        file_path_normalized = normalize_path(file_path)
        file_path_obj = Path(file_path_normalized).resolve()

        if not file_path_obj.exists():
            result = create_error_response(
                command="slide-add",
                error=f"프레젠테이션 파일을 찾을 수 없습니다: {file_path}",
                error_type="FileNotFoundError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 프레젠테이션 열기
        try:
            prs = Presentation(str(file_path_obj))
        except Exception as e:
            result = create_error_response(
                command="slide-add",
                error=f"프레젠테이션 파일을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        total_slides = len(prs.slides)

        # 레이아웃 결정 (기본값: 1 = Title and Content)
        if layout is None:
            layout_identifier = 1  # 기본 레이아웃 인덱스
        else:
            # 숫자 문자열이면 int로 변환
            try:
                layout_identifier = int(layout)
            except ValueError:
                layout_identifier = layout

        # 레이아웃 찾기
        try:
            slide_layout = get_layout_by_name_or_index(prs, layout_identifier)
            layout_name = slide_layout.name
        except (ValueError, IndexError) as e:
            result = create_error_response(
                command="slide-add",
                error=str(e),
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 위치 검증 (끝에 추가 허용)
        if position is not None:
            target_idx = validate_slide_number(position, total_slides, allow_append=True)
        else:
            target_idx = None  # 끝에 추가

        # 슬라이드 추가 (끝에)
        new_slide = prs.slides.add_slide(slide_layout)
        new_slide_idx = len(prs.slides) - 1  # 0-based index

        # 위치 지정이 있고, 끝이 아니면 순서 조정
        if position is not None and target_idx < new_slide_idx:
            # XML 레벨에서 순서 조정
            xml_slides = prs.slides._sldIdLst
            slide_id_element = xml_slides[new_slide_idx]
            xml_slides.remove(slide_id_element)
            xml_slides.insert(target_idx, slide_id_element)
            final_position = position
        else:
            final_position = len(prs.slides)

        # 파일 저장
        prs.save(str(file_path_obj))

        # 성공 응답
        data = {
            "slide_number": final_position,
            "layout": layout_name,
            "total_slides": len(prs.slides),
        }

        result = create_success_response(
            command="slide-add",
            data=data,
            message=f"슬라이드가 추가되었습니다 (위치: {final_position}, 레이아웃: {layout_name})",
        )

        if output_format == "json":
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✓ 슬라이드 추가 완료")
            typer.echo(f"  위치: {final_position}번")
            typer.echo(f"  레이아웃: {layout_name}")
            typer.echo(f"  총 슬라이드: {len(prs.slides)}개")

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="slide-add",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)


if __name__ == "__main__":
    typer.run(slide_add)
