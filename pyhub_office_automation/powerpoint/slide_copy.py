"""
PowerPoint 슬라이드 복사 명령어 (Typer 버전)
지정된 슬라이드를 다른 위치로 복사
"""

import json
from pathlib import Path

import typer

from pyhub_office_automation.version import get_version

from .utils import create_error_response, create_success_response, normalize_path, validate_slide_number


def slide_copy(
    file_path: str = typer.Option(..., "--file-path", help="프레젠테이션 파일 경로"),
    source: int = typer.Option(..., "--source", help="원본 슬라이드 번호 (1-based)"),
    destination: int = typer.Option(..., "--destination", help="대상 위치 (1-based)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택 (json/text)"),
):
    """
    PowerPoint 슬라이드를 지정된 위치로 복사합니다.

    현재 버전: 레이아웃만 복사 (빈 슬라이드 생성)
    향후 업데이트: 전체 내용 복사 예정

    예제:
        oa ppt slide-copy --file-path "report.pptx" --source 2 --destination 5
    """
    try:
        # python-pptx import
        try:
            from pptx import Presentation
        except ImportError:
            result = create_error_response(
                command="slide-copy",
                error="python-pptx 패키지가 설치되지 않았습니다",
                error_type="ImportError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 경로 정규화 및 검증
        file_path_normalized = normalize_path(file_path)
        file_path_obj = Path(file_path_normalized).resolve()

        if not file_path_obj.exists():
            result = create_error_response(
                command="slide-copy",
                error=f"프레젠테이션 파일을 찾을 수 없습니다: {file_path}",
                error_type="FileNotFoundError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 프레젠테이션 열기
        try:
            prs = Presentation(str(file_path_obj))
        except Exception as e:
            result = create_error_response(
                command="slide-copy",
                error=f"프레젠테이션 파일을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        total_slides = len(prs.slides)

        # 슬라이드 번호 검증
        source_idx = validate_slide_number(source, total_slides, allow_append=False)
        # destination은 끝에 추가 가능
        dest_idx = validate_slide_number(destination, total_slides, allow_append=True)

        # 원본 슬라이드 가져오기
        source_slide = prs.slides[source_idx]
        source_layout = source_slide.slide_layout
        layout_name = source_layout.name

        # 1. 원본 위치에 복제 (끝에 추가)
        new_slide = prs.slides.add_slide(source_layout)
        new_slide_idx = len(prs.slides) - 1

        # 2. 복제된 슬라이드를 대상 위치로 이동
        # destination 위치 조정 (원본 뒤에 추가되었으므로)
        if dest_idx <= source_idx:
            # destination이 source보다 앞이면 그대로 사용
            target_idx = dest_idx
        else:
            # destination이 source보다 뒤면 -1 (원본 뒤에 복제본이 추가되었으므로)
            target_idx = dest_idx

        if target_idx < new_slide_idx:
            xml_slides = prs.slides._sldIdLst
            slide_id_element = xml_slides[new_slide_idx]
            xml_slides.remove(slide_id_element)
            xml_slides.insert(target_idx, slide_id_element)

        # 파일 저장
        prs.save(str(file_path_obj))

        # 성공 응답
        data = {
            "source_slide": source,
            "destination_slide": destination,
            "layout": layout_name,
            "total_slides": len(prs.slides),
            "note": "현재 레이아웃만 복사됩니다 (내용 복사는 향후 업데이트 예정)",
        }

        result = create_success_response(
            command="slide-copy",
            data=data,
            message=f"슬라이드 {source}번이 {destination}번 위치로 복사되었습니다",
        )

        if output_format == "json":
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✓ 슬라이드 복사 완료")
            typer.echo(f"  원본: {source}번")
            typer.echo(f"  복사본 위치: {destination}번")
            typer.echo(f"  레이아웃: {layout_name}")
            typer.echo(f"  총 슬라이드: {len(prs.slides)}개")
            typer.echo(f"  ℹ️  현재 레이아웃만 복사됩니다")

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="slide-copy",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)


if __name__ == "__main__":
    typer.run(slide_copy)
