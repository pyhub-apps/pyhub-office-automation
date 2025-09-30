"""
PowerPoint 슬라이드 삭제 명령어 (Typer 버전)
지정된 슬라이드를 삭제
"""

import json
from pathlib import Path

import typer

from pyhub_office_automation.version import get_version

from .utils import create_error_response, create_success_response, normalize_path, validate_slide_number


def slide_delete(
    file_path: str = typer.Option(..., "--file-path", help="프레젠테이션 파일 경로"),
    slide_number: int = typer.Option(..., "--slide-number", help="삭제할 슬라이드 번호 (1-based)"),
    force: bool = typer.Option(False, "--force", help="확인 없이 강제 삭제"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택 (json/text)"),
):
    """
    PowerPoint 프레젠테이션에서 슬라이드를 삭제합니다.

    지정된 슬라이드를 프레젠테이션에서 제거합니다.

    예제:
        oa ppt slide-delete --file-path "report.pptx" --slide-number 3
        oa ppt slide-delete --file-path "report.pptx" --slide-number 1 --force
    """
    try:
        # python-pptx import
        try:
            from pptx import Presentation
        except ImportError:
            result = create_error_response(
                command="slide-delete",
                error="python-pptx 패키지가 설치되지 않았습니다",
                error_type="ImportError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 경로 정규화 및 검증
        file_path_normalized = normalize_path(file_path)
        file_path_obj = Path(file_path_normalized).resolve()

        if not file_path_obj.exists():
            result = create_error_response(
                command="slide-delete",
                error=f"프레젠테이션 파일을 찾을 수 없습니다: {file_path}",
                error_type="FileNotFoundError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 프레젠테이션 열기
        try:
            prs = Presentation(str(file_path_obj))
        except Exception as e:
            result = create_error_response(
                command="slide-delete",
                error=f"프레젠테이션 파일을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        total_slides = len(prs.slides)

        # 슬라이드 번호 검증
        slide_idx = validate_slide_number(slide_number, total_slides, allow_append=False)

        # Force 플래그 없으면 경고 (JSON에만 포함)
        warning = None
        if not force and total_slides == 1:
            warning = "마지막 슬라이드를 삭제하려고 합니다. --force 옵션을 사용하세요."

        # XML 슬라이드 리스트에서 삭제
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)

        # 관계 ID 가져오기 및 관계 제거 (메모리 누수 방지)
        try:
            rId = slides[slide_idx].rId
            prs.part.drop_rel(rId)
        except Exception:
            # 관계 제거 실패해도 계속 진행
            pass

        # XML 리스트에서 제거
        xml_slides.remove(slides[slide_idx])

        # 파일 저장
        prs.save(str(file_path_obj))

        # 성공 응답
        data = {
            "deleted_slide_number": slide_number,
            "remaining_slides": len(prs.slides),
            "warning": warning,
        }

        result = create_success_response(
            command="slide-delete",
            data=data,
            message=f"슬라이드 {slide_number}번이 삭제되었습니다 (남은 슬라이드: {len(prs.slides)}개)",
        )

        if output_format == "json":
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✓ 슬라이드 삭제 완료")
            typer.echo(f"  삭제된 슬라이드: {slide_number}번")
            typer.echo(f"  남은 슬라이드: {len(prs.slides)}개")
            if warning:
                typer.echo(f"  ⚠️  {warning}")

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="slide-delete",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)


if __name__ == "__main__":
    typer.run(slide_delete)
