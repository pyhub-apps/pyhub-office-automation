"""
PowerPoint 슬라이드 복제 명령어 (Typer 버전)
지정된 슬라이드를 바로 뒤에 복제 (Phase 1: 레이아웃만 복제)
"""

import json
from pathlib import Path

import typer

from pyhub_office_automation.version import get_version

from .utils import create_error_response, create_success_response, normalize_path, validate_slide_number


def slide_duplicate(
    file_path: str = typer.Option(..., "--file-path", help="프레젠테이션 파일 경로"),
    slide_number: int = typer.Option(..., "--slide-number", help="복제할 슬라이드 번호 (1-based)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택 (json/text)"),
):
    """
    PowerPoint 슬라이드를 복제합니다 (바로 뒤에 추가).

    현재 버전: 레이아웃만 복제 (빈 슬라이드 생성)
    향후 업데이트: 전체 내용 복제 예정

    예제:
        oa ppt slide-duplicate --file-path "report.pptx" --slide-number 2
    """
    try:
        # python-pptx import
        try:
            from pptx import Presentation
        except ImportError:
            result = create_error_response(
                command="slide-duplicate",
                error="python-pptx 패키지가 설치되지 않았습니다",
                error_type="ImportError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 경로 정규화 및 검증
        file_path_normalized = normalize_path(file_path)
        file_path_obj = Path(file_path_normalized).resolve()

        if not file_path_obj.exists():
            result = create_error_response(
                command="slide-duplicate",
                error=f"프레젠테이션 파일을 찾을 수 없습니다: {file_path}",
                error_type="FileNotFoundError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 프레젠테이션 열기
        try:
            prs = Presentation(str(file_path_obj))
        except Exception as e:
            result = create_error_response(
                command="slide-duplicate",
                error=f"프레젠테이션 파일을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        total_slides = len(prs.slides)

        # 슬라이드 번호 검증
        slide_idx = validate_slide_number(slide_number, total_slides, allow_append=False)

        # 원본 슬라이드 가져오기
        source_slide = prs.slides[slide_idx]
        source_layout = source_slide.slide_layout
        layout_name = source_layout.name

        # 간이 복제: 같은 레이아웃으로 새 슬라이드 생성 (빈 슬라이드)
        # Phase 1 구현 - 레이아웃만 복제
        new_slide = prs.slides.add_slide(source_layout)
        new_slide_idx = len(prs.slides) - 1

        # 원본 바로 뒤로 이동
        target_idx = slide_idx + 1
        if target_idx < new_slide_idx:
            xml_slides = prs.slides._sldIdLst
            slide_id_element = xml_slides[new_slide_idx]
            xml_slides.remove(slide_id_element)
            xml_slides.insert(target_idx, slide_id_element)

        # 파일 저장
        prs.save(str(file_path_obj))

        # 성공 응답
        data = {
            "source_slide": slide_number,
            "duplicated_slide": slide_number + 1,
            "layout": layout_name,
            "total_slides": len(prs.slides),
            "note": "현재 레이아웃만 복제됩니다 (내용 복제는 향후 업데이트 예정)",
        }

        result = create_success_response(
            command="slide-duplicate",
            data=data,
            message=f"슬라이드 {slide_number}번이 복제되었습니다 (위치: {slide_number + 1}번)",
        )

        if output_format == "json":
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✓ 슬라이드 복제 완료")
            typer.echo(f"  원본: {slide_number}번")
            typer.echo(f"  복제본: {slide_number + 1}번")
            typer.echo(f"  레이아웃: {layout_name}")
            typer.echo(f"  총 슬라이드: {len(prs.slides)}개")
            typer.echo(f"  ℹ️  현재 레이아웃만 복제됩니다")

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="slide-duplicate",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)


if __name__ == "__main__":
    typer.run(slide_duplicate)
