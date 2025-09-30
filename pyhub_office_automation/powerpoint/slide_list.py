"""
PowerPoint 슬라이드 목록 조회 명령어 (Typer 버전)
프레젠테이션의 모든 슬라이드 정보 제공
"""

import json
from pathlib import Path

import typer

from pyhub_office_automation.version import get_version

from .utils import create_error_response, create_success_response, get_slide_content_summary, get_slide_title, normalize_path


def slide_list(
    file_path: str = typer.Option(..., "--file-path", help="정보를 조회할 프레젠테이션 파일 경로"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택 (json/text)"),
):
    """
    PowerPoint 프레젠테이션의 모든 슬라이드 목록을 조회합니다.

    각 슬라이드의 번호, 레이아웃, 제목, 콘텐츠 요약을 제공합니다.

    예제:
        oa ppt slide-list --file-path "report.pptx"
        oa ppt slide-list --file-path "C:/Work/presentation.pptx" --format text
    """
    try:
        # python-pptx import
        try:
            from pptx import Presentation
        except ImportError:
            result = create_error_response(
                command="slide-list",
                error="python-pptx 패키지가 설치되지 않았습니다",
                error_type="ImportError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 경로 정규화 및 검증
        file_path_normalized = normalize_path(file_path)
        file_path_obj = Path(file_path_normalized).resolve()

        if not file_path_obj.exists():
            result = create_error_response(
                command="slide-list",
                error=f"프레젠테이션 파일을 찾을 수 없습니다: {file_path}",
                error_type="FileNotFoundError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 프레젠테이션 열기
        try:
            prs = Presentation(str(file_path_obj))
        except Exception as e:
            result = create_error_response(
                command="slide-list",
                error=f"프레젠테이션 파일을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 슬라이드 정보 수집
        total_slides = len(prs.slides)
        slides_info = []

        for idx, slide in enumerate(prs.slides):
            # 레이아웃 이름
            try:
                layout_name = slide.slide_layout.name
            except Exception:
                layout_name = "Unknown"

            # 슬라이드 제목
            title = get_slide_title(slide)

            # 콘텐츠 요약 (도형 타입별 개수)
            content_summary = get_slide_content_summary(slide)

            slides_info.append(
                {
                    "slide_number": idx + 1,
                    "layout": layout_name,
                    "title": title,
                    "shapes": content_summary,
                }
            )

        # 성공 응답 데이터
        data = {
            "file_name": file_path_obj.name,
            "file_path": str(file_path_obj),
            "total_slides": total_slides,
            "slides": slides_info,
        }

        result = create_success_response(
            command="slide-list",
            data=data,
            message=f"슬라이드 목록: {file_path_obj.name} ({total_slides}개)",
        )

        if output_format == "json":
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        else:
            # Text 형식 출력
            typer.echo(f"📊 프레젠테이션: {file_path_obj.name}")
            typer.echo(f"총 슬라이드: {total_slides}개")
            typer.echo()
            for slide_info in slides_info:
                typer.echo(f"슬라이드 {slide_info['slide_number']}: {slide_info['layout']}")
                if slide_info["title"]:
                    typer.echo(f"  제목: {slide_info['title']}")
                shapes = slide_info["shapes"]
                shape_summary = ", ".join([f"{key}:{val}" for key, val in shapes.items() if val > 0])
                if shape_summary:
                    typer.echo(f"  도형: {shape_summary}")
                typer.echo()

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="slide-list",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)


if __name__ == "__main__":
    typer.run(slide_list)
