"""
PowerPoint 슬라이드 순서 변경 명령어 (Typer 버전)
슬라이드를 다른 위치로 이동
"""

import json
from pathlib import Path

import typer

from pyhub_office_automation.version import get_version

from .utils import create_error_response, create_success_response, normalize_path, validate_slide_number


def slide_reorder(
    file_path: str = typer.Option(..., "--file-path", help="프레젠테이션 파일 경로"),
    from_position: int = typer.Option(..., "--from", help="이동할 슬라이드 번호 (1-based)"),
    to_position: int = typer.Option(..., "--to", help="이동할 위치 (1-based)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 선택 (json/text)"),
):
    """
    PowerPoint 슬라이드의 순서를 변경합니다 (이동).

    지정된 슬라이드를 다른 위치로 이동합니다.

    예제:
        oa ppt slide-reorder --file-path "report.pptx" --from 3 --to 1
        oa ppt slide-reorder --file-path "presentation.pptx" --from 1 --to 5
    """
    try:
        # python-pptx import
        try:
            from pptx import Presentation
        except ImportError:
            result = create_error_response(
                command="slide-reorder",
                error="python-pptx 패키지가 설치되지 않았습니다",
                error_type="ImportError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 경로 정규화 및 검증
        file_path_normalized = normalize_path(file_path)
        file_path_obj = Path(file_path_normalized).resolve()

        if not file_path_obj.exists():
            result = create_error_response(
                command="slide-reorder",
                error=f"프레젠테이션 파일을 찾을 수 없습니다: {file_path}",
                error_type="FileNotFoundError",
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        # 프레젠테이션 열기
        try:
            prs = Presentation(str(file_path_obj))
        except Exception as e:
            result = create_error_response(
                command="slide-reorder",
                error=f"프레젠테이션 파일을 열 수 없습니다: {str(e)}",
                error_type=type(e).__name__,
            )
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            raise typer.Exit(1)

        total_slides = len(prs.slides)

        # 슬라이드 번호 검증
        from_idx = validate_slide_number(from_position, total_slides, allow_append=False)
        to_idx = validate_slide_number(to_position, total_slides, allow_append=False)

        # 같은 위치면 아무것도 하지 않음
        if from_idx == to_idx:
            result = create_success_response(
                command="slide-reorder",
                data={
                    "from_position": from_position,
                    "to_position": to_position,
                    "total_slides": total_slides,
                    "changed": False,
                },
                message=f"슬라이드가 이미 위치 {from_position}에 있습니다 (변경 없음)",
            )
            if output_format == "json":
                typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
            else:
                typer.echo(f"✓ 슬라이드가 이미 위치 {from_position}에 있습니다")
            return

        # XML 슬라이드 리스트에서 순서 변경
        xml_slides = prs.slides._sldIdLst
        slide_id_element = xml_slides[from_idx]

        # 제거 후 새 위치에 삽입
        xml_slides.remove(slide_id_element)
        xml_slides.insert(to_idx, slide_id_element)

        # 파일 저장
        prs.save(str(file_path_obj))

        # 성공 응답
        data = {
            "from_position": from_position,
            "to_position": to_position,
            "total_slides": total_slides,
            "changed": True,
        }

        result = create_success_response(
            command="slide-reorder",
            data=data,
            message=f"슬라이드를 {from_position}번에서 {to_position}번으로 이동했습니다",
        )

        if output_format == "json":
            typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"✓ 슬라이드 이동 완료: {from_position}번 → {to_position}번")

    except typer.Exit:
        raise
    except Exception as e:
        result = create_error_response(
            command="slide-reorder",
            error=str(e),
            error_type=type(e).__name__,
        )
        typer.echo(json.dumps(result, ensure_ascii=False, indent=2))
        raise typer.Exit(1)


if __name__ == "__main__":
    typer.run(slide_reorder)
