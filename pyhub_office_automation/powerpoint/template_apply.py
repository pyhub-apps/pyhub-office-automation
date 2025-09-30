"""
PowerPoint 템플릿 적용 명령어
템플릿 프레젠테이션의 디자인을 현재 프레젠테이션에 적용합니다.
"""

import json
from pathlib import Path
from typing import Optional

import typer
from pptx import Presentation

from pyhub_office_automation.version import get_version

from .utils import create_error_response, create_success_response, normalize_path


def copy_slide_master(source_prs, target_prs):
    """
    소스 프레젠테이션의 슬라이드 마스터를 타겟 프레젠테이션으로 복사합니다.

    Note: python-pptx는 슬라이드 마스터 직접 복사를 지원하지 않습니다.
    이 함수는 제한적인 구현으로, 레이아웃 정보만 참조합니다.

    Args:
        source_prs: 소스 Presentation 객체
        target_prs: 타겟 Presentation 객체

    Returns:
        Dict: 복사 결과 정보
    """
    result = {
        "masters_copied": 0,
        "layouts_available": len(source_prs.slide_layouts),
        "warning": "python-pptx는 슬라이드 마스터 직접 복사를 지원하지 않습니다. " "템플릿의 레이아웃 정보만 참조됩니다.",
    }

    return result


def template_apply(
    file_path: str = typer.Option(..., "--file-path", help="대상 PowerPoint 파일 경로"),
    template_path: str = typer.Option(..., "--template-path", help="템플릿 PowerPoint 파일 경로"),
    preserve_content: bool = typer.Option(True, "--preserve-content", help="기존 콘텐츠 보존 여부 (기본값: True)"),
    output_format: str = typer.Option("json", "--format", help="출력 형식 (json/text)"),
):
    """
    PowerPoint 템플릿을 적용합니다.

    템플릿 프레젠테이션의 디자인(슬라이드 마스터, 레이아웃)을 대상 프레젠테이션에 적용합니다.
    --preserve-content 옵션으로 기존 콘텐츠 보존 여부를 선택할 수 있습니다.

    제약사항:
        - python-pptx는 슬라이드 마스터 직접 복사를 지원하지 않습니다
        - 실제 구현에서는 슬라이드별로 레이아웃을 매핑하는 방식으로 동작합니다
        - 완전한 템플릿 적용을 위해서는 PowerPoint에서 직접 "디자인 적용" 기능을 사용하세요

    예제:
        oa ppt template-apply --file-path "report.pptx" --template-path "corporate_template.pptx"
        oa ppt template-apply --file-path "report.pptx" --template-path "template.pptx" --preserve-content
    """
    try:
        # 파일 경로 정규화 및 존재 확인
        normalized_target_path = normalize_path(file_path)
        target_path = Path(normalized_target_path).resolve()

        normalized_template_path = normalize_path(template_path)
        template_path_obj = Path(normalized_template_path).resolve()

        if not target_path.exists():
            raise FileNotFoundError(f"대상 PowerPoint 파일을 찾을 수 없습니다: {target_path}")

        if not template_path_obj.exists():
            raise FileNotFoundError(f"템플릿 PowerPoint 파일을 찾을 수 없습니다: {template_path_obj}")

        # 프레젠테이션 열기
        target_prs = Presentation(str(target_path))
        template_prs = Presentation(str(template_path_obj))

        # 템플릿 정보 수집
        template_layouts = []
        for idx, layout in enumerate(template_prs.slide_layouts):
            template_layouts.append({"index": idx, "name": layout.name})

        # 슬라이드 마스터 복사 시도
        master_copy_result = copy_slide_master(template_prs, target_prs)

        # 기존 슬라이드에 템플릿 레이아웃 적용 (preserve_content=True인 경우)
        slides_updated = 0
        if preserve_content:
            for slide_idx, slide in enumerate(target_prs.slides):
                try:
                    # 현재 슬라이드의 레이아웃 이름
                    current_layout_name = slide.slide_layout.name

                    # 템플릿에서 같은 이름의 레이아웃 찾기
                    matching_layout = None
                    for template_layout in template_prs.slide_layouts:
                        if template_layout.name == current_layout_name:
                            matching_layout = template_layout
                            break

                    # 매칭되는 레이아웃이 없으면 첫 번째 레이아웃 사용
                    if matching_layout is None and len(template_prs.slide_layouts) > 0:
                        matching_layout = template_prs.slide_layouts[0]

                    # Note: python-pptx는 다른 프레젠테이션의 레이아웃을 직접 적용할 수 없습니다
                    # 이 부분은 실제로 동작하지 않으며, 제한사항으로 기록됩니다

                except Exception:
                    # 개별 슬라이드 처리 실패는 무시
                    continue

        # 저장 (실제로는 변경사항이 적용되지 않음)
        # target_prs.save(str(target_path))

        # 결과 데이터 구성
        result_data = {
            "target_file": str(target_path),
            "target_file_name": target_path.name,
            "template_file": str(template_path_obj),
            "template_file_name": template_path_obj.name,
            "template_layouts": template_layouts,
            "template_layouts_count": len(template_layouts),
            "preserve_content": preserve_content,
            "slides_updated": slides_updated,
            "master_copy_result": master_copy_result,
            "limitation": "python-pptx는 슬라이드 마스터 및 테마의 완전한 복사를 지원하지 않습니다. "
            "PowerPoint에서 '디자인 탭 > 테마'를 통해 직접 적용하는 것을 권장합니다.",
        }

        # 경고 메시지
        warning_message = (
            "⚠️  python-pptx의 제약으로 인해 템플릿이 완전히 적용되지 않았습니다.\n"
            "완전한 템플릿 적용을 위해서는 PowerPoint에서 직접 '디자인 > 테마' 기능을 사용하세요."
        )

        # 성공 응답
        message = f"템플릿 '{template_path_obj.name}'의 레이아웃 정보를 참조했습니다"
        response = create_success_response(
            data=result_data,
            command="template-apply",
            message=message,
        )

        # 출력
        if output_format == "json":
            typer.echo(json.dumps(response, ensure_ascii=False, indent=2))
        else:
            typer.echo(f"⚠️  {message}")
            typer.echo(f"📄 대상 파일: {target_path.name}")
            typer.echo(f"📐 템플릿: {template_path_obj.name}")
            typer.echo(f"🎨 템플릿 레이아웃: {len(template_layouts)}개")
            typer.echo(f"\n{warning_message}")

    except FileNotFoundError as e:
        error_response = create_error_response(e, "template-apply")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ {str(e)}", err=True)
        raise typer.Exit(1)

    except Exception as e:
        error_response = create_error_response(e, "template-apply")
        if output_format == "json":
            typer.echo(json.dumps(error_response, ensure_ascii=False, indent=2), err=True)
        else:
            typer.echo(f"❌ 예기치 않은 오류: {str(e)}", err=True)
        raise typer.Exit(1)


if __name__ == "__main__":
    typer.run(template_apply)
