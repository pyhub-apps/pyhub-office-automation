"""
PowerPoint 자동화를 위한 공통 유틸리티 함수들
python-pptx 기반 PowerPoint 조작 및 데이터 처리 지원
"""

import json
import platform
import sys
import unicodedata
from enum import Enum
from pathlib import Path
from typing import Any, Dict, Optional, Union

from pyhub_office_automation.version import get_version


# CLI 명령어 인자를 위한 Enum 클래스들
class OutputFormat(str, Enum):
    """출력 형식 선택지"""

    JSON = "json"
    TEXT = "text"
    MARKDOWN = "markdown"


class LayoutType(str, Enum):
    """PowerPoint 레이아웃 타입"""

    TITLE = "title"
    TITLE_AND_CONTENT = "title_and_content"
    SECTION_HEADER = "section_header"
    TWO_CONTENT = "two_content"
    COMPARISON = "comparison"
    TITLE_ONLY = "title_only"
    BLANK = "blank"
    CONTENT_WITH_CAPTION = "content_with_caption"
    PICTURE_WITH_CAPTION = "picture_with_caption"


class ExportFormat(str, Enum):
    """프레젠테이션 내보내기 형식"""

    PDF = "pdf"
    PNG = "png"
    JPG = "jpg"
    JPEG = "jpeg"
    GIF = "gif"


class PowerPointBackend(str, Enum):
    """PowerPoint 백엔드 선택지"""

    PYTHON_PPTX = "python-pptx"
    COM = "com"
    AUTO = "auto"


class PlaceholderType(str, Enum):
    """Placeholder 타입 (Issue #77)"""

    TITLE = "title"
    BODY = "body"
    SUBTITLE = "subtitle"


class ShapeType(str, Enum):
    """도형 타입 (Issue #77)"""

    RECTANGLE = "rectangle"
    ROUNDED_RECTANGLE = "rounded-rectangle"
    ELLIPSE = "ellipse"
    ARROW_RIGHT = "arrow-right"
    ARROW_LEFT = "arrow-left"
    ARROW_UP = "arrow-up"
    ARROW_DOWN = "arrow-down"
    STAR = "star"
    PENTAGON = "pentagon"
    HEXAGON = "hexagon"


class ChartType(str, Enum):
    """차트 타입 (Issue #78)"""

    COLUMN = "column"
    BAR = "bar"
    LINE = "line"
    PIE = "pie"
    AREA = "area"
    SCATTER = "scatter"
    DOUGHNUT = "doughnut"


def get_powerpoint_backend(force_backend: Optional[str] = None) -> str:
    """
    현재 플랫폼에서 사용 가능한 PowerPoint 백엔드를 반환합니다.

    COM-First 전략: Windows에서는 COM 우선, 기타는 python-pptx

    Args:
        force_backend: 강제로 사용할 백엔드 ("com", "python-pptx", None)

    Returns:
        str: 'python-pptx' (크로스플랫폼) 또는 'com' (Windows 전용)

    Raises:
        ImportError: 사용 가능한 백엔드가 없는 경우
        ValueError: 지원하지 않는 백엔드가 지정된 경우
    """
    # backend_selector 임포트 (순환 임포트 방지)
    from .backend_selector import detect_backend

    return detect_backend(force_backend)


def check_feature_availability(feature: str) -> Dict[str, Any]:
    """
    특정 기능의 플랫폼별 가용성을 체크합니다.

    Args:
        feature: 체크할 기능 이름

    Returns:
        Dict: 기능 가용성 정보
            - available: 현재 플랫폼에서 사용 가능 여부
            - backend: 사용할 백엔드 ('python-pptx' or 'com')
            - platform: 현재 플랫폼
            - limitations: 제한사항 목록
    """
    current_platform = platform.system()
    backend = get_powerpoint_backend()

    # 기본 기능 (모든 플랫폼 지원)
    basic_features = [
        "presentation-create",
        "presentation-open",
        "presentation-save",
        "slide-add",
        "slide-delete",
        "content-add-text",
        "content-add-image",
    ]

    # COM 전용 기능 (Windows만)
    com_only_features = ["macro-run", "animation-add", "transition-set"]

    limitations = []
    available = True

    if feature in com_only_features:
        if current_platform != "Windows":
            available = False
            limitations.append(f"{feature}는 Windows에서만 사용 가능합니다")
        elif backend != PowerPointBackend.COM.value:
            available = False
            limitations.append(f"{feature}는 Windows COM (pywin32)이 필요합니다")

    return {
        "available": available,
        "backend": backend,
        "platform": current_platform,
        "feature": feature,
        "limitations": limitations,
    }


def normalize_path(path: str) -> str:
    """
    경로의 한글 문자를 정규화합니다 (macOS 자소분리 문제 해결).

    Args:
        path: 정규화할 경로 문자열

    Returns:
        정규화된 경로 문자열
    """
    if not isinstance(path, str):
        return path

    # macOS에서 한글 자소분리 문제 해결
    if platform.system() == "Darwin":
        # NFD -> NFC 정규화 (자소 결합)
        return unicodedata.normalize("NFC", path)

    return path


def create_success_response(
    command: str,
    data: Any,
    message: Optional[str] = None,
    version: Optional[str] = None,
) -> Dict[str, Any]:
    """
    성공 응답을 표준 JSON 형식으로 생성합니다.

    Args:
        command: 실행한 명령어 이름
        data: 반환할 데이터
        message: 추가 메시지 (선택사항)
        version: 패키지 버전 (기본값: 현재 버전)

    Returns:
        Dict: 표준화된 성공 응답
    """
    response = {
        "success": True,
        "command": command,
        "data": data,
        "version": version or get_version(),
    }

    if message:
        response["message"] = message

    return response


def create_error_response(
    command: str,
    error: Union[str, Exception],
    error_type: Optional[str] = None,
    version: Optional[str] = None,
) -> Dict[str, Any]:
    """
    에러 응답을 표준 JSON 형식으로 생성합니다.

    Args:
        command: 실행한 명령어 이름
        error: 에러 메시지 또는 Exception 객체
        error_type: 에러 타입 (선택사항)
        version: 패키지 버전 (기본값: 현재 버전)

    Returns:
        Dict: 표준화된 에러 응답
    """
    error_message = str(error)
    if isinstance(error, Exception):
        error_type = error_type or type(error).__name__

    response = {
        "success": False,
        "command": command,
        "error": error_message,
        "error_type": error_type or "UnknownError",
        "version": version or get_version(),
    }

    return response


def get_or_open_presentation(
    file_path: Optional[str] = None,
    presentation_name: Optional[str] = None,
    backend: str = "auto",
) -> Any:
    """
    파일 경로 또는 프레젠테이션 이름으로 프레젠테이션을 열거나 가져옵니다.

    Args:
        file_path: 프레젠테이션 파일 경로
        presentation_name: 열려있는 프레젠테이션 이름
        backend: 사용할 백엔드 ('auto', 'python-pptx', 'com')

    Returns:
        Presentation 객체 (python-pptx) 또는 COM 객체

    Raises:
        ValueError: 파일 경로와 프레젠테이션 이름이 모두 없거나 둘 다 제공된 경우
        FileNotFoundError: 파일이 존재하지 않는 경우
        ImportError: 필요한 라이브러리가 설치되지 않은 경우
    """
    # 입력 검증
    if file_path and presentation_name:
        raise ValueError("file_path와 presentation_name 중 하나만 지정해야 합니다")

    if not file_path and not presentation_name:
        raise ValueError("file_path 또는 presentation_name 중 하나는 필수입니다")

    # 백엔드 선택
    if backend == "auto":
        backend = get_powerpoint_backend()

    # 경로 정규화
    if file_path:
        file_path = normalize_path(file_path)
        if not Path(file_path).exists():
            raise FileNotFoundError(f"프레젠테이션 파일을 찾을 수 없습니다: {file_path}")

    # python-pptx 백엔드
    if backend == PowerPointBackend.PYTHON_PPTX.value:
        try:
            from pptx import Presentation

            if file_path:
                return Presentation(file_path)
            else:
                raise NotImplementedError(
                    "python-pptx는 프레젠테이션 이름으로 열기를 지원하지 않습니다. file_path를 사용하세요."
                )

        except ImportError:
            raise ImportError("python-pptx 패키지가 설치되지 않았습니다. 'pip install python-pptx'로 설치하세요.")

    # COM 백엔드 (Windows 전용)
    elif backend == PowerPointBackend.COM.value:
        if platform.system() != "Windows":
            raise NotImplementedError("COM 백엔드는 Windows에서만 사용 가능합니다")

        try:
            import win32com.client

            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            powerpoint.Visible = True

            if file_path:
                return powerpoint.Presentations.Open(str(Path(file_path).resolve()))
            else:
                # 이름으로 찾기
                for pres in powerpoint.Presentations:
                    if pres.Name == presentation_name:
                        return pres
                raise FileNotFoundError(f"열려있는 프레젠테이션을 찾을 수 없습니다: {presentation_name}")

        except ImportError:
            raise ImportError("pywin32 패키지가 설치되지 않았습니다. 'pip install pywin32'로 설치하세요.")

    else:
        raise ValueError(f"지원하지 않는 백엔드입니다: {backend}")


def output_result(result: Dict[str, Any], output_format: str = "json") -> None:
    """
    결과를 지정된 형식으로 출력합니다.

    Args:
        result: 출력할 결과 딕셔너리
        output_format: 출력 형식 ('json', 'text', 'markdown')
    """
    if output_format == OutputFormat.JSON.value:
        try:
            # UTF-8 인코딩으로 JSON 출력
            json_output = json.dumps(result, ensure_ascii=False, indent=2)
            print(json_output)
        except UnicodeEncodeError:
            # 인코딩 실패 시 ASCII로 폴백
            json_output = json.dumps(result, ensure_ascii=True, indent=2)
            print(json_output)

    elif output_format == OutputFormat.TEXT.value:
        # 텍스트 형식 출력
        if result.get("success"):
            print(f"✓ {result.get('command')} - Success")
            if result.get("message"):
                print(f"  {result['message']}")
        else:
            print(f"✗ {result.get('command')} - Error")
            print(f"  {result.get('error')}")

    elif output_format == OutputFormat.MARKDOWN.value:
        # 마크다운 형식 출력
        print(f"## {result.get('command')}")
        print()
        if result.get("success"):
            print("**Status**: ✓ Success")
            if result.get("message"):
                print(f"**Message**: {result['message']}")
        else:
            print("**Status**: ✗ Error")
            print(f"**Error**: {result.get('error')}")
        print()


def validate_output_format(format_str: str) -> str:
    """
    출력 형식 문자열을 검증합니다.

    Args:
        format_str: 검증할 형식 문자열

    Returns:
        검증된 형식 문자열

    Raises:
        ValueError: 지원하지 않는 형식인 경우
    """
    format_lower = format_str.lower()
    valid_formats = [f.value for f in OutputFormat]

    if format_lower not in valid_formats:
        raise ValueError(f"지원하지 않는 출력 형식입니다: {format_str}. 사용 가능: {', '.join(valid_formats)}")

    return format_lower


# Slide 관리 유틸리티 함수들 (Issue #76)


def get_layout_by_name_or_index(prs, identifier: Union[str, int]):
    """
    레이아웃을 이름(str) 또는 인덱스(int)로 찾습니다.

    Args:
        prs: Presentation 객체
        identifier: 레이아웃 이름(문자열) 또는 인덱스(정수)

    Returns:
        SlideLayout 객체

    Raises:
        ValueError: 레이아웃을 찾을 수 없는 경우
        IndexError: 인덱스가 범위를 벗어난 경우
    """
    # 인덱스로 찾기
    if isinstance(identifier, int):
        try:
            return prs.slide_layouts[identifier]
        except IndexError:
            raise IndexError(f"레이아웃 인덱스 범위 초과: {identifier} (사용 가능: 0-{len(prs.slide_layouts)-1})")

    # 이름으로 찾기
    if isinstance(identifier, str):
        for layout in prs.slide_layouts:
            if layout.name == identifier:
                return layout

        # 사용 가능한 레이아웃 목록 제공
        available_layouts = [f"{i}: {layout.name}" for i, layout in enumerate(prs.slide_layouts)]
        raise ValueError(
            f"레이아웃을 찾을 수 없습니다: '{identifier}'\n" f"사용 가능한 레이아웃:\n" + "\n".join(available_layouts)
        )

    raise TypeError(f"레이아웃 식별자는 문자열 또는 정수여야 합니다: {type(identifier)}")


def validate_slide_number(slide_num: int, total_slides: int, allow_append: bool = False) -> int:
    """
    슬라이드 번호를 검증하고 0-based 인덱스로 변환합니다.

    Args:
        slide_num: 검증할 슬라이드 번호 (1-based)
        total_slides: 총 슬라이드 수
        allow_append: True면 total_slides+1도 허용 (끝에 추가 시)

    Returns:
        int: 0-based 인덱스

    Raises:
        TypeError: 슬라이드 번호가 정수가 아닌 경우
        ValueError: 슬라이드 번호가 범위를 벗어난 경우
    """
    if not isinstance(slide_num, int):
        raise TypeError(f"슬라이드 번호는 정수여야 합니다: {slide_num} ({type(slide_num).__name__})")

    max_value = total_slides + 1 if allow_append else total_slides

    if total_slides == 0:
        # 빈 프레젠테이션인 경우
        if allow_append and slide_num == 1:
            return 0
        raise ValueError("프레젠테이션에 슬라이드가 없습니다")

    if not (1 <= slide_num <= max_value):
        raise ValueError(f"슬라이드 번호 범위: 1-{max_value}, 입력값: {slide_num}")

    return slide_num - 1  # 0-based 인덱스 반환


def get_slide_content_summary(slide) -> Dict[str, int]:
    """
    슬라이드의 콘텐츠 요약을 생성합니다 (도형 타입별 개수).

    Args:
        slide: Slide 객체

    Returns:
        Dict[str, int]: 도형 타입별 개수
            - textbox: 텍스트 박스 수
            - picture: 이미지 수
            - chart: 차트 수
            - table: 표 수
            - other: 기타 도형 수

    Example:
        >>> summary = get_slide_content_summary(slide)
        >>> print(summary)
        {"textbox": 3, "picture": 1, "chart": 0, "table": 1, "other": 0}
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        # python-pptx 없으면 빈 요약 반환
        return {"textbox": 0, "picture": 0, "chart": 0, "table": 0, "other": 0}

    summary = {"textbox": 0, "picture": 0, "chart": 0, "table": 0, "other": 0}

    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                summary["textbox"] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                summary["picture"] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                summary["chart"] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                summary["table"] += 1
            else:
                summary["other"] += 1
        except Exception:
            # 도형 타입 확인 실패 시 other로 카운트
            summary["other"] += 1

    return summary


def get_slide_title(slide) -> str:
    """
    슬라이드의 제목을 추출합니다.

    Args:
        slide: Slide 객체

    Returns:
        str: 슬라이드 제목 (없으면 빈 문자열)

    Example:
        >>> title = get_slide_title(slide)
        >>> print(title)
        "Q4 Sales Report"
    """
    try:
        if hasattr(slide, "shapes") and hasattr(slide.shapes, "title"):
            title_shape = slide.shapes.title
            if title_shape and hasattr(title_shape, "text"):
                return title_shape.text.strip()
    except Exception:
        # 제목 추출 실패 시 빈 문자열 반환
        pass

    return ""


# Content 관리 유틸리티 함수들 (Issue #77)


def get_placeholder_by_type(slide, placeholder_type: str):
    """
    Placeholder 타입으로 슬라이드의 placeholder를 찾습니다.

    Args:
        slide: Slide 객체
        placeholder_type: placeholder 타입 (title, body, subtitle)

    Returns:
        Placeholder 객체

    Raises:
        ValueError: placeholder를 찾을 수 없는 경우
    """
    try:
        if placeholder_type == PlaceholderType.TITLE.value:
            if hasattr(slide.shapes, "title"):
                return slide.shapes.title
            raise ValueError("이 슬라이드에는 title placeholder가 없습니다")

        elif placeholder_type == PlaceholderType.BODY.value:
            # Body는 일반적으로 index 1
            if 1 in slide.placeholders:
                return slide.placeholders[1]
            raise ValueError("이 슬라이드에는 body placeholder가 없습니다")

        elif placeholder_type == PlaceholderType.SUBTITLE.value:
            # Subtitle은 일반적으로 index 2
            if 2 in slide.placeholders:
                return slide.placeholders[2]
            raise ValueError("이 슬라이드에는 subtitle placeholder가 없습니다")

        else:
            raise ValueError(f"지원하지 않는 placeholder 타입: {placeholder_type}")

    except Exception as e:
        raise ValueError(f"Placeholder 접근 실패: {str(e)}")


def parse_color(color_str: str):
    """
    색상 문자열을 RGBColor 객체로 파싱합니다.

    Args:
        color_str: 색상 문자열 (red, blue, #FF0000 등)

    Returns:
        RGBColor 객체

    Raises:
        ValueError: 지원하지 않는 색상 형식인 경우

    Example:
        >>> color = parse_color("red")
        >>> color = parse_color("#FF0000")
    """
    try:
        from pptx.util import RGBColor
    except ImportError:
        raise ImportError("python-pptx 패키지가 필요합니다")

    # 색상 이름 매핑
    color_names = {
        "red": RGBColor(255, 0, 0),
        "green": RGBColor(0, 255, 0),
        "blue": RGBColor(0, 0, 255),
        "yellow": RGBColor(255, 255, 0),
        "orange": RGBColor(255, 165, 0),
        "purple": RGBColor(128, 0, 128),
        "black": RGBColor(0, 0, 0),
        "white": RGBColor(255, 255, 255),
        "gray": RGBColor(128, 128, 128),
    }

    color_lower = color_str.lower().strip()

    # 색상 이름으로 찾기
    if color_lower in color_names:
        return color_names[color_lower]

    # Hex 색상 파싱 (#RGB 또는 #RRGGBB)
    if color_str.startswith("#"):
        hex_color = color_str.lstrip("#")

        # #RGB 형식 (3자리)
        if len(hex_color) == 3:
            r = int(hex_color[0] * 2, 16)
            g = int(hex_color[1] * 2, 16)
            b = int(hex_color[2] * 2, 16)
            return RGBColor(r, g, b)

        # #RRGGBB 형식 (6자리)
        elif len(hex_color) == 6:
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)

    raise ValueError(
        f"지원하지 않는 색상 형식: {color_str}\n" f"사용 가능한 형식: 색상 이름 (red, blue 등), Hex (#FF0000, #F00)"
    )


def calculate_aspect_ratio_size(
    original_width: int, original_height: int, target_width: Optional[float] = None, target_height: Optional[float] = None
) -> tuple:
    """
    Aspect ratio를 유지하면서 크기를 계산합니다.

    Args:
        original_width: 원본 너비 (픽셀)
        original_height: 원본 높이 (픽셀)
        target_width: 목표 너비 (Inches, 선택)
        target_height: 목표 높이 (Inches, 선택)

    Returns:
        tuple: (width, height) in Inches

    Example:
        >>> # 원본 이미지 800x600, 너비 4인치로 조정
        >>> width, height = calculate_aspect_ratio_size(800, 600, target_width=4.0)
        >>> # 결과: (4.0, 3.0)
    """
    if original_width <= 0 or original_height <= 0:
        raise ValueError(f"유효하지 않은 크기: {original_width}x{original_height}")

    aspect_ratio = original_width / original_height

    # 너비만 지정된 경우
    if target_width and not target_height:
        return (target_width, target_width / aspect_ratio)

    # 높이만 지정된 경우
    elif target_height and not target_width:
        return (target_height * aspect_ratio, target_height)

    # 둘 다 지정된 경우 (aspect ratio 무시)
    elif target_width and target_height:
        return (target_width, target_height)

    # 둘 다 없는 경우 기본 크기 (6 inches width)
    else:
        default_width = 6.0
        return (default_width, default_width / aspect_ratio)


# Chart 관련 유틸리티 함수들 (Issue #78)


def parse_excel_range(excel_data: str) -> Dict[str, str]:
    """
    Excel 데이터 참조 문자열을 파싱합니다.

    Args:
        excel_data: Excel 참조 문자열 (예: "data.xlsx!A1:C10")

    Returns:
        Dict: {"file_path": "data.xlsx", "range": "A1:C10"}

    Raises:
        ValueError: 잘못된 형식인 경우

    Example:
        >>> result = parse_excel_range("sales.xlsx!Sheet1!B2:D20")
        >>> print(result)
        {"file_path": "sales.xlsx", "sheet": "Sheet1", "range": "B2:D20"}
    """
    if "!" not in excel_data:
        raise ValueError(
            f"Excel 데이터 참조 형식이 잘못되었습니다: {excel_data}\n예: 'file.xlsx!A1:C10' 또는 'file.xlsx!Sheet1!A1:C10'"
        )

    parts = excel_data.split("!")

    if len(parts) == 2:
        # file.xlsx!A1:C10 형식
        return {"file_path": parts[0].strip(), "sheet": None, "range": parts[1].strip()}
    elif len(parts) == 3:
        # file.xlsx!Sheet1!A1:C10 형식
        return {"file_path": parts[0].strip(), "sheet": parts[1].strip(), "range": parts[2].strip()}
    else:
        raise ValueError(
            f"Excel 데이터 참조 형식이 잘못되었습니다: {excel_data}\n예: 'file.xlsx!A1:C10' 또는 'file.xlsx!Sheet1!A1:C10'"
        )


def load_data_from_csv(csv_path: str):
    """
    CSV 파일에서 데이터를 로드합니다.

    Args:
        csv_path: CSV 파일 경로

    Returns:
        pandas.DataFrame

    Raises:
        FileNotFoundError: 파일이 없는 경우
        ValueError: CSV 파일 읽기 실패
    """
    try:
        import pandas as pd
    except ImportError:
        raise ImportError("pandas 패키지가 필요합니다. 'pip install pandas'로 설치하세요.")

    csv_path = normalize_path(csv_path)
    csv_file = Path(csv_path).resolve()

    if not csv_file.exists():
        raise FileNotFoundError(f"CSV 파일을 찾을 수 없습니다: {csv_file}")

    try:
        df = pd.read_csv(str(csv_file))
        return df
    except Exception as e:
        raise ValueError(f"CSV 파일 읽기 실패: {str(e)}")


def load_data_from_excel(file_path: str, sheet_name: Optional[str] = None, range_addr: Optional[str] = None):
    """
    Excel 파일에서 데이터를 로드합니다 (xlwings 사용).

    Args:
        file_path: Excel 파일 경로
        sheet_name: 시트 이름 (None이면 첫 번째 시트)
        range_addr: 범위 주소 (예: "A1:C10", None이면 사용된 범위 전체)

    Returns:
        pandas.DataFrame

    Raises:
        FileNotFoundError: 파일이 없는 경우
        ValueError: Excel 파일 읽기 실패
    """
    try:
        import pandas as pd
        import xlwings as xw
    except ImportError as e:
        raise ImportError(f"필요한 패키지가 없습니다: {str(e)}\n'pip install xlwings pandas'로 설치하세요.")

    file_path = normalize_path(file_path)
    excel_file = Path(file_path).resolve()

    if not excel_file.exists():
        raise FileNotFoundError(f"Excel 파일을 찾을 수 없습니다: {excel_file}")

    try:
        # xlwings로 Excel 파일 열기
        with xw.App(visible=False) as app:
            wb = app.books.open(str(excel_file))
            try:
                # 시트 선택
                if sheet_name:
                    sheet = wb.sheets[sheet_name]
                else:
                    sheet = wb.sheets[0]

                # 데이터 읽기
                if range_addr:
                    data = sheet.range(range_addr).options(pd.DataFrame, header=True, index=False).value
                else:
                    # 사용된 범위 전체 읽기
                    data = sheet.used_range.options(pd.DataFrame, header=True, index=False).value

                return data
            finally:
                wb.close()

    except Exception as e:
        raise ValueError(f"Excel 파일 읽기 실패: {str(e)}")


def create_chart_data(df, chart_type: str):
    """
    pandas DataFrame을 python-pptx ChartData로 변환합니다.

    Args:
        df: pandas DataFrame
        chart_type: 차트 타입 (column, bar, line, pie 등)

    Returns:
        ChartData 객체

    Raises:
        ImportError: python-pptx가 설치되지 않은 경우
        ValueError: 데이터 형식이 잘못된 경우
    """
    try:
        from pptx.chart.data import CategoryChartData
    except ImportError:
        raise ImportError("python-pptx 패키지가 필요합니다. 'pip install python-pptx'로 설치하세요.")

    if df.empty:
        raise ValueError("DataFrame이 비어있습니다")

    chart_data = CategoryChartData()

    # 카테고리 설정 (첫 번째 열)
    categories = df.iloc[:, 0].tolist()
    chart_data.categories = categories

    # 시리즈 추가 (나머지 열들)
    for col_name in df.columns[1:]:
        series_values = df[col_name].tolist()
        chart_data.add_series(col_name, series_values)

    return chart_data


def get_shape_by_index_or_name(slide, identifier: Union[int, str]):
    """
    슬라이드에서 shape를 인덱스 또는 이름으로 찾습니다.

    Args:
        slide: Slide 객체
        identifier: shape 인덱스(int) 또는 이름(str)

    Returns:
        Shape 객체

    Raises:
        ValueError: shape를 찾을 수 없는 경우
        IndexError: 인덱스가 범위를 벗어난 경우
    """
    # 인덱스로 찾기
    if isinstance(identifier, int):
        if 0 <= identifier < len(slide.shapes):
            return slide.shapes[identifier]
        else:
            raise IndexError(f"Shape 인덱스 범위 초과: {identifier} (사용 가능: 0-{len(slide.shapes)-1})")

    # 이름으로 찾기
    if isinstance(identifier, str):
        for shape in slide.shapes:
            if shape.name == identifier:
                return shape

        # 사용 가능한 shape 목록 제공
        available_shapes = [f"{i}: {shape.name}" for i, shape in enumerate(slide.shapes)]
        raise ValueError(f"Shape를 찾을 수 없습니다: '{identifier}'\n사용 가능한 Shape:\n" + "\n".join(available_shapes))

    raise TypeError(f"Shape 식별자는 문자열 또는 정수여야 합니다: {type(identifier)}")


# COM 백엔드 관련 유틸리티 함수들 (Issue #84)


def get_or_open_presentation(
    file_path: Optional[str] = None,
    presentation_name: Optional[str] = None,
    backend: str = "auto",
) -> tuple:
    """
    파일 경로 또는 프레젠테이션 이름으로 프레젠테이션을 가져오거나 엽니다.

    COM-First 전략: backend="auto"일 때 자동으로 최적 백엔드 선택

    Args:
        file_path: 프레젠테이션 파일 경로
        presentation_name: 열려있는 프레젠테이션 이름
        backend: 사용할 백엔드 ('auto', 'python-pptx', 'com')

    Note:
        PowerPoint COM은 항상 visible=True로 실행됩니다 (API 제약사항)

    Returns:
        (backend_instance, presentation_object) 튜플
        - COM: (PowerPointCOM, Presentation COM 객체)
        - python-pptx: (None, Presentation 객체)

    Raises:
        ValueError: 파일 경로와 프레젠테이션 이름이 모두 없거나 둘 다 제공된 경우
        FileNotFoundError: 파일이 존재하지 않는 경우
        ImportError: 필요한 라이브러리가 설치되지 않은 경우

    Example:
        >>> # 자동 백엔드 선택 (COM 우선)
        >>> backend_inst, prs = get_or_open_presentation(file_path="report.pptx")
        >>> # ... 작업 수행 ...
        >>> if backend_inst:  # COM 백엔드인 경우
        ...     backend_inst.close()
    """
    # 입력 검증
    if file_path and presentation_name:
        raise ValueError("file_path와 presentation_name 중 하나만 지정해야 합니다")

    # 백엔드 선택
    if backend == "auto":
        backend = get_powerpoint_backend()

    backend_lower = backend.lower()

    # python-pptx 백엔드
    if backend_lower == PowerPointBackend.PYTHON_PPTX.value:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx 패키지가 설치되지 않았습니다. " "'pip install python-pptx'로 설치하세요")

        if file_path:
            # 경로 정규화 및 검증
            file_path_norm = normalize_path(file_path)
            if not Path(file_path_norm).exists():
                raise FileNotFoundError(f"프레젠테이션 파일을 찾을 수 없습니다: {file_path}")

            prs = Presentation(file_path_norm)
            return None, prs
        else:
            raise NotImplementedError(
                "python-pptx는 프레젠테이션 이름으로 열기를 지원하지 않습니다. "
                "file_path를 사용하거나 COM 백엔드(--backend com)를 사용하세요"
            )

    # COM 백엔드 (Windows 전용)
    elif backend_lower == PowerPointBackend.COM.value:
        from .com_backend import get_or_open_presentation_com

        return get_or_open_presentation_com(file_path=file_path, presentation_name=presentation_name)

    else:
        raise ValueError(f"지원하지 않는 백엔드입니다: {backend}")


def create_presentation_with_backend(save_path: Optional[str] = None, backend: str = "auto") -> tuple:
    """
    새 프레젠테이션을 생성합니다.

    Args:
        save_path: 저장 경로 (선택)
        backend: 사용할 백엔드 ('auto', 'python-pptx', 'com')

    Returns:
        (backend_instance, presentation_object) 튜플

    Example:
        >>> backend_inst, prs = create_presentation_with_backend(save_path="new.pptx")
    """
    # 백엔드 선택
    if backend == "auto":
        backend = get_powerpoint_backend()

    backend_lower = backend.lower()

    # python-pptx 백엔드
    if backend_lower == PowerPointBackend.PYTHON_PPTX.value:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx 패키지가 필요합니다")

        prs = Presentation()

        if save_path:
            save_path_norm = normalize_path(save_path)
            prs.save(save_path_norm)

        return None, prs

    # COM 백엔드
    elif backend_lower == PowerPointBackend.COM.value:
        from .com_backend import PowerPointCOM

        ppt = PowerPointCOM()
        prs = ppt.create_presentation(save_path)

        return ppt, prs

    else:
        raise ValueError(f"지원하지 않는 백엔드입니다: {backend}")
