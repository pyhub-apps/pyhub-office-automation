"""
Interactive PowerPoint shell for stateful PowerPoint automation
Issue #85 Phase 5: https://github.com/pyhub-apps/pyhub-office-automation/issues/85
"""

import json
import shlex
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

import typer
from prompt_toolkit import PromptSession
from prompt_toolkit.completion import Completer, Completion
from prompt_toolkit.history import FileHistory
from prompt_toolkit.styles import Style
from rich.console import Console
from rich.table import Table

# Other Commands (2)
from pyhub_office_automation.powerpoint.animation_add import animation_add

# Content Commands (11)
from pyhub_office_automation.powerpoint.content_add_audio import content_add_audio
from pyhub_office_automation.powerpoint.content_add_chart import content_add_chart
from pyhub_office_automation.powerpoint.content_add_equation import content_add_equation
from pyhub_office_automation.powerpoint.content_add_excel_chart import content_add_excel_chart
from pyhub_office_automation.powerpoint.content_add_image import content_add_image
from pyhub_office_automation.powerpoint.content_add_shape import content_add_shape
from pyhub_office_automation.powerpoint.content_add_smartart import content_add_smartart
from pyhub_office_automation.powerpoint.content_add_table import content_add_table
from pyhub_office_automation.powerpoint.content_add_text import content_add_text
from pyhub_office_automation.powerpoint.content_add_video import content_add_video
from pyhub_office_automation.powerpoint.content_update import content_update

# Export Commands (3)
from pyhub_office_automation.powerpoint.export_images import export_images
from pyhub_office_automation.powerpoint.export_notes import export_notes
from pyhub_office_automation.powerpoint.export_pdf import export_pdf

# Layout & Theme Commands (4)
from pyhub_office_automation.powerpoint.layout_apply import layout_apply
from pyhub_office_automation.powerpoint.layout_list import layout_list

# PowerPoint command imports
# Presentation Commands (5)
from pyhub_office_automation.powerpoint.presentation_create import presentation_create
from pyhub_office_automation.powerpoint.presentation_info import presentation_info
from pyhub_office_automation.powerpoint.presentation_list import presentation_list
from pyhub_office_automation.powerpoint.presentation_open import presentation_open
from pyhub_office_automation.powerpoint.presentation_save import presentation_save
from pyhub_office_automation.powerpoint.run_macro import run_macro

# Slide Commands (6)
from pyhub_office_automation.powerpoint.slide_add import slide_add
from pyhub_office_automation.powerpoint.slide_copy import slide_copy
from pyhub_office_automation.powerpoint.slide_delete import slide_delete
from pyhub_office_automation.powerpoint.slide_duplicate import slide_duplicate
from pyhub_office_automation.powerpoint.slide_list import slide_list
from pyhub_office_automation.powerpoint.slide_reorder import slide_reorder

# Slideshow Commands (2)
from pyhub_office_automation.powerpoint.slideshow_control import slideshow_control
from pyhub_office_automation.powerpoint.slideshow_start import slideshow_start
from pyhub_office_automation.powerpoint.template_apply import template_apply
from pyhub_office_automation.powerpoint.theme_apply import theme_apply

console = Console()


@dataclass
class PptShellContext:
    """PowerPoint shell session context"""

    presentation_path: Optional[str] = None  # Full file path
    presentation_name: Optional[str] = None  # Just filename for display
    slide_number: Optional[int] = None  # Current slide number (1-indexed)
    prs: Optional[object] = None  # python-pptx Presentation object

    def get_prompt_text(self) -> str:
        """Generate prompt text with current context"""
        prs_name = self.presentation_name or "None"
        slide = str(self.slide_number) if self.slide_number else "None"
        return f"[PPT: {prs_name} > Slide {slide}] > "

    def update_presentation(self, file_path: str):
        """Update presentation context"""
        try:
            from pptx import Presentation

            self.presentation_path = str(Path(file_path).resolve())
            self.presentation_name = Path(file_path).name
            self.prs = Presentation(self.presentation_path)

            # Set to first slide if exists
            if len(self.prs.slides) > 0:
                self.slide_number = 1
            else:
                self.slide_number = None

            console.print(f"✓ Presentation set: {self.presentation_name}", style="green")
            if self.slide_number:
                console.print(f"✓ Active slide: {self.slide_number}/{len(self.prs.slides)}", style="green")

        except FileNotFoundError:
            console.print(f"Error: File not found: {file_path}", style="bold red")
        except Exception as e:
            console.print(f"Error loading presentation: {e}", style="bold red")

    def update_slide(self, slide_number: int):
        """Update active slide context"""
        if not self.prs:
            console.print("Error: No presentation loaded", style="bold red")
            return

        total_slides = len(self.prs.slides)
        if slide_number < 1 or slide_number > total_slides:
            console.print(f"Error: Slide number must be between 1 and {total_slides}", style="bold red")
            return

        self.slide_number = slide_number
        console.print(f"✓ Active slide: {self.slide_number}/{total_slides}", style="green")


class PptShellCompleter(Completer):
    """Tab completion for PowerPoint shell commands"""

    def __init__(self, context: PptShellContext):
        self.context = context
        # Shell-specific commands (8)
        self.shell_commands = [
            "help",
            "show",
            "use",
            "clear",
            "exit",
            "quit",
            "slides",
            "presentation-info",
        ]

        # All PPT CLI commands (33)
        self.ppt_commands = [
            # Presentation (5)
            "presentation-create",
            "presentation-open",
            "presentation-save",
            "presentation-list",
            "presentation-info",
            # Slide (6)
            "slide-list",
            "slide-add",
            "slide-delete",
            "slide-duplicate",
            "slide-copy",
            "slide-reorder",
            # Content (11)
            "content-add-text",
            "content-add-image",
            "content-add-shape",
            "content-add-table",
            "content-add-chart",
            "content-add-video",
            "content-add-smartart",
            "content-update",
            "content-add-excel-chart",
            "content-add-audio",
            "content-add-equation",
            # Layout & Theme (4)
            "layout-list",
            "layout-apply",
            "template-apply",
            "theme-apply",
            # Export (3)
            "export-pdf",
            "export-images",
            "export-notes",
            # Slideshow (2)
            "slideshow-start",
            "slideshow-control",
            # Other (2)
            "run-macro",
            "animation-add",
        ]

        self.all_commands = self.shell_commands + self.ppt_commands

    def get_completions(self, document, complete_event):
        """Generate completions based on current input"""
        text = document.text_before_cursor.lower()
        words = text.split()

        # Complete command names
        if len(words) <= 1:
            for cmd in self.all_commands:
                if cmd.startswith(text):
                    yield Completion(cmd, start_position=-len(text))

        # Complete "use" subcommands
        elif words[0] == "use":
            if len(words) == 1 or not words[1]:
                for sub in ["presentation", "slide"]:
                    yield Completion(sub, start_position=0)
            elif len(words) == 2:
                prefix = words[1].lower()
                for sub in ["presentation", "slide"]:
                    if sub.startswith(prefix):
                        yield Completion(sub, start_position=-len(prefix))


def parse_shell_args(command_line: str) -> list[str]:
    """Parse shell command line into arguments using shlex"""
    try:
        return shlex.split(command_line)
    except ValueError:
        # Handle unclosed quotes gracefully
        return command_line.split()


def execute_ppt_command(ctx: PptShellContext, command: str, args: list[str]) -> bool:
    """
    Execute a PowerPoint CLI command with context injection
    Returns True to continue shell, False to exit
    """
    # Map command names to functions (33 total)
    ppt_commands = {
        # Presentation (5)
        "presentation-create": presentation_create,
        "presentation-open": presentation_open,
        "presentation-save": presentation_save,
        "presentation-list": presentation_list,
        "presentation-info": presentation_info,
        # Slide (6)
        "slide-list": slide_list,
        "slide-add": slide_add,
        "slide-delete": slide_delete,
        "slide-duplicate": slide_duplicate,
        "slide-copy": slide_copy,
        "slide-reorder": slide_reorder,
        # Content (11)
        "content-add-text": content_add_text,
        "content-add-image": content_add_image,
        "content-add-shape": content_add_shape,
        "content-add-table": content_add_table,
        "content-add-chart": content_add_chart,
        "content-add-video": content_add_video,
        "content-add-smartart": content_add_smartart,
        "content-update": content_update,
        "content-add-excel-chart": content_add_excel_chart,
        "content-add-audio": content_add_audio,
        "content-add-equation": content_add_equation,
        # Layout & Theme (4)
        "layout-list": layout_list,
        "layout-apply": layout_apply,
        "template-apply": template_apply,
        "theme-apply": theme_apply,
        # Export (3)
        "export-pdf": export_pdf,
        "export-images": export_images,
        "export-notes": export_notes,
        # Slideshow (2)
        "slideshow-start": slideshow_start,
        "slideshow-control": slideshow_control,
        # Other (2)
        "run-macro": run_macro,
        "animation-add": animation_add,
    }

    if command not in ppt_commands:
        console.print(f"Error: Unknown PPT command: {command}", style="bold red")
        console.print("Hint: Type 'help' for available commands", style="yellow")
        return True

    # Context injection
    injected_args = list(args)

    # Check if --file-path is already specified
    has_file_arg = "--file-path" in args or "-f" in args

    # Inject --file-path if not specified and we have context
    if not has_file_arg and ctx.presentation_path:
        injected_args.extend(["--file-path", ctx.presentation_path])

    # Check if --slide-number is already specified
    has_slide_arg = "--slide-number" in args or "--slide" in args

    # Inject --slide-number if not specified and we have context
    # (only for commands that accept slide number)
    slide_commands = [
        "slide-delete",
        "slide-duplicate",
        "slide-copy",
        "content-add-text",
        "content-add-image",
        "content-add-shape",
        "content-add-table",
        "content-add-chart",
        "content-add-video",
        "content-add-smartart",
        "content-update",
        "content-add-excel-chart",
        "content-add-audio",
        "content-add-equation",
        "layout-apply",
        "animation-add",
    ]
    if command in slide_commands and not has_slide_arg and ctx.slide_number:
        injected_args.extend(["--slide-number", str(ctx.slide_number)])

    # Execute command
    try:
        func = ppt_commands[command]
        # Use typer's testing runner to execute command
        from typer.testing import CliRunner

        runner = CliRunner()
        app = typer.Typer()
        app.command()(func)

        result = runner.invoke(app, injected_args)

        if result.exit_code == 0:
            if result.stdout:
                console.print(result.stdout)
        else:
            console.print(f"Error executing {command}", style="bold red")
            if result.stdout:
                console.print(result.stdout)
            if result.stderr:
                console.print(result.stderr, style="red")

    except Exception as e:
        console.print(f"Error executing {command}: {e}", style="bold red")

    return True


def show_help():
    """Show available commands grouped by category"""
    console.print("\n[bold cyan]PowerPoint Shell Commands[/bold cyan]\n")

    categories = {
        "Shell Commands (8)": [
            ("help", "Show this help message"),
            ("show context", "Display current presentation and slide"),
            ("use presentation <path>", "Switch to a presentation"),
            ("use slide <number>", "Switch to a slide (1-indexed)"),
            ("slides", "List all slides in current presentation"),
            ("presentation-info", "Show current presentation details"),
            ("clear", "Clear the screen"),
            ("exit / quit", "Exit the shell"),
        ],
        "Presentation Commands (5)": [
            ("presentation-create", "Create new presentation"),
            ("presentation-open", "Open existing presentation"),
            ("presentation-save", "Save presentation"),
            ("presentation-list", "List open presentations"),
            ("presentation-info", "Get presentation details"),
        ],
        "Slide Commands (6)": [
            ("slide-list", "List all slides"),
            ("slide-add", "Add new slide"),
            ("slide-delete", "Delete slide"),
            ("slide-duplicate", "Duplicate slide"),
            ("slide-copy", "Copy slide to position"),
            ("slide-reorder", "Reorder slides"),
        ],
        "Content Commands (11)": [
            ("content-add-text", "Add text box"),
            ("content-add-image", "Add image"),
            ("content-add-shape", "Add shape"),
            ("content-add-table", "Add table"),
            ("content-add-chart", "Add chart"),
            ("content-add-video", "Add video"),
            ("content-add-smartart", "Add SmartArt"),
            ("content-add-excel-chart", "Add Excel chart"),
            ("content-add-audio", "Add audio"),
            ("content-add-equation", "Add equation"),
            ("content-update", "Update content"),
        ],
        "Layout & Theme (4)": [
            ("layout-list", "List available layouts"),
            ("layout-apply", "Apply layout to slide"),
            ("template-apply", "Apply template"),
            ("theme-apply", "Apply theme"),
        ],
        "Export Commands (3)": [
            ("export-pdf", "Export to PDF"),
            ("export-images", "Export slides as images"),
            ("export-notes", "Export slide notes"),
        ],
        "Slideshow (2)": [
            ("slideshow-start", "Start slideshow"),
            ("slideshow-control", "Control running slideshow"),
        ],
        "Other (2)": [
            ("run-macro", "Run VBA macro"),
            ("animation-add", "Add animation effect"),
        ],
    }

    for category, commands in categories.items():
        console.print(f"\n[yellow]{category}[/yellow]")
        for cmd, desc in commands:
            console.print(f"  [green]{cmd:30}[/green] {desc}")

    console.print("\n[dim]Tip: Use Tab for command completion, Up/Down arrows for history[/dim]\n")


def show_context(ctx: PptShellContext):
    """Display current context"""
    console.print("\n[bold]Current Context:[/bold]")
    console.print(f"  Presentation: {ctx.presentation_name or 'None'}")
    if ctx.presentation_path:
        console.print(f"  Path: {ctx.presentation_path}")
    if ctx.prs:
        total = len(ctx.prs.slides)
        console.print(f"  Total Slides: {total}")
    console.print(f"  Active Slide: {ctx.slide_number or 'None'}")
    console.print("\n[dim]All PowerPoint commands will use this context automatically.[/dim]\n")


def execute_shell_command(ctx: PptShellContext, command_line: str) -> bool:
    """
    Execute shell command
    Returns True to continue loop, False to exit
    """
    command_line = command_line.strip()
    if not command_line:
        return True

    # Parse command
    args = parse_shell_args(command_line)
    if not args:
        return True

    command = args[0].lower()

    # Shell-specific commands
    if command in ["exit", "quit"]:
        console.print("Goodbye!", style="cyan")
        return False

    elif command == "help":
        show_help()
        return True

    elif command == "clear":
        console.clear()
        return True

    elif command == "show":
        if len(args) > 1 and args[1].lower() == "context":
            show_context(ctx)
        else:
            console.print("Usage: show context", style="yellow")
        return True

    elif command == "use":
        if len(args) < 3:
            console.print("Usage: use presentation <path> | use slide <number>", style="yellow")
            return True

        subcommand = args[1].lower()
        value = " ".join(args[2:])  # Handle paths with spaces

        if subcommand == "presentation":
            # Remove quotes if present
            value = value.strip("\"'")
            ctx.update_presentation(value)
        elif subcommand == "slide":
            try:
                slide_num = int(value)
                ctx.update_slide(slide_num)
            except ValueError:
                console.print("Error: Slide number must be an integer", style="bold red")
        else:
            console.print(f"Error: Unknown use subcommand: {subcommand}", style="bold red")
            console.print("Usage: use presentation <path> | use slide <number>", style="yellow")

        return True

    elif command == "slides":
        # Shortcut for slide-list with context
        return execute_ppt_command(ctx, "slide-list", [])

    elif command == "presentation-info":
        # Shortcut for presentation-info with context
        return execute_ppt_command(ctx, "presentation-info", [])

    # Otherwise, try as PowerPoint command
    else:
        return execute_ppt_command(ctx, command, args[1:])


def ppt_shell(
    file_path: Optional[str] = typer.Option(None, "--file-path", "-f", help="Initial presentation file path"),
    presentation_name: Optional[str] = typer.Option(
        None, "--presentation-name", "-p", help="Initial presentation name (if already open)"
    ),
):
    """
    Interactive PowerPoint Shell Mode

    Provides a stateful REPL environment for PowerPoint automation where you can:
    - Set presentation and slide context once
    - Run commands without repeating --file-path and --slide-number
    - Use Tab completion for all 33 PowerPoint commands
    - Access command history with Up/Down arrows

    Examples:
        $ oa ppt shell
        $ oa ppt shell --file-path "presentation.pptx"
        $ oa ppt shell --presentation-name "MyPresentation.pptx"
    """
    console.print("\n[bold cyan]PowerPoint Shell Mode[/bold cyan]", style="bold")
    console.print("Type 'help' for available commands, 'exit' to quit.\n")

    # Initialize context
    ctx = PptShellContext()

    # Set initial presentation if provided
    if file_path:
        ctx.update_presentation(file_path)
    elif presentation_name:
        # Try to find presentation by name
        # (This would require PowerPoint COM automation on Windows)
        console.print("[yellow]Warning: --presentation-name not yet implemented, use --file-path[/yellow]")

    # Create prompt session with history and completion
    history_file = Path.home() / ".oa_ppt_shell_history"
    session = PromptSession(
        history=FileHistory(str(history_file)),
        completer=PptShellCompleter(ctx),
        style=Style.from_dict({"prompt": "cyan bold"}),
    )

    # REPL loop
    while True:
        try:
            # Get prompt text from context
            prompt_text = ctx.get_prompt_text()

            # Read command
            command_line = session.prompt(prompt_text)

            # Execute command
            should_continue = execute_shell_command(ctx, command_line)
            if not should_continue:
                break

        except KeyboardInterrupt:
            # Ctrl+C: Cancel current command, continue shell
            console.print("\n[dim](Use 'exit' or 'quit' to exit shell)[/dim]")
            continue
        except EOFError:
            # Ctrl+D: Exit shell
            console.print("\nGoodbye!", style="cyan")
            break
        except Exception as e:
            console.print(f"\nError: {e}", style="bold red")
            console.print("[dim]Shell continues...[/dim]")
            continue


if __name__ == "__main__":
    typer.run(ppt_shell)
