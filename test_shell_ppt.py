"""
PowerPoint Shell Mode 실전 테스트 스크립트
테스트용 PowerPoint 파일을 생성하고 기본 명령어를 검증합니다.
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

# Windows 터미널 인코딩 문제 해결
if sys.platform == "win32":
    import codecs

    sys.stdout = codecs.getwriter("utf-8")(sys.stdout.detach())
    sys.stderr = codecs.getwriter("utf-8")(sys.stderr.detach())


def create_test_pptx():
    """테스트용 PowerPoint 파일 생성"""
    print("=" * 60)
    print("PowerPoint Shell Mode 테스트 파일 생성")
    print("=" * 60)

    try:
        # 새 프레젠테이션 생성
        print("\n1. 새 프레젠테이션 생성 중...")
        prs = Presentation()

        # 슬라이드 1: 제목 슬라이드
        print("2. 슬라이드 1 (제목 슬라이드) 생성 중...")
        slide1_layout = prs.slide_layouts[0]  # Title Slide
        slide1 = prs.slides.add_slide(slide1_layout)
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
        title.text = "PowerPoint Shell Mode Test"
        subtitle.text = "Interactive Testing Presentation"

        # 슬라이드 2: 제목 + 내용
        print("3. 슬라이드 2 (제목 + 내용) 생성 중...")
        slide2_layout = prs.slide_layouts[1]  # Title and Content
        slide2 = prs.slides.add_slide(slide2_layout)
        title2 = slide2.shapes.title
        title2.text = "Sample Content Slide"
        body2 = slide2.placeholders[1]
        tf2 = body2.text_frame
        tf2.text = "This is a test slide"
        p = tf2.add_paragraph()
        p.text = "Second bullet point"
        p.level = 0

        # 슬라이드 3: 빈 슬라이드
        print("4. 슬라이드 3 (빈 슬라이드) 생성 중...")
        slide3_layout = prs.slide_layouts[6]  # Blank
        slide3 = prs.slides.add_slide(slide3_layout)

        # 텍스트 박스 추가
        left = Inches(2)
        top = Inches(2)
        width = Inches(4)
        height = Inches(1)
        txBox = slide3.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "This is a blank slide with a text box"

        # 슬라이드 4: 두 개의 내용
        print("5. 슬라이드 4 (두 개의 내용) 생성 중...")
        slide4_layout = prs.slide_layouts[3]  # Two Content
        slide4 = prs.slides.add_slide(slide4_layout)
        title4 = slide4.shapes.title
        title4.text = "Two Content Slide"

        # 슬라이드 5: 섹션 헤더
        print("6. 슬라이드 5 (섹션 헤더) 생성 중...")
        slide5_layout = prs.slide_layouts[2]  # Section Header
        slide5 = prs.slides.add_slide(slide5_layout)
        title5 = slide5.shapes.title
        title5.text = "Section Header"

        # 파일 저장
        test_file = Path("test_shell.pptx").absolute()
        print(f"\n7. 파일 저장 중: {test_file}")
        prs.save(test_file)

        print("\n" + "=" * 60)
        print("✓ 테스트 파일 생성 완료!")
        print("=" * 60)
        print(f"\n파일 경로: {test_file}")
        print(f"슬라이드 수: {len(prs.slides)}")

        slide_info = []
        for i, slide in enumerate(prs.slides, 1):
            slide_info.append(f"{i}. {slide.slide_layout.name}")
        print(f"슬라이드 정보:\n  " + "\n  ".join(slide_info))

        print("\n" + "=" * 60)
        print("PowerPoint Shell Mode 테스트 시작")
        print("=" * 60)
        print("\n다음 명령을 실행하세요:")
        print(f'  uv run oa ppt shell --file-path "{test_file}"')
        print("\nShell 내부에서 테스트할 명령어:")
        print("  1. show context")
        print("  2. slides")
        print("  3. presentation-info")
        print("  4. use slide 2")
        print("  5. content-add-text --text 'Hello Shell' --left 100 --top 100")
        print("  6. use slide 3")
        print("  7. content-add-shape --shape-type RECTANGLE --left 50 --top 50 --width 200 --height 100")
        print("  8. help")
        print("  9. exit")
        print("\n✓ 파일 생성 완료. 이제 Shell을 실행하세요!")

    except Exception as e:
        print(f"\n✗ 에러 발생: {e}")
        import traceback

        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    create_test_pptx()
